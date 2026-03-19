"""Google Drive folder → Dify Knowledge Base sync service.

Quét toàn bộ thư mục Google Drive, tự động:
- Google Sheets → convert → upload Dify
- DOCX/PDF/XLSX → download → pre-process → upload Dify
- Theo dõi thay đổi → chỉ sync file mới/đã sửa

Setup:
1. Create Service Account (same as Google Sheets sync)
2. Enable BOTH Google Sheets API + Google Drive API
3. Share the FOLDER with service account email (Viewer)
4. Call sync_folder(folder_id, dataset_id)
"""

import hashlib
import io
import json
import logging
import os
import re
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

import httpx

logger = logging.getLogger(__name__)

SYNC_STATE_PATH = Path(__file__).parent.parent.parent / "data" / "drive_sync_state.json"

# Supported file types
GOOGLE_SHEETS_MIME = "application/vnd.google-apps.spreadsheet"
GOOGLE_DOCS_MIME = "application/vnd.google-apps.document"
GOOGLE_SLIDES_MIME = "application/vnd.google-apps.presentation"
SUPPORTED_MIMES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/pdf": "pdf",
    "text/plain": "txt",
    "text/markdown": "md",
    GOOGLE_SHEETS_MIME: "gsheet",
    GOOGLE_DOCS_MIME: "gdoc",
    GOOGLE_SLIDES_MIME: "gslides",
}


class GoogleDriveSync:
    """Sync entire Google Drive folder to Dify Knowledge Base."""

    def __init__(
        self,
        credentials_path: str,
        dify_base_url: str,
        dify_dataset_api_key: str,
    ):
        self.credentials_path = credentials_path
        self.dify_base_url = dify_base_url
        self.dify_api_key = dify_dataset_api_key
        self._drive_service = None
        self._sheets_service = None
        self._sync_state = self._load_sync_state()

    def _get_credentials(self):
        from google.oauth2.service_account import Credentials
        return Credentials.from_service_account_file(
            self.credentials_path,
            scopes=[
                "https://www.googleapis.com/auth/drive.readonly",
                "https://www.googleapis.com/auth/spreadsheets.readonly",
            ],
        )

    def _get_drive_service(self):
        if self._drive_service is None:
            from googleapiclient.discovery import build
            self._drive_service = build("drive", "v3", credentials=self._get_credentials())
        return self._drive_service

    def _get_sheets_service(self):
        if self._sheets_service is None:
            from googleapiclient.discovery import build
            self._sheets_service = build("sheets", "v4", credentials=self._get_credentials())
        return self._sheets_service

    def _load_sync_state(self) -> dict:
        if SYNC_STATE_PATH.exists():
            try:
                return json.loads(SYNC_STATE_PATH.read_text())
            except Exception:
                pass
        return {}

    def _save_sync_state(self):
        SYNC_STATE_PATH.parent.mkdir(parents=True, exist_ok=True)
        SYNC_STATE_PATH.write_text(json.dumps(self._sync_state, ensure_ascii=False, indent=2))

    def list_folder(self, folder_id: str) -> list[dict]:
        """List all supported files in a Google Drive folder (recursive)."""
        drive = self._get_drive_service()
        all_files = []

        def _scan(fid: str, path: str = ""):
            query = f"'{fid}' in parents and trashed = false"
            results = drive.files().list(
                q=query,
                fields="files(id, name, mimeType, modifiedTime, size)",
                pageSize=100,
                supportsAllDrives=True,
                includeItemsFromAllDrives=True,
            ).execute()

            for f in results.get("files", []):
                f["path"] = f"{path}/{f['name']}" if path else f["name"]

                if f["mimeType"] == "application/vnd.google-apps.folder":
                    _scan(f["id"], f["path"])
                elif f["mimeType"] in SUPPORTED_MIMES:
                    all_files.append(f)

        _scan(folder_id)
        return all_files

    def _file_hash(self, file_info: dict) -> str:
        """Hash based on modifiedTime + size to detect changes."""
        key = f"{file_info['id']}:{file_info.get('modifiedTime', '')}:{file_info.get('size', '')}"
        return hashlib.sha256(key.encode()).hexdigest()[:16]

    def _download_file(self, file_id: str, mime_type: str) -> tuple[str, bytes]:
        """Download a file from Google Drive. Returns (extension, content_bytes)."""
        drive = self._get_drive_service()

        if mime_type == GOOGLE_SHEETS_MIME:
            # Export as xlsx
            content = drive.files().export(fileId=file_id, mimeType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet").execute()
            return "xlsx", content
        elif mime_type == GOOGLE_DOCS_MIME:
            # Try docx first, fallback to plain text for large files
            try:
                content = drive.files().export(fileId=file_id, mimeType="application/vnd.openxmlformats-officedocument.wordprocessingml.document").execute()
                return "docx", content
            except Exception:
                # Fallback: export as plain text (no size limit)
                content = drive.files().export(fileId=file_id, mimeType="text/plain").execute()
                if isinstance(content, str):
                    content = content.encode("utf-8")
                return "txt", content
        elif mime_type == GOOGLE_SLIDES_MIME:
            # Export as pptx
            content = drive.files().export(fileId=file_id, mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation").execute()
            return "pptx", content
        else:
            # Download directly
            content = drive.files().get_media(fileId=file_id).execute()
            ext = SUPPORTED_MIMES.get(mime_type, "bin")
            return ext, content

    def _sheet_to_markdown(self, spreadsheet_id: str, title: str) -> str:
        """Convert Google Sheet to markdown using Sheets API."""
        service = self._get_sheets_service()
        meta = service.spreadsheets().get(spreadsheetId=spreadsheet_id).execute()

        parts = [f"# {title}\n"]
        for sheet in meta.get("sheets", []):
            tab_title = sheet["properties"]["title"]
            data = service.spreadsheets().values().get(
                spreadsheetId=spreadsheet_id, range=tab_title
            ).execute()
            rows = data.get("values", [])
            if not rows:
                continue

            parts.append(f"## {tab_title}\n")
            for row in rows:
                cells = [str(c).strip() if c else "" for c in row]
                non_empty = [c for c in cells if c]
                if non_empty:
                    parts.append(" | ".join(non_empty))
            parts.append("")

        return "\n".join(parts)

    def _preprocess_file(self, ext: str, content: bytes, title: str) -> Optional[str]:
        """Pre-process file content to markdown text."""
        try:
            if ext in ("txt", "md"):
                return content.decode("utf-8", errors="replace")

            # Save to temp file for processing
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name

            try:
                if ext == "xlsx":
                    return self._process_excel(tmp_path, title)
                elif ext == "docx":
                    return self._process_docx(tmp_path, title)
                elif ext == "pdf":
                    return self._process_pdf(tmp_path, title)
                elif ext == "pptx":
                    return self._process_pptx(tmp_path, title)
            finally:
                os.unlink(tmp_path)
        except Exception as e:
            logger.error(f"Preprocess error for {title}: {e}")
            return None

    def _process_excel(self, path: str, title: str) -> str:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        parts = [f"# {title}\n"]
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## {sheet_name}\n")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() if c is not None and str(c).strip() != "None" else "" for c in row]
                non_empty = [c for c in cells if c]
                if non_empty:
                    parts.append(" | ".join(non_empty))
            parts.append("")
        return "\n".join(parts)

    def _process_docx(self, path: str, title: str) -> str:
        from docx import Document
        doc = Document(path)
        parts = [f"# {title}\n"]

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            if "Heading" in style:
                level = 2
                try:
                    level = int(style[-1])
                except Exception:
                    pass
                parts.append(f"{'#' * level} {text}")
            else:
                parts.append(text)

        for table in doc.tables:
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip().replace("\n", " | ") for cell in row.cells]
                cleaned = []
                prev = None
                for c in cells:
                    if c != prev:
                        cleaned.append(c)
                    prev = c
                if row_idx == 0:
                    parts.append("| " + " | ".join(cleaned) + " |")
                    parts.append("| " + " | ".join(["---"] * len(cleaned)) + " |")
                else:
                    parts.append("| " + " | ".join(cleaned) + " |")
            parts.append("")

        return "\n".join(parts)

    def _process_pdf(self, path: str, title: str) -> str:
        try:
            import pdfplumber
            parts = [f"# {title}\n"]
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text and len(text.strip()) > 20:
                        parts.append(f"## Trang {i + 1}\n{text}")
            return "\n".join(parts)
        except ImportError:
            # Fallback: just return filename
            return f"# {title}\n\n(PDF file - install pdfplumber for text extraction)"

    def _process_pptx(self, path: str, title: str) -> str:
        """Extract text from PowerPoint slides."""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            parts = [f"# {title}\n"]
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    parts.append(f"## Slide {i + 1}\n" + "\n".join(slide_text))
            return "\n".join(parts)
        except ImportError:
            return f"# {title}\n\n(PPTX file - install python-pptx for text extraction)"

    def _split_sections(self, markdown: str, title: str) -> list[dict]:
        """Split markdown into semantic sections for better RAG."""
        lines = markdown.split("\n")
        sections = []
        current = []
        heading = title
        pattern = re.compile(r"^(#{1,3}\s+|[IVX]+\.\s+|\d+\.\s+[A-ZĐ])")

        for line in lines:
            if pattern.match(line.strip()) and current:
                text = "\n".join(current).strip()
                if text and len(text) > 50:
                    sections.append({"name": f"{title} — {heading}", "text": text})
                current = []
                heading = line.strip().lstrip("#").strip()
            current.append(line)

        if current:
            text = "\n".join(current).strip()
            if text and len(text) > 50:
                sections.append({"name": f"{title} — {heading}", "text": text})

        # Merge small sections
        merged = []
        buf = None
        for sec in sections:
            if buf and len(buf["text"]) < 200:
                buf["text"] += "\n\n" + sec["text"]
                buf["name"] = sec["name"]
            else:
                if buf:
                    merged.append(buf)
                buf = sec
        if buf:
            merged.append(buf)

        return merged if merged else [{"name": title, "text": markdown}]

    async def sync_folder(
        self,
        folder_id: str,
        dataset_id: str,
        force: bool = False,
    ) -> dict:
        """Sync all files in a Google Drive folder to a Dify KB.

        Returns: {synced: int, skipped: int, errors: int, details: [...]}
        """
        files = self.list_folder(folder_id)
        results = {"synced": 0, "skipped": 0, "errors": 0, "details": []}
        state_key = f"folder:{folder_id}:{dataset_id}"

        if state_key not in self._sync_state:
            self._sync_state[state_key] = {"files": {}}

        folder_state = self._sync_state[state_key]["files"]

        async with httpx.AsyncClient(timeout=120.0) as client:
            for f in files:
                file_id = f["id"]
                file_name = f["name"]
                mime_type = f["mimeType"]
                file_hash = self._file_hash(f)

                # Skip if unchanged
                if not force and folder_state.get(file_id, {}).get("hash") == file_hash:
                    results["skipped"] += 1
                    results["details"].append({"file": file_name, "status": "skipped", "reason": "unchanged"})
                    continue

                try:
                    # Get markdown content
                    if mime_type == GOOGLE_SHEETS_MIME:
                        markdown = self._sheet_to_markdown(file_id, file_name)
                    else:
                        ext, content = self._download_file(file_id, mime_type)
                        markdown = self._preprocess_file(ext, content, file_name)
                        if not markdown:
                            results["errors"] += 1
                            results["details"].append({"file": file_name, "status": "error", "reason": "preprocess failed"})
                            continue

                    # Delete old docs for this file
                    old_doc_ids = folder_state.get(file_id, {}).get("doc_ids", [])
                    for old_id in old_doc_ids:
                        try:
                            await client.delete(
                                f"{self.dify_base_url}/datasets/{dataset_id}/documents/{old_id}",
                                headers={"Authorization": f"Bearer {self.dify_api_key}"},
                            )
                        except Exception:
                            pass

                    # Split into sections and upload
                    sections = self._split_sections(markdown, file_name)
                    new_doc_ids = []

                    for section in sections:
                        resp = await client.post(
                            f"{self.dify_base_url}/datasets/{dataset_id}/document/create_by_text",
                            headers={
                                "Authorization": f"Bearer {self.dify_api_key}",
                                "Content-Type": "application/json",
                            },
                            json={
                                "name": section["name"],
                                "text": section["text"],
                                "indexing_technique": "high_quality",
                                "process_rule": {"mode": "automatic"},
                            },
                        )
                        if resp.status_code == 200:
                            doc_id = resp.json().get("document", {}).get("id")
                            if doc_id:
                                new_doc_ids.append(doc_id)

                    # Update state
                    folder_state[file_id] = {
                        "hash": file_hash,
                        "name": file_name,
                        "doc_ids": new_doc_ids,
                        "sections": len(sections),
                        "synced_at": datetime.now(timezone.utc).isoformat(),
                    }

                    results["synced"] += 1
                    results["details"].append({
                        "file": file_name,
                        "status": "synced",
                        "sections": len(sections),
                    })

                except Exception as e:
                    logger.error(f"Sync error for {file_name}: {e}")
                    results["errors"] += 1
                    results["details"].append({"file": file_name, "status": "error", "reason": str(e)})

        self._sync_state[state_key]["last_sync"] = datetime.now(timezone.utc).isoformat()
        self._sync_state[state_key]["folder_id"] = folder_id
        self._sync_state[state_key]["dataset_id"] = dataset_id
        self._save_sync_state()

        return results

    def get_sync_status(self) -> list[dict]:
        """Get status of all synced folders."""
        return [
            {
                "folder_id": v.get("folder_id", ""),
                "dataset_id": v.get("dataset_id", ""),
                "last_sync": v.get("last_sync", ""),
                "files_count": len(v.get("files", {})),
                "files": [
                    {"name": f["name"], "synced_at": f.get("synced_at", ""), "sections": f.get("sections", 0)}
                    for f in v.get("files", {}).values()
                ],
            }
            for k, v in self._sync_state.items()
            if k.startswith("folder:")
        ]
