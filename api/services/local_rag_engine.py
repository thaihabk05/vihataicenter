"""
Local RAG Engine — replaces Dify with embedded Qdrant + Knowledge Graph.

Features:
- Qdrant embedded vector store (no external server needed)
- Hybrid search: semantic (vector) + keyword (BM25-like)
- Knowledge Graph for entity queries (people, products, policies)
- Contextual chunking with metadata headers
- Claude API for answer generation with strict grounding
"""

import os
import json
import hashlib
import re
import uuid
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from datetime import datetime

import networkx as nx

# ─── Configuration ───
DATA_DIR = Path(__file__).parent.parent.parent / "data"
RAG_DIR = DATA_DIR / "rag"
RAG_DIR.mkdir(parents=True, exist_ok=True)

QDRANT_PATH = str(RAG_DIR / "qdrant_db")
KG_PATH = RAG_DIR / "knowledge_graph.json"
CHUNK_INDEX_PATH = RAG_DIR / "chunk_index.json"

COLLECTION_NAME = "vihat_knowledge"
EMBED_MODEL = "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2"
EMBED_DIM = 384  # dimension for this model
CHUNK_SIZE = 2000  # chars per chunk (larger = fewer chunks, better for local mode)
CHUNK_OVERLAP = 100

# ─── Lazy-loaded globals ───
_qdrant_client = None
_embed_model = None
_knowledge_graph = None
_chunk_index = {}  # chunk_id -> metadata


import threading
_qdrant_lock = threading.Lock()
_qdrant_clients = {}  # thread_id -> client

def _get_qdrant():
    """Get or create Qdrant client per thread (embedded mode SQLite requires same-thread access)."""
    tid = threading.current_thread().ident
    if tid not in _qdrant_clients:
        with _qdrant_lock:
            if tid not in _qdrant_clients:
                from qdrant_client import QdrantClient
                client = QdrantClient(path=QDRANT_PATH)
                try:
                    client.get_collection(COLLECTION_NAME)
                except Exception:
                    from qdrant_client.models import VectorParams, Distance
                    client.create_collection(
                        collection_name=COLLECTION_NAME,
                        vectors_config=VectorParams(size=EMBED_DIM, distance=Distance.COSINE),
                    )
                _qdrant_clients[tid] = client
            print(f"[RAG] Created Qdrant client for thread {tid}")
    return _qdrant_clients[tid]


def _get_embedder():
    """Get or create sentence transformer model."""
    global _embed_model
    if _embed_model is None:
        from sentence_transformers import SentenceTransformer
        _embed_model = SentenceTransformer(EMBED_MODEL)
        print(f"[RAG] Loaded embedding model: {EMBED_MODEL}")
    return _embed_model


def _embed_texts(texts: List[str]) -> List[List[float]]:
    """Embed a list of texts."""
    model = _get_embedder()
    embeddings = model.encode(texts, show_progress_bar=False, normalize_embeddings=True)
    return embeddings.tolist()


def _load_chunk_index():
    """Load chunk index from disk."""
    global _chunk_index
    if CHUNK_INDEX_PATH.exists():
        try:
            _chunk_index = json.loads(CHUNK_INDEX_PATH.read_text(encoding="utf-8"))
        except Exception:
            _chunk_index = {}
    return _chunk_index


def _save_chunk_index():
    """Save chunk index to disk."""
    CHUNK_INDEX_PATH.write_text(
        json.dumps(_chunk_index, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


# ═══════════════════════════════════════════════════════════════
# DOCUMENT CHUNKING
# ═══════════════════════════════════════════════════════════════

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Split text into overlapping chunks, respecting paragraph boundaries."""
    if not text or not text.strip():
        return []

    # Split by double newlines (paragraphs) first
    paragraphs = re.split(r'\n\s*\n', text)
    chunks = []
    current_chunk = ""

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        if len(current_chunk) + len(para) + 2 <= chunk_size:
            current_chunk += ("\n\n" + para if current_chunk else para)
        else:
            if current_chunk:
                chunks.append(current_chunk)
            # If single paragraph is too long, split by sentences
            if len(para) > chunk_size:
                sentences = re.split(r'(?<=[.!?。])\s+', para)
                current_chunk = ""
                for sent in sentences:
                    if len(current_chunk) + len(sent) + 1 <= chunk_size:
                        current_chunk += (" " + sent if current_chunk else sent)
                    else:
                        if current_chunk:
                            chunks.append(current_chunk)
                        current_chunk = sent
            else:
                current_chunk = para

    if current_chunk:
        chunks.append(current_chunk)

    return chunks


def create_context_header(doc_meta: dict) -> str:
    """Create context header for a chunk."""
    parts = []
    if doc_meta.get("title"):
        parts.append(f"Tài liệu: {doc_meta['title']}")
    if doc_meta.get("description"):
        parts.append(f"Mô tả: {doc_meta['description'][:200]}")
    if doc_meta.get("tags"):
        parts.append(f"Sản phẩm: {', '.join(doc_meta['tags'])}")
    if doc_meta.get("knowledge_base"):
        kb_names = {"sales": "Kinh doanh", "hr": "Nhân sự", "accounting": "Kế toán",
                    "general": "Chung", "management": "Quản lý"}
        parts.append(f"Phòng ban: {kb_names.get(doc_meta['knowledge_base'], doc_meta['knowledge_base'])}")
    if doc_meta.get("drive_url"):
        parts.append(f"Link: {doc_meta['drive_url']}")
    return "\n".join(parts)


# ═══════════════════════════════════════════════════════════════
# INDEXING
# ═══════════════════════════════════════════════════════════════

def index_document(doc_id: str, text: str, doc_meta: dict) -> int:
    """Index a document into Qdrant. Returns number of chunks indexed."""
    from qdrant_client.models import PointStruct

    if not text or not text.strip():
        return 0

    # Remove old chunks for this doc
    remove_document(doc_id)

    # Chunk the text
    chunks = chunk_text(text)
    if not chunks:
        return 0

    # Create context header
    header = create_context_header(doc_meta)

    # Prepare chunk data (not PointStruct yet — build after embedding)
    chunk_data = []
    for i, chunk in enumerate(chunks):
        contextualized = f"{header}\n\n{chunk}" if header else chunk
        tags = doc_meta.get("tags") or []

        # Store in chunk index
        chunk_id = hashlib.md5(f"{doc_id}:{i}:{chunk[:50]}".encode()).hexdigest()
        _chunk_index[chunk_id] = {
            "doc_id": doc_id,
            "chunk_idx": i,
            "title": doc_meta.get("title", ""),
            "knowledge_base": doc_meta.get("knowledge_base", ""),
            "tags": tags,
            "drive_url": doc_meta.get("drive_url", ""),
            "text": chunk,
            "contextualized": contextualized,
        }

        chunk_data.append({
            "doc_id": doc_id,
            "chunk_idx": i,
            "title": doc_meta.get("title", ""),
            "text": contextualized,
            "original_text": chunk,
            "knowledge_base": doc_meta.get("knowledge_base", ""),
            "tags": tags,
            "drive_url": doc_meta.get("drive_url", ""),
        })

    # Embed all chunks
    texts_to_embed = [c["text"] for c in chunk_data]
    embeddings = _embed_texts(texts_to_embed)

    # Build PointStruct with proper IDs and vectors
    final_points = []
    for idx, (cd, emb) in enumerate(zip(chunk_data, embeddings)):
        int_id = int(hashlib.md5(f"{doc_id}:{idx}".encode()).hexdigest()[:15], 16)
        final_points.append(PointStruct(
            id=int_id,
            vector=emb,
            payload=cd,
        ))

    # Upsert to Qdrant
    client = _get_qdrant()
    BATCH_SIZE = 64
    for i in range(0, len(final_points), BATCH_SIZE):
        batch = final_points[i:i + BATCH_SIZE]
        client.upsert(collection_name=COLLECTION_NAME, points=batch)

    _save_chunk_index()
    return len(final_points)


def remove_document(doc_id: str):
    """Remove all chunks for a document from Qdrant."""
    from qdrant_client.models import Filter, FieldCondition, MatchValue

    try:
        client = _get_qdrant()
        client.delete(
            collection_name=COLLECTION_NAME,
            points_selector=Filter(
                must=[FieldCondition(key="doc_id", match=MatchValue(value=doc_id))]
            ),
        )
    except Exception as e:
        print(f"[RAG] Remove doc {doc_id} error: {e}")

    # Remove from chunk index
    to_remove = [cid for cid, meta in _chunk_index.items() if meta.get("doc_id") == doc_id]
    for cid in to_remove:
        del _chunk_index[cid]


def get_index_stats() -> dict:
    """Get indexing statistics."""
    try:
        client = _get_qdrant()
        info = client.get_collection(COLLECTION_NAME)
        return {
            "total_chunks": info.points_count,
            "total_docs": len(set(m.get("doc_id") for m in _chunk_index.values())),
            "collection": COLLECTION_NAME,
            "embed_model": EMBED_MODEL,
        }
    except Exception:
        return {"total_chunks": 0, "total_docs": 0}


# ═══════════════════════════════════════════════════════════════
# SEARCH / RETRIEVAL
# ═══════════════════════════════════════════════════════════════

def search(query: str, top_k: int = 8, knowledge_base: str = "",
           tags: List[str] = None) -> List[dict]:
    """
    Hybrid search: semantic (Qdrant vector) + keyword boost.
    Returns list of {text, score, title, doc_id, drive_url, ...}
    """
    from qdrant_client.models import Filter, FieldCondition, MatchValue

    client = _get_qdrant()

    # Build filter
    must_conditions = []
    if knowledge_base:
        must_conditions.append(
            FieldCondition(key="knowledge_base", match=MatchValue(value=knowledge_base))
        )

    query_filter = Filter(must=must_conditions) if must_conditions else None

    # Semantic search
    query_embedding = _embed_texts([query])[0]
    results = client.query_points(
        collection_name=COLLECTION_NAME,
        query=query_embedding,
        query_filter=query_filter,
        limit=top_k * 2,  # get more for keyword reranking
        with_payload=True,
    ).points

    # Keyword boost — boost results that contain query keywords
    query_words = set(re.findall(r'\w+', query.lower()))
    scored_results = []
    for r in results:
        text = r.payload.get("original_text", "").lower()
        # Count keyword matches
        keyword_matches = sum(1 for w in query_words if w in text and len(w) > 2)
        keyword_boost = min(keyword_matches * 0.05, 0.2)  # max 0.2 boost
        final_score = r.score + keyword_boost

        scored_results.append({
            "text": r.payload.get("text", ""),
            "original_text": r.payload.get("original_text", ""),
            "score": round(final_score, 4),
            "semantic_score": round(r.score, 4),
            "keyword_boost": round(keyword_boost, 4),
            "title": r.payload.get("title", ""),
            "doc_id": r.payload.get("doc_id", ""),
            "knowledge_base": r.payload.get("knowledge_base", ""),
            "tags": r.payload.get("tags", []),
            "drive_url": r.payload.get("drive_url", ""),
        })

    # Sort by final score and take top_k
    scored_results.sort(key=lambda x: x["score"], reverse=True)
    return scored_results[:top_k]


# ═══════════════════════════════════════════════════════════════
# KNOWLEDGE GRAPH
# ═══════════════════════════════════════════════════════════════

def _get_kg() -> nx.DiGraph:
    """Get or create knowledge graph."""
    global _knowledge_graph
    if _knowledge_graph is None:
        _knowledge_graph = nx.DiGraph()
        if KG_PATH.exists():
            try:
                data = json.loads(KG_PATH.read_text(encoding="utf-8"))
                for node in data.get("nodes", []):
                    _knowledge_graph.add_node(node["id"], **node.get("attrs", {}))
                for edge in data.get("edges", []):
                    _knowledge_graph.add_edge(edge["from"], edge["to"], **edge.get("attrs", {}))
                print(f"[KG] Loaded {len(_knowledge_graph.nodes)} nodes, {len(_knowledge_graph.edges)} edges")
            except Exception as e:
                print(f"[KG] Load error: {e}")
    return _knowledge_graph


def _save_kg():
    """Save knowledge graph to disk."""
    G = _get_kg()
    data = {
        "nodes": [{"id": n, "attrs": dict(G.nodes[n])} for n in G.nodes],
        "edges": [{"from": u, "to": v, "attrs": dict(G.edges[u, v])} for u, v in G.edges],
    }
    KG_PATH.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def add_entity(entity_id: str, entity_type: str, name: str, **attrs):
    """Add an entity to the knowledge graph."""
    G = _get_kg()
    G.add_node(entity_id, type=entity_type, name=name, **attrs)


def add_relation(from_id: str, to_id: str, relation: str, **attrs):
    """Add a relation between entities."""
    G = _get_kg()
    G.add_edge(from_id, to_id, relation=relation, **attrs)


def query_kg(query: str) -> List[dict]:
    """Query knowledge graph for entity information."""
    G = _get_kg()
    if not G.nodes:
        return []

    results = []
    query_lower = query.lower()

    # Search nodes by name
    for node_id in G.nodes:
        node = G.nodes[node_id]
        name = node.get("name", "").lower()
        if not name:
            continue

        # Fuzzy match: any word in query matches node name
        query_words = re.findall(r'\w+', query_lower)
        name_words = re.findall(r'\w+', name)
        matches = sum(1 for qw in query_words if any(qw in nw or nw in qw for nw in name_words) and len(qw) > 2)

        if matches > 0:
            # Get all relations
            relations = []
            for _, target, edge_data in G.out_edges(node_id, data=True):
                target_node = G.nodes.get(target, {})
                relations.append({
                    "relation": edge_data.get("relation", ""),
                    "target": target_node.get("name", target),
                    "target_type": target_node.get("type", ""),
                })
            for source, _, edge_data in G.in_edges(node_id, data=True):
                source_node = G.nodes.get(source, {})
                relations.append({
                    "relation": edge_data.get("relation", "") + " (ngược)",
                    "source": source_node.get("name", source),
                    "source_type": source_node.get("type", ""),
                })

            results.append({
                "entity": node.get("name", node_id),
                "type": node.get("type", ""),
                "match_score": matches,
                "attributes": {k: v for k, v in node.items() if k not in ("type", "name")},
                "relations": relations,
            })

    results.sort(key=lambda x: x["match_score"], reverse=True)
    return results[:5]


def build_kg_from_products(products: dict, solutions: dict):
    """Build knowledge graph from product/solution data."""
    G = _get_kg()

    # Add ViHAT Group as root
    add_entity("vihat_group", "company", "ViHAT Group",
               description="Công ty CPaaS hàng đầu Việt Nam")

    for pid, p in products.items():
        if p.get("status") == "deprecated":
            continue
        add_entity(f"product_{pid}", "product", p["name"],
                   slug=p["slug"],
                   description=p.get("short_description", ""),
                   features=p.get("features", []),
                   pricing=p.get("pricing_model", ""))
        add_relation(f"product_{pid}", "vihat_group", "thuộc về")

        # Add solutions
        for sid, s in solutions.items():
            if s.get("product_id") == pid and s.get("status") == "active":
                add_entity(f"solution_{sid}", "solution", s["name"],
                           slug=s["slug"],
                           description=s.get("description", ""),
                           aliases=s.get("aliases", []))
                add_relation(f"solution_{sid}", f"product_{pid}", "là giải pháp của")

    _save_kg()
    print(f"[KG] Built from products: {len(G.nodes)} nodes, {len(G.edges)} edges")


def build_kg_from_text(text: str, doc_id: str, anthropic_key: str = ""):
    """Extract entities and relations from text using Claude, add to KG."""
    if not anthropic_key or len(text) < 50:
        return

    import httpx

    prompt = f"""Phân tích đoạn text sau và trích xuất các thực thể (entity) và quan hệ (relation).

Text:
{text[:4000]}

Trả về JSON:
{{
  "entities": [
    {{"id": "unique_id", "type": "person|company|product|policy|position", "name": "Tên", "attributes": {{"key": "value"}}}}
  ],
  "relations": [
    {{"from_id": "id1", "to_id": "id2", "relation": "là CEO của|thuộc về|cung cấp|..."}}
  ]
}}

Chỉ trích xuất thông tin CÓ TRONG TEXT, KHÔNG bịa thêm. Nếu không có entity nào, trả {{"entities":[], "relations":[]}}."""

    try:
        resp = httpx.post(
            "https://api.anthropic.com/v1/messages",
            headers={
                "x-api-key": anthropic_key,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": "claude-sonnet-4-20250514",
                "max_tokens": 2000,
                "messages": [{"role": "user", "content": prompt}],
            },
            timeout=30,
        )
        content = resp.json().get("content", [{}])[0].get("text", "")
        # Extract JSON
        match = re.search(r'\{[\s\S]*\}', content)
        if match:
            data = json.loads(match.group())
            for ent in data.get("entities", []):
                add_entity(ent["id"], ent.get("type", ""), ent.get("name", ""),
                           **ent.get("attributes", {}), source_doc=doc_id)
            for rel in data.get("relations", []):
                add_relation(rel["from_id"], rel["to_id"], rel["relation"], source_doc=doc_id)
            _save_kg()
    except Exception as e:
        print(f"[KG] Entity extraction error: {e}")


# ═══════════════════════════════════════════════════════════════
# CHAT / ANSWER GENERATION
# ═══════════════════════════════════════════════════════════════

def generate_answer(query: str, anthropic_key: str, knowledge_base: str = "",
                    user_id: str = "", conversation_id: str = "") -> dict:
    """
    Full RAG pipeline: search → retrieve → generate answer.
    Returns {answer, sources, kg_results}
    """
    import httpx

    # 1. Knowledge Graph query (fast, exact)
    kg_results = query_kg(query)
    kg_context = ""
    if kg_results:
        kg_parts = []
        for r in kg_results:
            info = f"• {r['entity']} ({r['type']})"
            for attr_k, attr_v in r.get("attributes", {}).items():
                if attr_v and attr_k not in ("source_doc",):
                    info += f"\n  - {attr_k}: {attr_v}"
            for rel in r.get("relations", []):
                target = rel.get("target", rel.get("source", ""))
                info += f"\n  - {rel['relation']}: {target}"
            kg_parts.append(info)
        kg_context = "THÔNG TIN TỪ KNOWLEDGE GRAPH:\n" + "\n".join(kg_parts)

    # 2. Vector search (semantic + keyword)
    search_results = search(query, top_k=8, knowledge_base=knowledge_base)

    # 3. Build context from search results
    rag_context = ""
    sources = []
    if search_results:
        context_parts = []
        seen_docs = set()
        for r in search_results:
            if r["score"] < 0.25:
                continue
            context_parts.append(f"--- Nguồn: {r['title']} (score: {r['score']}) ---\n{r['original_text']}")
            if r["doc_id"] not in seen_docs:
                sources.append({
                    "title": r["title"],
                    "score": r["score"],
                    "drive_url": r.get("drive_url", ""),
                    "knowledge_base": r.get("knowledge_base", ""),
                })
                seen_docs.add(r["doc_id"])
        rag_context = "\n\n".join(context_parts)

    # 4. Generate answer with Claude
    full_context = ""
    if kg_context:
        full_context += kg_context + "\n\n"
    if rag_context:
        full_context += "THÔNG TIN TỪ KHO TRI THỨC:\n" + rag_context

    if not full_context.strip():
        return {
            "answer": "Em không tìm thấy thông tin liên quan trong kho tri thức. Anh/chị có thể mô tả rõ hơn câu hỏi không?",
            "sources": [],
            "kg_results": [],
        }

    system_prompt = """Bạn là Tiểu My — trợ lý kiến thức nội bộ của ViHAT Group.

NGUYÊN TẮC TUYỆT ĐỐI:
1. CHỈ trả lời dựa trên thông tin trong [context] bên dưới. TUYỆT ĐỐI KHÔNG bịa thêm bất kỳ thông tin nào.
2. Nếu thông tin không có trong context, nói rõ: "Em không tìm thấy thông tin này trong kho tri thức."
3. KHÔNG tự thêm học vấn, kinh nghiệm, số liệu, chức vụ cũ hay bất kỳ thông tin nào không có trong tài liệu.
4. Khi trích dẫn tài liệu, nếu có link thì tạo hyperlink markdown: [Tên tài liệu](link)
5. Trả lời bằng tiếng Việt, thân thiện, chuyên nghiệp.
6. Sử dụng emoji phù hợp để tạo cảm giác thân thiện.
7. Nếu có thông tin từ Knowledge Graph, ưu tiên sử dụng vì đây là thông tin chính xác nhất."""

    try:
        resp = httpx.post(
            "https://api.anthropic.com/v1/messages",
            headers={
                "x-api-key": anthropic_key,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": "claude-sonnet-4-20250514",
                "max_tokens": 2000,
                "system": system_prompt,
                "messages": [
                    {"role": "user", "content": f"Dựa trên thông tin sau:\n\n{full_context}\n\nCâu hỏi: {query}"}
                ],
            },
            timeout=60,
        )
        answer = resp.json().get("content", [{}])[0].get("text", "Có lỗi xảy ra.")
    except Exception as e:
        answer = f"Lỗi khi tạo câu trả lời: {str(e)}"

    return {
        "answer": answer,
        "sources": sources[:6],
        "kg_results": [{"entity": r["entity"], "type": r["type"]} for r in kg_results],
    }


# ═══════════════════════════════════════════════════════════════
# TEXT EXTRACTION (from various file types)
# ═══════════════════════════════════════════════════════════════

def extract_text_from_file(file_path: str) -> str:
    """Extract text from various file types."""
    path = Path(file_path)
    if not path.exists():
        return ""

    ext = path.suffix.lower()
    try:
        if ext == ".txt":
            return path.read_text(encoding="utf-8", errors="ignore")

        elif ext == ".docx":
            from docx import Document
            doc = Document(str(path))
            parts = []
            # Paragraphs
            for p in doc.paragraphs:
                if p.text.strip():
                    parts.append(p.text)
            # Tables (important! Dify misses these)
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            return "\n".join(parts)

        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                parts.append(f"=== Sheet: {sheet} ===")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        parts.append(" | ".join(cells))
            return "\n".join(parts)

        elif ext == ".pdf":
            try:
                from pdfminer.high_level import extract_text
                return extract_text(str(path))
            except ImportError:
                # Fallback: try PyPDF2
                try:
                    import PyPDF2
                    with open(str(path), "rb") as f:
                        reader = PyPDF2.PdfReader(f)
                        return "\n".join(page.extract_text() or "" for page in reader.pages)
                except Exception:
                    return ""

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            if p.text.strip():
                                parts.append(p.text)
            return "\n".join(parts)

        else:
            # Try as plain text
            return path.read_text(encoding="utf-8", errors="ignore")

    except Exception as e:
        print(f"[RAG] Extract text error for {path.name}: {e}")
        return ""


# ═══════════════════════════════════════════════════════════════
# INITIALIZATION
# ═══════════════════════════════════════════════════════════════

def initialize():
    """Initialize RAG engine (load index, KG)."""
    _load_chunk_index()
    _get_kg()
    print(f"[RAG] Initialized: {len(_chunk_index)} chunks in index, {len(_get_kg().nodes)} KG nodes")
