"""
Mock API Server for Admin Panel development.
No PostgreSQL/Redis/Dify required - pure in-memory mock data.

Usage:
  pip install fastapi uvicorn python-jose passlib[bcrypt]
  python dev_mock_server.py

Then open: http://localhost:3000 (Admin Panel)
API runs at: http://localhost:8000
"""

from __future__ import annotations

import os
import uuid
import random
import httpx
import json
from datetime import datetime, timedelta
from contextlib import asynccontextmanager
from typing import Optional, List

from fastapi import FastAPI, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

# Load .env
from pathlib import Path
env_path = Path(__file__).parent / ".env"
if env_path.exists():
    for line in env_path.read_text().splitlines():
        line = line.strip()
        if line and not line.startswith("#") and "=" in line:
            key, _, val = line.partition("=")
            key = key.strip()
            val = val.split("#")[0].strip()  # remove inline comments
            if key and val:
                # Use setdefault but override empty values
                if not os.environ.get(key):
                    os.environ[key] = val

DIFY_BASE_URL = os.environ.get("DIFY_BASE_URL", "http://103.29.27.91/v1")
DIFY_API_KEY = os.environ.get("DIFY_API_KEY_GENERAL", "")
DIFY_DATASET_API_KEY = os.environ.get("DIFY_DATASET_API_KEY", "")
# Public URL for file downloads (used in chatbot responses to external channels)
PUBLIC_API_URL = os.environ.get("PUBLIC_API_URL", "http://localhost:8000")
DIFY_DATASET_IDS = {
    "sales": os.environ.get("DIFY_DATASET_ID_SALES", ""),
    "hr": os.environ.get("DIFY_DATASET_ID_HR", ""),
    "accounting": os.environ.get("DIFY_DATASET_ID_ACCOUNTING", ""),
    "general": os.environ.get("DIFY_DATASET_ID_GENERAL", ""),
    "management": os.environ.get("DIFY_DATASET_ID_MANAGEMENT", ""),
}

# ---- Fake Auth ----
SECRET_KEY = "dev-secret-key"
ALGORITHM = "HS256"

try:
    from jose import jwt
    HAS_JWT = True
except ImportError:
    HAS_JWT = False
    print("WARNING: python-jose not installed. Auth will use simple token.")

# Simple password verification (no bcrypt dependency needed for dev)
import hashlib

def simple_hash(password: str) -> str:
    return hashlib.sha256(password.encode()).hexdigest()

def simple_verify(password: str, hashed: str) -> bool:
    return hashlib.sha256(password.encode()).hexdigest() == hashed


def create_token(user_id: str, role: str, tenant_id: str | None = None) -> str:
    tid = tenant_id or DEFAULT_TENANT_ID
    if HAS_JWT:
        return jwt.encode(
            {"sub": user_id, "role": role, "tenant_id": tid, "exp": datetime.utcnow() + timedelta(hours=24)},
            SECRET_KEY, algorithm=ALGORITHM
        )
    return f"mock-token-{user_id}"


def verify_token(token: str):
    if HAS_JWT:
        try:
            return jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        except Exception:
            return None
    if token.startswith("mock-token-"):
        user_id = token.replace("mock-token-", "")
        return {"sub": user_id, "role": "super_admin"}
    return None


# ---- Mock Data ----
ADMIN_ID = str(uuid.uuid4())

MOCK_USERS = [
    {
        "id": ADMIN_ID,
        "name": "Đinh Thái Hà",
        "email": "admin@vihat.vn",
        "department": "management",
        "role": "super_admin",
        "password_hash": simple_hash("vihat@2026"),
        "zalo_id": "zalo_admin_001",
        "telegram_id": 123456789,
        "knowledge_access": ["sales", "hr", "accounting", "general", "management"],
        "is_active": True,
        "created_at": "2026-03-01T09:00:00",
    },
    {
        "id": str(uuid.uuid4()),
        "name": "Nguyễn Văn Bình",
        "email": "binh.nguyen@vihat.vn",
        "department": "sales",
        "role": "lead",
        "password_hash": simple_hash("password123"),
        "zalo_id": "zalo_sales_001",
        "telegram_id": None,
        "knowledge_access": ["sales", "general"],
        "is_active": True,
        "created_at": "2026-03-05T10:00:00",
    },
    {
        "id": str(uuid.uuid4()),
        "name": "Trần Thị Cúc",
        "email": "cuc.tran@vihat.vn",
        "department": "hr",
        "role": "admin",
        "password_hash": simple_hash("password123"),
        "zalo_id": None,
        "telegram_id": 987654321,
        "knowledge_access": ["hr", "general"],
        "is_active": True,
        "created_at": "2026-03-03T08:30:00",
    },
    {
        "id": str(uuid.uuid4()),
        "name": "Lê Minh Đức",
        "email": "duc.le@vihat.vn",
        "department": "accounting",
        "role": "member",
        "password_hash": simple_hash("password123"),
        "zalo_id": "zalo_acct_001",
        "telegram_id": None,
        "knowledge_access": ["accounting", "general"],
        "is_active": True,
        "created_at": "2026-03-10T14:00:00",
    },
]

# ═══════════════════════════════════════════════════════════════════════════════
# TENANT & RICH PRODUCT MODEL (SaaS-Ready + Versioning)
# ═══════════════════════════════════════════════════════════════════════════════

DEFAULT_TENANT_ID = "t-vihat-001"

TENANTS: dict = {
    DEFAULT_TENANT_ID: {
        "id": DEFAULT_TENANT_ID,
        "slug": "vihat",
        "name": "ViHAT Group",
        "logo_url": None,
        "primary_color": "#1E40AF",
        "config": {
            "legal_entities": [
                {"id": "omijsc", "label": "OMIJSC", "template": "proposal_omijsc.pptx"},
                {"id": "vihat_solutions", "label": "ViHAT Solutions", "template": "proposal_vihat_solutions.pptx"},
                {"id": "vihat_group", "label": "ViHAT GROUP", "template": "proposal_vihat_group.pptx"},
            ],
            "departments": ["sales", "hr", "accounting", "general", "management"],
            "channels": ["zalo_oa", "telegram", "web_admin"],
            "features": {"proposals": True, "knowledge": True, "chat": True},
        },
        "is_active": True,
        "created_at": "2026-03-01T00:00:00",
        "updated_at": "2026-03-01T00:00:00",
    }
}

# Add tenant_id to all mock users
for _u in MOCK_USERS:
    _u["tenant_id"] = DEFAULT_TENANT_ID

# Rich Product data (persisted to JSON)
PRODUCTS: dict = {}  # product_id -> product dict
PRODUCT_VERSIONS: dict = {}  # product_id -> [version dicts]
PRODUCTS_JSON_PATH = Path(__file__).parent / "data" / "products.json"
PRODUCT_VERSIONS_JSON_PATH = Path(__file__).parent / "data" / "product_versions.json"

def _save_products():
    """Persist PRODUCTS to JSON."""
    try:
        PRODUCTS_JSON_PATH.write_text(json.dumps(PRODUCTS, ensure_ascii=False, indent=2, default=str), encoding="utf-8")
    except Exception as e:
        print(f"[Products] Save error: {e}")

def _save_product_versions():
    """Persist PRODUCT_VERSIONS to JSON."""
    try:
        PRODUCT_VERSIONS_JSON_PATH.write_text(json.dumps(PRODUCT_VERSIONS, ensure_ascii=False, indent=2, default=str), encoding="utf-8")
    except Exception as e:
        print(f"[Products] Save versions error: {e}")

def _init_products():
    """Load products from JSON file, or seed from defaults if no file exists."""
    global PRODUCTS, PRODUCT_VERSIONS

    # Try loading persisted data first
    if PRODUCTS_JSON_PATH.exists():
        try:
            PRODUCTS = json.loads(PRODUCTS_JSON_PATH.read_text(encoding="utf-8"))
            print(f"[Products] Loaded {len(PRODUCTS)} products from {PRODUCTS_JSON_PATH.name}")
        except Exception as e:
            print(f"[Products] Load error: {e}, will seed defaults")
            PRODUCTS = {}

    if PRODUCT_VERSIONS_JSON_PATH.exists():
        try:
            PRODUCT_VERSIONS = json.loads(PRODUCT_VERSIONS_JSON_PATH.read_text(encoding="utf-8"))
        except Exception:
            PRODUCT_VERSIONS = {}

    if PRODUCTS:
        return  # Already loaded from file

    now = datetime.utcnow().isoformat()

    seed_products = [
        {
            "slug": "tong_dai",
            "name": "Tổng đài",
            "short_description": "Giải pháp tổng đài IP đám mây cho doanh nghiệp",
            "full_description": "Hệ thống tổng đài IP đám mây (Cloud PBX/Contact Center) cho phép doanh nghiệp quản lý cuộc gọi inbound/outbound, hỗ trợ phân phối cuộc gọi thông minh ACD, IVR đa cấp, ghi âm, giám sát realtime, và tích hợp CRM.",
            "features": ["Cloud PBX", "ACD - Phân phối cuộc gọi thông minh", "IVR đa cấp", "Ghi âm cuộc gọi", "Giám sát realtime", "Báo cáo thống kê", "Click-to-call", "Softphone/WebRTC", "API tích hợp CRM"],
            "use_cases": ["Tổng đài CSKH inbound", "Telesales outbound", "Hotline doanh nghiệp", "Contact Center đa kênh"],
            "target_industries": ["chung", "ban_le", "y_te", "bat_dong_san", "fmcg"],
            "pricing_model": "Gói thuê bao theo số lượng agent/tháng. Có gói Starter, Business, Enterprise.",
            "competitive_advantages": ["Triển khai nhanh trong 24h", "Không cần đầu tư hạ tầng", "Tích hợp sẵn eSMS, Zalo OA", "SLA 99.9%", "AI routing thông minh"],
            "integration_options": ["CRM (Salesforce, HubSpot, Zoho)", "ERP", "Zalo OA", "Facebook Messenger", "Email"],
        },
        {
            "slug": "da_kenh",
            "name": "Đa kênh",
            "short_description": "Nền tảng chăm sóc khách hàng đa kênh hợp nhất",
            "full_description": "Giải pháp Omnichannel Contact Center hợp nhất tất cả kênh liên lạc (Thoại, Zalo, Facebook, Email, Live Chat, SMS) vào 1 giao diện duy nhất. Agent có thể xử lý đồng thời nhiều kênh, lịch sử tương tác khách hàng được lưu trữ xuyên suốt.",
            "features": ["Hợp nhất đa kênh", "Lịch sử tương tác 360°", "Phân phối ticket thông minh", "SLA management", "Báo cáo đa kênh", "Chatbot tích hợp", "Knowledge Base nội bộ"],
            "use_cases": ["CSKH đa kênh", "Help Desk nội bộ", "Social Commerce support", "Quản lý ticket"],
            "target_industries": ["chung", "ban_le", "thoi_trang", "fmcg"],
            "pricing_model": "Gói theo số agent + số kênh kết nối. Thanh toán tháng/năm.",
            "competitive_advantages": ["1 giao diện cho tất cả kênh", "Không bỏ lỡ tin nhắn", "Tích hợp Zalo OA chính thức", "AI auto-reply"],
            "integration_options": ["Zalo OA", "Facebook Page", "Instagram", "Email SMTP", "Website Live Chat", "CRM"],
        },
        {
            "slug": "p_zalo",
            "name": "Zalo cá nhân (P-Zalo)",
            "short_description": "Quản lý bán hàng qua Zalo cá nhân tập trung",
            "full_description": "Giải pháp quản lý tập trung các tài khoản Zalo cá nhân của đội ngũ sales. Giám sát hội thoại, phân tích hiệu suất, tự động phân phối lead, và đảm bảo data khách hàng thuộc về doanh nghiệp.",
            "features": ["Quản lý nhiều Zalo cá nhân", "Giám sát hội thoại", "Phân phối lead tự động", "Báo cáo hiệu suất sale", "Backup dữ liệu chat", "Template tin nhắn"],
            "use_cases": ["Quản lý đội sales qua Zalo", "Bán hàng online", "CSKH cá nhân hóa"],
            "target_industries": ["ban_le", "lam_dep", "thoi_trang", "bat_dong_san"],
            "pricing_model": "Gói theo số tài khoản Zalo quản lý.",
            "competitive_advantages": ["Không mất data khi sale nghỉ việc", "Giám sát chất lượng CSKH", "Tự động gán khách hàng mới"],
            "integration_options": ["CRM", "Zalo OA", "Hệ thống ERP"],
        },
        {
            "slug": "sms_brandname",
            "name": "SMS Brandname",
            "short_description": "Gửi SMS thương hiệu hàng loạt với tên doanh nghiệp",
            "full_description": "Dịch vụ gửi tin nhắn SMS hiển thị tên thương hiệu doanh nghiệp (không hiện số điện thoại). Hỗ trợ SMS CSKH, SMS Marketing, OTP, và thông báo giao dịch với tỉ lệ gửi thành công cao.",
            "features": ["SMS Brandname CSKH", "SMS Marketing", "SMS OTP", "API gửi tin tự động", "Lập lịch gửi", "Báo cáo delivery", "Template quản lý"],
            "use_cases": ["Gửi OTP xác thực", "Thông báo giao dịch", "Chăm sóc khách hàng", "Marketing promotion"],
            "target_industries": ["chung", "ban_le", "y_te", "fmcg", "bat_dong_san"],
            "pricing_model": "Trả theo số lượng tin nhắn gửi thành công. Giá theo đầu số và nhà mạng.",
            "competitive_advantages": ["Kết nối trực tiếp nhà mạng", "Tỉ lệ gửi thành công >98%", "API đơn giản", "Hỗ trợ Unicode tiếng Việt"],
            "integration_options": ["API RESTful", "SMPP", "Webhook callback", "CRM", "ERP"],
        },
        {
            "slug": "zbs",
            "name": "ZBS",
            "short_description": "Zalo Business Solution - Gửi thông báo qua Zalo OA",
            "full_description": "Giải pháp gửi tin nhắn ZNS (Zalo Notification Service) qua Zalo Official Account. Gửi thông báo giao dịch, CSKH, remarketing tới khách hàng qua Zalo với tỉ lệ đọc cao.",
            "features": ["ZNS Template Message", "Zalo OA Follower management", "Broadcast message", "API tích hợp", "Báo cáo tương tác", "Mini App integration"],
            "use_cases": ["Thông báo đơn hàng", "Nhắc lịch hẹn", "CSKH qua Zalo", "Remarketing"],
            "target_industries": ["chung", "ban_le", "lam_dep", "suc_khoe", "y_te"],
            "pricing_model": "Theo số ZNS gửi thành công + phí duy trì Zalo OA.",
            "competitive_advantages": ["Tỉ lệ đọc >90%", "Chi phí thấp hơn SMS", "Rich media (hình ảnh, nút bấm)", "Tích hợp Mini App"],
            "integration_options": ["API RESTful", "Zalo OA SDK", "Webhook", "CRM"],
        },
        {
            "slug": "topup",
            "name": "Topup",
            "short_description": "Nạp tiền điện thoại/data tự động cho khách hàng",
            "full_description": "Dịch vụ nạp tiền điện thoại, data mobile tự động qua API. Dùng cho chương trình loyalty, khuyến mãi, hoàn tiền, tặng quà khách hàng bằng giá trị nạp tiền điện thoại.",
            "features": ["Nạp tiền tự động qua API", "Hỗ trợ tất cả nhà mạng VN", "Nạp data mobile", "Báo cáo realtime", "Webhook callback"],
            "use_cases": ["Chương trình loyalty", "Khuyến mãi nạp thẻ", "Hoàn tiền cho KH", "Thưởng nhân viên"],
            "target_industries": ["ban_le", "fmcg", "chung"],
            "pricing_model": "Chiết khấu theo tổng giá trị nạp/tháng.",
            "competitive_advantages": ["Kết nối trực tiếp nhà mạng", "Nạp tức thì <5s", "API đơn giản", "Không giới hạn số lượng"],
            "integration_options": ["API RESTful", "Webhook callback", "CRM", "Loyalty platform"],
        },
        {
            "slug": "miniapp",
            "name": "MiniApp",
            "short_description": "Xây dựng Mini App trên Zalo cho doanh nghiệp",
            "full_description": "Dịch vụ phát triển Mini App chạy trên nền tảng Zalo. Cho phép doanh nghiệp tạo ứng dụng nhỏ gọn trong Zalo để bán hàng, đặt lịch, tích điểm, thanh toán mà không cần khách hàng cài app riêng.",
            "features": ["E-commerce Mini App", "Booking/đặt lịch", "Loyalty/tích điểm", "Thanh toán ZaloPay", "Push notification", "Quản lý thành viên"],
            "use_cases": ["Bán hàng online trên Zalo", "Đặt lịch dịch vụ", "Chương trình khách hàng thân thiết", "Menu nhà hàng"],
            "target_industries": ["ban_le", "lam_dep", "suc_khoe", "thoi_trang"],
            "pricing_model": "Phí phát triển + phí duy trì hàng tháng.",
            "competitive_advantages": ["Không cần cài app riêng", "75M+ người dùng Zalo", "Tích hợp thanh toán ZaloPay", "UX native trên Zalo"],
            "integration_options": ["Zalo OA", "ZaloPay", "API backend", "CRM"],
        },
        {
            "slug": "callbot_ai",
            "name": "Callbot AI",
            "short_description": "Robot gọi điện tự động bằng AI",
            "full_description": "Giải pháp Callbot sử dụng AI (NLP, TTS, ASR) để thực hiện cuộc gọi tự động. Callbot có thể gọi hàng ngàn cuộc/ngày cho các kịch bản: nhắc nợ, xác nhận đơn hàng, khảo sát, telesales, nhắc lịch hẹn.",
            "features": ["Gọi tự động hàng ngàn cuộc/ngày", "Nhận diện giọng nói (ASR)", "Tổng hợp giọng nói tự nhiên (TTS)", "Kịch bản linh hoạt (drag-drop)", "Chuyển agent khi cần", "Báo cáo kết quả chi tiết"],
            "use_cases": ["Nhắc nợ/thanh toán", "Xác nhận đơn hàng", "Khảo sát khách hàng", "Telesales tự động", "Nhắc lịch hẹn"],
            "target_industries": ["chung", "ban_le", "y_te", "bat_dong_san", "fmcg"],
            "pricing_model": "Theo số phút gọi thành công + phí setup kịch bản.",
            "competitive_advantages": ["Giọng nói tự nhiên như người thật", "Xử lý 10,000+ cuộc gọi/ngày", "Tỉ lệ hoàn thành >85%", "Tiết kiệm 70% chi phí nhân sự"],
            "integration_options": ["API RESTful", "CRM", "Tổng đài OmiCall", "Webhook callback"],
        },
        {
            "slug": "chatbot_ai",
            "name": "Chatbot AI",
            "short_description": "Chatbot AI tự động trả lời khách hàng đa kênh",
            "full_description": "Giải pháp Chatbot AI sử dụng NLP/LLM để tự động trả lời khách hàng trên nhiều kênh (Zalo, Facebook, Website, App). Chatbot học từ knowledge base doanh nghiệp, hỗ trợ 24/7, và chuyển agent khi cần.",
            "features": ["AI NLP/LLM hiểu ngữ cảnh", "Đa kênh (Zalo, FB, Web, App)", "Học từ Knowledge Base", "Chuyển agent thông minh", "Phân tích intent", "Đa ngôn ngữ", "Analytics dashboard"],
            "use_cases": ["CSKH tự động 24/7", "FAQ automation", "Lead qualification", "Đặt hàng qua chat", "Hỗ trợ kỹ thuật"],
            "target_industries": ["chung", "ban_le", "lam_dep", "thoi_trang", "fmcg"],
            "pricing_model": "Gói theo số lượng conversation/tháng + phí training model.",
            "competitive_advantages": ["Hiểu tiếng Việt tự nhiên", "Học từ dữ liệu doanh nghiệp", "Tích hợp sẵn OmiCall", "Chuyển đổi liền mạch Bot→Agent"],
            "integration_options": ["Zalo OA", "Facebook Messenger", "Website Widget", "API", "CRM"],
        },
    ]

    for i, p in enumerate(seed_products):
        pid = str(uuid.uuid4())[:8]
        PRODUCTS[pid] = {
            "id": pid,
            "tenant_id": DEFAULT_TENANT_ID,
            "slug": p["slug"],
            "name": p["name"],
            "short_description": p.get("short_description", ""),
            "full_description": p.get("full_description", ""),
            "features": p.get("features", []),
            "use_cases": p.get("use_cases", []),
            "target_industries": p.get("target_industries", []),
            "pricing_model": p.get("pricing_model", ""),
            "competitive_advantages": p.get("competitive_advantages", []),
            "integration_options": p.get("integration_options", []),
            "status": "active",
            "sort_order": i,
            "created_at": now,
            "updated_at": now,
        }
        PRODUCT_VERSIONS[pid] = []  # no history yet for initial seed

    _save_products()
    print(f"[Products] Seeded {len(PRODUCTS)} default products")

_init_products()

# ═══════════════════════════════════════════════════════════════════════════════
# SOLUTIONS (customer-facing names/use-cases linking to Products)
# ═══════════════════════════════════════════════════════════════════════════════
SOLUTIONS: dict = {}  # solution_id -> solution dict
SOLUTIONS_JSON_PATH = Path(__file__).parent / "data" / "solutions.json"

def _save_solutions():
    """Persist SOLUTIONS to JSON."""
    try:
        SOLUTIONS_JSON_PATH.write_text(json.dumps(SOLUTIONS, ensure_ascii=False, indent=2, default=str), encoding="utf-8")
    except Exception as e:
        print(f"[Solutions] Save error: {e}")

def _init_solutions():
    """Load solutions from JSON file, or seed defaults if no file exists."""
    global SOLUTIONS

    if SOLUTIONS_JSON_PATH.exists():
        try:
            SOLUTIONS = json.loads(SOLUTIONS_JSON_PATH.read_text(encoding="utf-8"))
            print(f"[Solutions] Loaded {len(SOLUTIONS)} solutions from {SOLUTIONS_JSON_PATH.name}")
            return
        except Exception as e:
            print(f"[Solutions] Load error: {e}, will seed defaults")
            SOLUTIONS = {}

    # Seed default solutions by mapping old product slugs
    now = datetime.utcnow().isoformat()

    # Helper: find product by slug
    def _find_product_id(slug: str) -> str:
        for pid, p in PRODUCTS.items():
            if p.get("slug") == slug:
                return pid
        return ""

    # First, ensure OMICall and eSMS products exist (merge from old tong_dai/da_kenh, sms_brandname/zbs)
    omicall_id = _find_product_id("omicall")
    if not omicall_id:
        # Create OMICall product by merging tong_dai + da_kenh info
        tong_dai = next((p for p in PRODUCTS.values() if p.get("slug") == "tong_dai"), None)
        da_kenh = next((p for p in PRODUCTS.values() if p.get("slug") == "da_kenh"), None)
        omicall_id = str(uuid.uuid4())[:8]
        merged_features = list(dict.fromkeys(
            (tong_dai or {}).get("features", []) + (da_kenh or {}).get("features", [])
        ))
        merged_use_cases = list(dict.fromkeys(
            (tong_dai or {}).get("use_cases", []) + (da_kenh or {}).get("use_cases", [])
        ))
        merged_industries = list(dict.fromkeys(
            (tong_dai or {}).get("target_industries", []) + (da_kenh or {}).get("target_industries", [])
        ))
        merged_advantages = list(dict.fromkeys(
            (tong_dai or {}).get("competitive_advantages", []) + (da_kenh or {}).get("competitive_advantages", [])
        ))
        merged_integrations = list(dict.fromkeys(
            (tong_dai or {}).get("integration_options", []) + (da_kenh or {}).get("integration_options", [])
        ))
        PRODUCTS[omicall_id] = {
            "id": omicall_id,
            "tenant_id": DEFAULT_TENANT_ID,
            "slug": "omicall",
            "name": "OMICall",
            "short_description": "Nền tảng tổng đài đa kênh thông minh cho doanh nghiệp",
            "full_description": "OMICall là nền tảng tổng đài đa kênh (Omnichannel Contact Center) của ViHAT, tích hợp thoại, Zalo, Facebook, Email, Live Chat vào 1 giao diện. Hỗ trợ Cloud PBX, ACD, IVR, ghi âm, giám sát realtime, và tích hợp CRM.",
            "features": merged_features,
            "use_cases": merged_use_cases,
            "target_industries": merged_industries,
            "pricing_model": "Gói thuê bao theo số lượng agent/tháng. Có gói Starter, Business, Enterprise.",
            "competitive_advantages": merged_advantages,
            "integration_options": merged_integrations,
            "status": "active",
            "sort_order": 0,
            "created_at": now,
            "updated_at": now,
        }
        PRODUCT_VERSIONS[omicall_id] = []
        # Remove old tong_dai and da_kenh products
        for pid in list(PRODUCTS.keys()):
            if PRODUCTS[pid].get("slug") in ("tong_dai", "da_kenh"):
                del PRODUCTS[pid]
                PRODUCT_VERSIONS.pop(pid, None)
        _save_products()
        print("[Solutions] Created OMICall product (merged from tong_dai + da_kenh)")

    esms_id = _find_product_id("esms")
    if not esms_id:
        # Create eSMS product by merging sms_brandname + zbs info
        sms = next((p for p in PRODUCTS.values() if p.get("slug") == "sms_brandname"), None)
        zbs = next((p for p in PRODUCTS.values() if p.get("slug") == "zbs"), None)
        esms_id = str(uuid.uuid4())[:8]
        merged_features = list(dict.fromkeys(
            (sms or {}).get("features", []) + (zbs or {}).get("features", [])
        ))
        merged_use_cases = list(dict.fromkeys(
            (sms or {}).get("use_cases", []) + (zbs or {}).get("use_cases", [])
        ))
        merged_industries = list(dict.fromkeys(
            (sms or {}).get("target_industries", []) + (zbs or {}).get("target_industries", [])
        ))
        merged_advantages = list(dict.fromkeys(
            (sms or {}).get("competitive_advantages", []) + (zbs or {}).get("competitive_advantages", [])
        ))
        merged_integrations = list(dict.fromkeys(
            (sms or {}).get("integration_options", []) + (zbs or {}).get("integration_options", [])
        ))
        PRODUCTS[esms_id] = {
            "id": esms_id,
            "tenant_id": DEFAULT_TENANT_ID,
            "slug": "esms",
            "name": "eSMS",
            "short_description": "Nền tảng gửi tin nhắn đa kênh: SMS Brandname, ZNS, Zalo OA",
            "full_description": "eSMS là nền tảng nhắn tin của ViHAT, hỗ trợ gửi SMS Brandname (CSKH, Marketing, OTP), ZNS qua Zalo OA, và thông báo giao dịch. Kết nối trực tiếp nhà mạng, tỉ lệ gửi thành công cao.",
            "features": merged_features,
            "use_cases": merged_use_cases,
            "target_industries": merged_industries,
            "pricing_model": "Trả theo số lượng tin nhắn gửi thành công.",
            "competitive_advantages": merged_advantages,
            "integration_options": merged_integrations,
            "status": "active",
            "sort_order": 1,
            "created_at": now,
            "updated_at": now,
        }
        PRODUCT_VERSIONS[esms_id] = []
        # Remove old sms_brandname and zbs products
        for pid in list(PRODUCTS.keys()):
            if PRODUCTS[pid].get("slug") in ("sms_brandname", "zbs"):
                del PRODUCTS[pid]
                PRODUCT_VERSIONS.pop(pid, None)
        _save_products()
        print("[Solutions] Created eSMS product (merged from sms_brandname + zbs)")

    # Now seed solutions
    seed_solutions = [
        {"name": "Tổng đài ảo", "slug": "tong_dai", "description": "Giải pháp tổng đài IP đám mây (Cloud PBX) cho doanh nghiệp", "product_slug": "omicall", "aliases": ["cloud pbx", "tổng đài ip", "tổng đài cloud", "virtual pbx"], "sort_order": 0},
        {"name": "Đa kênh / Omnichannel", "slug": "da_kenh", "description": "Giải pháp chăm sóc khách hàng đa kênh hợp nhất", "product_slug": "omicall", "aliases": ["omnichannel", "contact center", "đa kênh", "tổng đài đa kênh"], "sort_order": 1},
        {"name": "SMS Brandname", "slug": "sms_brandname", "description": "Gửi SMS thương hiệu hàng loạt với tên doanh nghiệp", "product_slug": "esms", "aliases": ["tin nhắn thương hiệu", "sms otp", "sms marketing"], "sort_order": 2},
        {"name": "ZNS / Zalo Business", "slug": "zbs", "description": "Gửi thông báo ZNS qua Zalo Official Account", "product_slug": "esms", "aliases": ["zalo notification", "zns", "zalo oa", "zalo business solution"], "sort_order": 3},
        {"name": "Quản lý Zalo cá nhân", "slug": "p_zalo", "description": "Quản lý bán hàng qua Zalo cá nhân tập trung", "product_slug": "p_zalo", "aliases": ["p-zalo", "zalo sales"], "sort_order": 4},
        {"name": "Nạp tiền tự động", "slug": "topup", "description": "Nạp tiền điện thoại/data tự động cho khách hàng", "product_slug": "topup", "aliases": ["nạp thẻ", "topup mobile"], "sort_order": 5},
        {"name": "Mini App Zalo", "slug": "miniapp", "description": "Xây dựng Mini App trên Zalo cho doanh nghiệp", "product_slug": "miniapp", "aliases": ["zalo miniapp", "mini app"], "sort_order": 6},
        {"name": "Robot gọi điện AI", "slug": "callbot_ai", "description": "Robot gọi điện tự động bằng AI cho nhắc nợ, telesales, khảo sát", "product_slug": "callbot_ai", "aliases": ["callbot", "auto call", "ai calling"], "sort_order": 7},
        {"name": "Chatbot tự động", "slug": "chatbot_ai", "description": "Chatbot AI tự động trả lời khách hàng đa kênh 24/7", "product_slug": "chatbot_ai", "aliases": ["chatbot", "ai chat", "trả lời tự động"], "sort_order": 8},
    ]

    for s in seed_solutions:
        sid = str(uuid.uuid4())[:8]
        product_id = _find_product_id(s["product_slug"])
        if not product_id:
            print(f"[Solutions] Warning: product '{s['product_slug']}' not found for solution '{s['name']}'")
            continue
        SOLUTIONS[sid] = {
            "id": sid,
            "tenant_id": DEFAULT_TENANT_ID,
            "name": s["name"],
            "slug": s["slug"],
            "description": s["description"],
            "product_id": product_id,
            "aliases": s.get("aliases", []),
            "status": "active",
            "sort_order": s.get("sort_order", 0),
            "created_at": now,
            "updated_at": now,
        }

    _save_solutions()
    print(f"[Solutions] Seeded {len(SOLUTIONS)} default solutions")

_init_solutions()

# ═══════════════════════════════════════════════════════════════════════════════
# VERSION SERVICE (reusable for products, rfi_templates, etc.)
# ═══════════════════════════════════════════════════════════════════════════════

def save_version(entity_id: str, versions_store: dict, current_snapshot: dict,
                 changed_by: str = "", version_label: str = "", change_summary: str = ""):
    """Save current state as a version before updating."""
    versions = versions_store.setdefault(entity_id, [])
    version_number = len(versions) + 1
    versions.append({
        "id": str(uuid.uuid4())[:8],
        "version_number": version_number,
        "version_label": version_label or f"v{version_number}",
        "changed_by": changed_by,
        "change_summary": change_summary,
        "snapshot": {k: v for k, v in current_snapshot.items() if k not in ("id", "tenant_id", "created_at")},
        "created_at": datetime.utcnow().isoformat(),
    })
    return version_number


SAMPLE_QUERIES = [
    ("Cho tôi bảng giá OmiCall Enterprise", "sales", "zalo_oa"),
    ("Quy trình xin nghỉ phép dài ngày", "hr", "telegram"),
    ("Cách xuất hóa đơn GTGT cho khách hàng nước ngoài", "accounting", "web_admin"),
    ("So sánh OmiCall với Stringee", "sales", "zalo_oa"),
    ("Chính sách bảo hiểm xã hội cho nhân viên thử việc", "hr", "telegram"),
    ("Quy trình thanh toán công nợ", "accounting", "zalo_oa"),
    ("Hướng dẫn dùng SMS Brandname", "sales", "telegram"),
    ("Thời gian nghỉ lễ 30/4 và 1/5", "general", "zalo_oa"),
    ("Giá gói Zalo ZNS Premium", "sales", "web_admin"),
    ("Quy trình onboarding nhân viên mới", "hr", "zalo_oa"),
]

SAMPLE_ANSWERS = [
    "Gói OmiCall Enterprise có giá từ 5.000.000đ/tháng, hỗ trợ tối đa 50 agents, tích hợp CRM, recording không giới hạn.",
    "Quy trình nghỉ phép dài ngày (>3 ngày): 1) Gửi đơn qua HRM trước 7 ngày. 2) Trưởng phòng phê duyệt. 3) HR xác nhận.",
    "Để xuất hóa đơn GTGT cho KH nước ngoài: Vào hệ thống EInvoice → Chọn 'Hóa đơn xuất khẩu' → Điền thông tin theo mẫu.",
    "So sánh OmiCall vs Stringee: OmiCall hỗ trợ AI routing, giá cạnh tranh hơn 20%, tích hợp sẵn với eSMS và Zalo OA.",
    "NV thử việc được đóng BHXH sau khi ký HĐLĐ chính thức. Trong thời gian thử việc, công ty đóng BHYT bắt buộc.",
]


def generate_mock_logs(count: int = 50) -> list:
    logs = []
    for i in range(count):
        q_idx = i % len(SAMPLE_QUERIES)
        query, dept, channel = SAMPLE_QUERIES[q_idx]
        a_idx = i % len(SAMPLE_ANSWERS)
        user = random.choice(MOCK_USERS)
        created = datetime.utcnow() - timedelta(hours=random.randint(1, 720))

        logs.append({
            "id": str(uuid.uuid4()),
            "user_id": user["id"],
            "user_name": user["name"],
            "channel": channel,
            "query_text": query,
            "answer_text": SAMPLE_ANSWERS[a_idx],
            "department_routed": dept,
            "sources": [
                {"document": f"doc_{dept}_{random.randint(1,10)}.pdf", "chunk": "...", "score": round(random.uniform(0.7, 0.98), 2)},
            ],
            "confidence_score": round(random.uniform(0.65, 0.98), 2),
            "tokens_prompt": random.randint(800, 2000),
            "tokens_completion": random.randint(200, 600),
            "processing_time_ms": random.randint(1200, 4500),
            "feedback_rating": random.choice([None, None, None, 4, 5]),
            "created_at": created.isoformat(),
        })
    logs.sort(key=lambda x: x["created_at"], reverse=True)
    return logs


MOCK_LOGS = generate_mock_logs(50)

MOCK_DOCUMENTS = [
    {
        "id": str(uuid.uuid4()),
        "knowledge_base": "sales",
        "title": "Bảng giá OmiCall Q1 2026",
        "file_name": "pricing_omicall_q1_2026.xlsx",
        "file_type": "xlsx",
        "file_size_bytes": 245760,
        "tags": ["omicall", "pricing"],
        "chunks_count": 24,
        "status": "ready",
        "created_at": "2026-03-01T10:00:00",
    },
    {
        "id": str(uuid.uuid4()),
        "knowledge_base": "sales",
        "title": "Catalog sản phẩm eSMS 2026",
        "file_name": "esms_catalog_2026.pdf",
        "file_type": "pdf",
        "file_size_bytes": 1548000,
        "tags": ["esms", "catalog"],
        "chunks_count": 42,
        "status": "ready",
        "created_at": "2026-03-02T11:00:00",
    },
    {
        "id": str(uuid.uuid4()),
        "knowledge_base": "hr",
        "title": "Nội quy lao động ViHAT Group 2026",
        "file_name": "noi_quy_lao_dong_2026.docx",
        "file_type": "docx",
        "file_size_bytes": 385024,
        "tags": ["noi-quy", "lao-dong"],
        "chunks_count": 18,
        "status": "ready",
        "created_at": "2026-03-03T09:00:00",
    },
    {
        "id": str(uuid.uuid4()),
        "knowledge_base": "accounting",
        "title": "Quy trình xuất hóa đơn GTGT",
        "file_name": "quy_trinh_hoa_don.pdf",
        "file_type": "pdf",
        "file_size_bytes": 520000,
        "tags": ["hoa-don", "thue"],
        "chunks_count": 15,
        "status": "ready",
        "created_at": "2026-03-05T14:00:00",
    },
    {
        "id": str(uuid.uuid4()),
        "knowledge_base": "general",
        "title": "SOP liên phòng ban - Phối hợp Sales-Kế toán",
        "file_name": "sop_lien_phong_ban.pdf",
        "file_type": "pdf",
        "file_size_bytes": 290000,
        "tags": ["sop", "lien-phong-ban"],
        "chunks_count": 10,
        "status": "ready",
        "created_at": "2026-03-07T16:00:00",
    },
]

# ---- App ----

@asynccontextmanager
async def lifespan(app: FastAPI):
    print("\n" + "=" * 60)
    print("  ViHAT Knowledge System — Mock Dev Server")
    print("=" * 60)
    print(f"  API:   http://localhost:8000")
    print(f"  Docs:  http://localhost:8000/docs")
    print(f"  Login: admin@vihat.vn / vihat@2026")
    print("=" * 60 + "\n")
    yield


app = FastAPI(
    title="ViHAT Knowledge System (Mock Dev)",
    version="dev",
    lifespan=lifespan,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def get_current_user(request: Request) -> dict:
    auth = request.headers.get("Authorization", "")
    if not auth.startswith("Bearer "):
        raise HTTPException(401, "Not authenticated")
    token = auth.split(" ", 1)[1]
    payload = verify_token(token)
    if not payload:
        raise HTTPException(401, "Invalid token")
    user_id = payload["sub"]
    for u in MOCK_USERS:
        if u["id"] == user_id:
            return u
    raise HTTPException(401, "User not found")


# --- Auth ---

class LoginRequest(BaseModel):
    email: str
    password: str


@app.post("/api/v1/auth/login")
async def login(req: LoginRequest):
    for user in MOCK_USERS:
        if user["email"] == req.email:
            if simple_verify(req.password, user["password_hash"]):
                token = create_token(user["id"], user["role"])
                return {
                    "access_token": token,
                    "token_type": "bearer",
                    "user": {k: v for k, v in user.items() if k != "password_hash"},
                }
    raise HTTPException(401, "Email hoặc mật khẩu không đúng")


@app.get("/api/v1/auth/me")
async def get_me(request: Request):
    user = get_current_user(request)
    return {k: v for k, v in user.items() if k != "password_hash"}


# --- Stats ---

@app.get("/api/v1/admin/stats")
async def get_stats(days: int = 30, request: Request = None):
    return {
        "period": f"{(datetime.utcnow() - timedelta(days=days)).date()} to {datetime.utcnow().date()}",
        "total_queries": 3542,
        "by_department": {"sales": 1890, "hr": 620, "accounting": 480, "general": 552},
        "by_channel": {"zalo_oa": 2100, "telegram": 1200, "web_admin": 242},
        "avg_response_time_ms": 2150,
        "avg_confidence_score": 0.87,
        "tokens_used": {"total": 2450000, "prompt": 1680000, "completion": 770000},
    }


# --- Users ---

@app.get("/api/v1/admin/users")
async def list_users(department: Optional[str] = None):
    users = MOCK_USERS
    if department:
        users = [u for u in users if u["department"] == department]
    return [{k: v for k, v in u.items() if k != "password_hash"} for u in users]


class UserCreate(BaseModel):
    name: str
    email: Optional[str] = None
    department: str
    role: str = "member"
    zalo_id: Optional[str] = None
    telegram_id: Optional[int] = None
    knowledge_access: List[str] = []


@app.post("/api/v1/admin/users")
async def create_user(data: UserCreate):
    new_user = {
        "id": str(uuid.uuid4()),
        "name": data.name,
        "email": data.email,
        "department": data.department,
        "role": data.role,
        "zalo_id": data.zalo_id,
        "telegram_id": data.telegram_id,
        "knowledge_access": data.knowledge_access,
        "is_active": True,
        "created_at": datetime.utcnow().isoformat(),
    }
    MOCK_USERS.append(new_user)
    return new_user


@app.put("/api/v1/admin/users/{user_id}")
async def update_user(user_id: str, data: dict):
    for u in MOCK_USERS:
        if u["id"] == user_id:
            for k, v in data.items():
                if k != "id" and k != "password_hash":
                    u[k] = v
            return {k: v for k, v in u.items() if k != "password_hash"}
    raise HTTPException(404, "User not found")


@app.delete("/api/v1/admin/users/{user_id}")
async def delete_user(user_id: str):
    for u in MOCK_USERS:
        if u["id"] == user_id:
            u["is_active"] = False
            return {"status": "ok", "message": f"User {u['name']} deactivated"}
    raise HTTPException(404, "User not found")


# --- Knowledge ---

KB_NAME_MAP = {
    "sales": "ViHAT Sales Knowledge",
    "hr": "ViHAT HR Knowledge",
    "accounting": "ViHAT Accounting Knowledge",
    "general": "ViHAT General Knowledge",
    "management": "ViHAT Management Knowledge",
}
KB_NAME_REVERSE = {v: k for k, v in KB_NAME_MAP.items()}

# --- Dify status cache (avoid fetching every request) ---
_dify_status_cache: dict = {}  # doc_id → {kb, status, name, tokens}
_dify_status_cache_ts: float = 0  # timestamp of last fetch
_DIFY_CACHE_TTL = 30  # seconds — refresh every 30s


async def _refresh_dify_status_cache(knowledge_base: Optional[str] = None):
    """Fetch Dify doc statuses with pagination and cache them."""
    global _dify_status_cache, _dify_status_cache_ts
    import time

    now = time.time()
    if now - _dify_status_cache_ts < _DIFY_CACHE_TTL and _dify_status_cache:
        return _dify_status_cache

    datasets_to_check = {}
    if knowledge_base:
        ds_id = DIFY_DATASET_IDS.get(knowledge_base, "")
        if ds_id:
            datasets_to_check[knowledge_base] = ds_id
    else:
        datasets_to_check = {k: v for k, v in DIFY_DATASET_IDS.items() if v}

    new_cache = {}
    async with httpx.AsyncClient(timeout=30.0) as client:
        for kb_key, ds_id in datasets_to_check.items():
            try:
                page = 1
                while True:
                    resp = await client.get(
                        f"{DIFY_BASE_URL}/datasets/{ds_id}/documents",
                        headers={"Authorization": f"Bearer {DIFY_DATASET_API_KEY}"},
                        params={"page": page, "limit": 100},
                    )
                    resp.raise_for_status()
                    data = resp.json().get("data", [])
                    if not data:
                        break
                    for doc in data:
                        new_cache[doc["id"]] = {
                            "kb": kb_key,
                            "status": doc.get("indexing_status", "unknown"),
                            "name": doc.get("name", ""),
                            "tokens": doc.get("tokens", 0),
                        }
                    if len(data) < 100:
                        break
                    page += 1
            except Exception as e:
                print(f"Error fetching docs from {kb_key}: {e}")

    _dify_status_cache = new_cache
    _dify_status_cache_ts = now
    return new_cache


@app.get("/api/v1/admin/knowledge/list")
async def list_knowledge(knowledge_base: Optional[str] = None, status: Optional[str] = None):
    """List parent documents (grouped). Each item = 1 original file with section count."""
    if not DIFY_DATASET_API_KEY:
        return []

    # Use cached Dify statuses (refreshes every 30s)
    all_dify_docs = await _refresh_dify_status_cache(knowledge_base)

    # Build parent document list from registry
    docs = []
    for parent_id, entry in FILE_REGISTRY.items():
        kb = entry.get("knowledge_base", "")
        if knowledge_base and kb != knowledge_base:
            continue

        section_ids = entry.get("section_doc_ids", [])

        # Determine overall status from sections
        section_statuses = [all_dify_docs.get(sid, {}).get("status", "unknown") for sid in section_ids]
        if not section_ids:
            # No section IDs tracked → treat as ready (legacy data or upload without tracking)
            overall_status = "ready"
        elif all(s == "completed" for s in section_statuses) and section_statuses:
            overall_status = "ready"
        elif any(s in ("indexing", "splitting", "parsing", "waiting") for s in section_statuses):
            overall_status = "indexing"
        elif any(s == "error" for s in section_statuses):
            overall_status = "error"
        elif all(s == "unknown" for s in section_statuses):
            # All sections not found in Dify → likely completed & archived, treat as ready
            overall_status = "ready"
        else:
            overall_status = "ready"

        file_name = entry.get("file_name", "")
        ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else "txt"

        source_id = entry.get("source_id")
        docs.append({
            "id": parent_id,
            "knowledge_base": kb,
            "title": entry.get("title", file_name),
            "description": entry.get("description", ""),
            "file_name": file_name,
            "file_type": ext.upper(),
            "file_size_bytes": Path(entry.get("file_path", "")).stat().st_size if Path(entry.get("file_path", "")).exists() else 0,
            "tags": entry.get("tags", []),
            "sections_count": len(section_ids),
            "status": overall_status,
            "source_type": entry.get("source_type", "upload"),
            "source_id": source_id,
            "source_name": _get_source_name(source_id),
            "drive_url": entry.get("drive_url", ""),
            "uploaded_by": entry.get("uploaded_by"),
            "uploaded_by_name": _get_user_name(entry.get("uploaded_by")),
            "created_at": entry.get("uploaded_at", ""),
            "download_url": f"http://localhost:8000/api/v1/files/{parent_id}/download",
        })

    # Sort by created_at desc
    docs.sort(key=lambda d: d.get("created_at", ""), reverse=True)
    return docs


from fastapi import UploadFile, File, Form
import tempfile


def preprocess_excel(file_content: bytes, file_name: str) -> str:
    """Convert Excel to markdown for better RAG indexing."""
    try:
        import openpyxl
        import io
        wb = openpyxl.load_workbook(io.BytesIO(file_content), data_only=True)
        result = [f"# {file_name}\n"]
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            result.append(f"## {sheet_name}\n")
            for row in ws.iter_rows(values_only=True):
                cells = []
                for c in row:
                    s = str(c).strip() if c is not None else ""
                    if s == "None":
                        s = ""
                    cells.append(s)
                non_empty = [c for c in cells if c]
                if non_empty:
                    result.append(" | ".join(non_empty))
            result.append("")
        return "\n".join(result)
    except Exception as e:
        print(f"Excel preprocess error: {e}")
        return None


def preprocess_docx(file_content: bytes, file_name: str) -> str:
    """Convert DOCX to markdown for better RAG indexing."""
    try:
        from docx import Document as DocxDoc
        import io
        doc = DocxDoc(io.BytesIO(file_content))
        parts = [f"# {file_name}\n"]
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            if "Heading" in style:
                level = 2
                try:
                    level = int(style[-1])
                except:
                    pass
                parts.append(f"{'#' * level} {text}")
            else:
                parts.append(text)
        for i, table in enumerate(doc.tables):
            parts.append(f"\n### Bảng {i+1}\n")
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip().replace("\n", " | ") for cell in row.cells]
                cleaned = []
                prev = None
                for c in cells:
                    if c != prev:
                        cleaned.append(c)
                    prev = c
                if row_idx == 0:
                    parts.append("| " + " | ".join(cleaned) + " |")
                    parts.append("| " + " | ".join(["---"] * len(cleaned)) + " |")
                else:
                    parts.append("| " + " | ".join(cleaned) + " |")
            parts.append("")
        return "\n".join(parts)
    except Exception as e:
        print(f"DOCX preprocess error: {e}")
        return None


def split_into_sections(markdown_text: str, doc_title: str, max_tokens: int = 1500) -> list[dict]:
    """Split markdown into semantic sections for better RAG chunking.

    Each section keeps its heading context and stays under max_tokens.
    Returns list of {name, text} dicts ready for Dify upload.
    """
    lines = markdown_text.split("\n")
    sections = []
    current_section = []
    current_heading = doc_title

    # Detect section boundaries: Roman numerals, ##, numbered headings
    import re
    section_pattern = re.compile(
        r'^(#{1,3}\s+|[IVX]+\.\s+|\d+\.\s+[A-ZĐ]|[A-ZĐ]{2,}.*:$)',
        re.MULTILINE
    )

    for line in lines:
        is_heading = bool(section_pattern.match(line.strip()))

        if is_heading and current_section:
            # Save current section
            text = "\n".join(current_section).strip()
            if text and len(text) > 50:  # Skip tiny sections
                sections.append({
                    "name": f"{doc_title} — {current_heading}",
                    "text": f"# {doc_title}\n## {current_heading}\n\n{text}",
                })
            current_section = []
            current_heading = line.strip().lstrip("#").strip()

        current_section.append(line)

    # Don't forget last section
    if current_section:
        text = "\n".join(current_section).strip()
        if text and len(text) > 50:
            sections.append({
                "name": f"{doc_title} — {current_heading}",
                "text": f"# {doc_title}\n## {current_heading}\n\n{text}",
            })

    # Merge small sections (< 200 chars) with next one
    merged = []
    buffer = None
    for sec in sections:
        if buffer:
            if len(buffer["text"]) < 200:
                buffer["text"] += "\n\n" + sec["text"]
                buffer["name"] = sec["name"]
            else:
                merged.append(buffer)
                buffer = sec
        else:
            buffer = sec
    if buffer:
        merged.append(buffer)

    return merged if merged else [{"name": doc_title, "text": markdown_text}]


# File storage for original files (for download)
UPLOAD_DIR = Path(__file__).parent / "data" / "uploads"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

# Parent document registry:
# parent_id → {file_name, title, file_path, knowledge_base, uploaded_at, section_doc_ids: [...], drive_url?, source_type}
FILE_REGISTRY: dict[str, dict] = {}

# Load existing files from disk
_registry_path = UPLOAD_DIR / "_registry.json"
if _registry_path.exists():
    try:
        FILE_REGISTRY = json.loads(_registry_path.read_text())
    except Exception:
        pass


def save_registry():
    _registry_path.write_text(json.dumps(FILE_REGISTRY, ensure_ascii=False, indent=2))


# ---- Backfill FILE_REGISTRY from drive_sync_state.json ----
# Files synced before the registration fix won't appear in the knowledge list.
# This reads the sync state and creates missing FILE_REGISTRY entries at startup.
_drive_state_path = Path(__file__).parent / "data" / "drive_sync_state.json"
if _drive_state_path.exists():
    try:
        _drive_state = json.loads(_drive_state_path.read_text())
        # Build reverse map: dataset_id → knowledge_base name
        _dataset_to_kb = {v: k for k, v in DIFY_DATASET_IDS.items() if v}
        _backfill_count = 0
        _existing_drive_ids = {
            e.get("drive_file_id") for e in FILE_REGISTRY.values() if e.get("drive_file_id")
        }

        for _state_key, _state_entry in _drive_state.items():
            _dataset_id = _state_entry.get("dataset_id", "")
            _kb_name = _dataset_to_kb.get(_dataset_id, "general")
            _files = _state_entry.get("files", {})

            for _fid, _finfo in _files.items():
                _doc_ids = _finfo.get("doc_ids", [])
                if not _doc_ids:
                    continue
                if _fid in _existing_drive_ids:
                    continue

                _fname = _finfo.get("name", _fid)
                _parent_id = f"drive-{_fid[:12]}"

                # Guess file type from name
                _ext = _fname.rsplit(".", 1)[-1].upper() if "." in _fname else "FILE"
                # Google native types don't have extensions — detect from name patterns
                _is_sheet = False
                if _ext == _fname.upper():
                    _ext = "GSHEET"  # extensionless files from Drive are mostly Sheets
                    _is_sheet = True

                # Build proper Google URL based on detected type
                if _is_sheet or _ext == "GSHEET":
                    _drive_url = f"https://docs.google.com/spreadsheets/d/{_fid}"
                elif _ext == "GDOC":
                    _drive_url = f"https://docs.google.com/document/d/{_fid}"
                elif _ext == "GSLIDES":
                    _drive_url = f"https://docs.google.com/presentation/d/{_fid}"
                else:
                    _drive_url = f"https://drive.google.com/file/d/{_fid}/view"

                FILE_REGISTRY[_parent_id] = {
                    "file_name": _fname,
                    "title": _fname,
                    "file_path": "",
                    "knowledge_base": _kb_name,
                    "uploaded_at": _finfo.get("synced_at", datetime.now().isoformat()),
                    "section_doc_ids": _doc_ids,
                    "source_type": "google_drive",
                    "drive_url": _drive_url,
                    "drive_file_id": _fid,
                    "file_type": _ext,
                    "file_size_bytes": 0,
                    "tags": [],
                }
                _existing_drive_ids.add(_fid)
                _backfill_count += 1

        if _backfill_count > 0:
            save_registry()
            print(f"[Backfill] Registered {_backfill_count} Drive files from sync state into FILE_REGISTRY")
        # Cleanup temp vars
        del _drive_state, _dataset_to_kb, _backfill_count, _existing_drive_ids
    except Exception as _e:
        print(f"[Backfill] Warning: could not backfill from drive_sync_state.json: {_e}")

# --- Fix drive_url for existing GSHEET entries (one-time migration) ---
_url_fixed = 0
for _pid, _entry in FILE_REGISTRY.items():
    _ft = _entry.get("file_type", "")
    _du = _entry.get("drive_url", "")
    _fid = _entry.get("drive_file_id", "")
    if not _fid or not _du:
        continue
    # Fix GSHEET: /file/d/ → /spreadsheets/d/
    if _ft == "GSHEET" and "/file/d/" in _du:
        _entry["drive_url"] = f"https://docs.google.com/spreadsheets/d/{_fid}"
        _url_fixed += 1
    elif _ft == "GDOC" and "/file/d/" in _du:
        _entry["drive_url"] = f"https://docs.google.com/document/d/{_fid}"
        _url_fixed += 1
    elif _ft == "GSLIDES" and "/file/d/" in _du:
        _entry["drive_url"] = f"https://docs.google.com/presentation/d/{_fid}"
        _url_fixed += 1
if _url_fixed > 0:
    save_registry()
    print(f"[Migration] Fixed drive_url for {_url_fixed} Google native files")


# ---- Drive Sources Registry ----
# Tracks imported Drive folders/links as "sources" for grouping and bulk delete.
DRIVE_SOURCES: dict[str, dict] = {}
_sources_path = UPLOAD_DIR / "_sources.json"
if _sources_path.exists():
    try:
        DRIVE_SOURCES = json.loads(_sources_path.read_text())
    except Exception:
        pass


def save_sources():
    _sources_path.write_text(json.dumps(DRIVE_SOURCES, ensure_ascii=False, indent=2))


# Helper: get user name from MOCK_USERS by user_id
def _get_user_name(user_id: str | None) -> str:
    if not user_id:
        return ""
    for u in MOCK_USERS:
        if u["id"] == user_id:
            return u["name"]
    return ""


# Helper: get source name from DRIVE_SOURCES
def _get_source_name(source_id: str | None) -> str:
    if not source_id:
        return ""
    src = DRIVE_SOURCES.get(source_id)
    return src.get("name", "") if src else ""


from fastapi.responses import FileResponse


@app.get("/api/v1/files/{file_id}/download")
async def download_file(file_id: str):
    """Download original file by Dify document ID or file registry ID."""
    # Try file registry first
    entry = FILE_REGISTRY.get(file_id)
    if entry:
        fpath = Path(entry["file_path"])
        if fpath.exists():
            return FileResponse(
                str(fpath),
                filename=entry["file_name"],
                media_type="application/octet-stream",
            )

    raise HTTPException(404, "File không tồn tại")


@app.get("/api/v1/files")
async def list_files():
    """List all downloadable files."""
    files = []
    for doc_id, entry in FILE_REGISTRY.items():
        files.append({
            "id": doc_id,
            "file_name": entry["file_name"],
            "knowledge_base": entry.get("knowledge_base", ""),
            "uploaded_at": entry.get("uploaded_at", ""),
            "download_url": f"/api/v1/files/{doc_id}/download",
        })
    return files


@app.post("/api/v1/admin/knowledge/upload")
async def upload_knowledge(
    request: Request,
    file: UploadFile = File(...),
    knowledge_base: str = Form("general"),
    title: str = Form(""),
    description: str = Form(""),
    tags: str = Form(""),
    auto_chunk: bool = Form(True),
    chunk_size: int = Form(800),
    chunk_overlap: int = Form(100),
):
    """Upload document to Dify with auto pre-processing for Excel/DOCX."""
    import traceback as tb
    user = get_current_user(request)
    uploaded_by = user.get("sub", "")
    print(f"[Upload] Received: kb={knowledge_base}, title={title}, file={file.filename}, size={file.size}")
    ds_id = DIFY_DATASET_IDS.get(knowledge_base, "")
    if not ds_id or not DIFY_DATASET_API_KEY:
        raise HTTPException(400, f"Dataset ID not configured for '{knowledge_base}'")

    # File size check — max 50MB
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    file_content = await file.read()
    file_size = len(file_content)
    if file_size > MAX_FILE_SIZE:
        raise HTTPException(413, f"File quá lớn ({file_size // (1024*1024)}MB). Giới hạn tối đa 50MB.")
    print(f"[Upload] File size: {file_size / (1024*1024):.1f}MB")
    file_name = file.filename or title or "document"

    # Save original file for download
    stored_name = f"{uuid.uuid4()}_{file_name}"
    stored_path = UPLOAD_DIR / stored_name
    stored_path.write_bytes(file_content)
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""

    # Auto pre-process for better RAG quality
    # Offload CPU-heavy parsing to thread to keep event loop responsive
    preprocessed_text = None
    if ext in ("xlsx", "xls"):
        preprocessed_text = await asyncio.to_thread(preprocess_excel, file_content, file_name)
    elif ext in ("docx", "doc"):
        preprocessed_text = await asyncio.to_thread(preprocess_docx, file_content, file_name)
    elif ext == "pdf":
        def _process_pdf():
            import pdfplumber, tempfile as tf
            with tf.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                tmp.write(file_content)
                tmp_path = tmp.name
            parts = [f"# {title or file_name}\n"]
            try:
                with pdfplumber.open(tmp_path) as pdf:
                    total_pages = len(pdf.pages)
                    print(f"[Upload] PDF has {total_pages} pages, size={file_size/(1024*1024):.1f}MB")
                    for i, page in enumerate(pdf.pages):
                        try:
                            text = page.extract_text()
                            if text and len(text.strip()) > 20:
                                parts.append(f"## Trang {i+1}\n{text}")
                        except Exception as pe:
                            print(f"[Upload] PDF page {i+1} error: {pe}")
                            continue
            finally:
                import os; os.unlink(tmp_path)
            if len(parts) > 1:
                print(f"[Upload] PDF extracted {len(parts)-1} pages text")
                return "\n\n".join(parts)
            print(f"[Upload] PDF is image-based, uploading raw to Dify")
            return None
        try:
            preprocessed_text = await asyncio.to_thread(_process_pdf)
        except Exception as e:
            print(f"[Upload] PDF preprocess failed: {e}, uploading raw")
    elif ext in ("txt", "md", "csv"):
        try:
            raw_text = file_content.decode("utf-8", errors="replace")
            preprocessed_text = f"# {title or file_name}\n\n{raw_text}"
            print(f"[Upload] Text file: {len(raw_text)} chars")
        except Exception as e:
            print(f"[Upload] Text decode error: {e}")

    # For files >15MB, MUST use create_by_text (Dify has 15MB file limit)
    DIFY_FILE_LIMIT = 15 * 1024 * 1024
    if file_size > DIFY_FILE_LIMIT and not preprocessed_text:
        raise HTTPException(
            413,
            f"File {file_size/(1024*1024):.0f}MB vượt giới hạn Dify (15MB) và không thể trích xuất nội dung text. "
            f"Vui lòng chia nhỏ file hoặc dùng định dạng khác (DOCX, XLSX, TXT)."
        )

    chunk_config = {
        "indexing_technique": "high_quality",
        "process_rule": {"mode": "automatic"},
    }

    section_doc_ids = []
    result = None

    try:
        async with httpx.AsyncClient(timeout=300.0) as client:
            if preprocessed_text:
                sections = split_into_sections(preprocessed_text, title or file_name)
                print(f"[Upload] Split '{title}' into {len(sections)} sections")

                for i, section in enumerate(sections):
                    resp = await client.post(
                        f"{DIFY_BASE_URL}/datasets/{ds_id}/document/create_by_text",
                        headers={
                            "Authorization": f"Bearer {DIFY_DATASET_API_KEY}",
                            "Content-Type": "application/json",
                        },
                        json={
                            "name": section["name"],
                            "text": section["text"],
                            "indexing_technique": "high_quality",
                            "process_rule": {"mode": "automatic"},
                        },
                    )
                    if resp.status_code == 200:
                        sec_id = resp.json().get("document", {}).get("id")
                        if sec_id:
                            section_doc_ids.append(sec_id)
                        if result is None:
                            result = resp.json()
                    else:
                        print(f"[Upload] Section {i+1} error: {resp.status_code}")
            else:
                resp = await client.post(
                    f"{DIFY_BASE_URL}/datasets/{ds_id}/document/create_by_file",
                    headers={"Authorization": f"Bearer {DIFY_DATASET_API_KEY}"},
                    files={"file": (file_name, file_content, file.content_type or "application/octet-stream")},
                    data={"data": json.dumps(chunk_config)},
                )
                if resp.status_code == 200:
                    result = resp.json()
                    sec_id = result.get("document", {}).get("id")
                    if sec_id:
                        section_doc_ids.append(sec_id)

            if not result:
                raise HTTPException(500, "Dify upload failed")
    except HTTPException:
        raise
    except Exception as e:
        print(f"[Upload] ERROR: {e}")
        tb.print_exc()
        raise HTTPException(500, f"Upload error: {e}")

    # Parent document ID = first section ID or UUID
    parent_id = section_doc_ids[0] if section_doc_ids else str(uuid.uuid4())

    parsed_tags = json.loads(tags) if tags else []
    FILE_REGISTRY[parent_id] = {
        "file_name": file_name,
        "title": title or file_name,
        "description": description,
        "file_path": str(stored_path),
        "knowledge_base": knowledge_base,
        "source_type": "upload",
        "uploaded_at": datetime.utcnow().isoformat(),
        "section_doc_ids": section_doc_ids,
        "sections_count": len(section_doc_ids),
        "tags": parsed_tags,
        "uploaded_by": uploaded_by,
        "source_id": None,
    }
    save_registry()

    return {
        "status": "success",
        "document_id": parent_id,
        "sections_count": len(section_doc_ids),
        "processing_time_ms": 0,
        "preprocessed": preprocessed_text is not None,
        "download_url": f"http://localhost:8000/api/v1/files/{parent_id}/download",
    }


@app.delete("/api/v1/admin/knowledge/{doc_id}")
async def delete_knowledge(doc_id: str, knowledge_base: Optional[str] = None):
    """Delete parent document and ALL its sections from Dify (cascade delete)."""
    if not DIFY_DATASET_API_KEY:
        raise HTTPException(500, "DIFY_DATASET_API_KEY not configured")

    # Find parent in registry
    parent = FILE_REGISTRY.get(doc_id)
    section_ids_to_delete = []

    if parent:
        section_ids_to_delete = parent.get("section_doc_ids", [doc_id])
        kb = parent.get("knowledge_base", knowledge_base or "")
    else:
        section_ids_to_delete = [doc_id]
        kb = knowledge_base

    # Build datasets to try
    datasets_to_try = {}
    if kb and kb in DIFY_DATASET_IDS:
        datasets_to_try[kb] = DIFY_DATASET_IDS[kb]
    else:
        datasets_to_try = {k: v for k, v in DIFY_DATASET_IDS.items() if v}

    deleted_count = 0
    async with httpx.AsyncClient(timeout=30.0) as client:
        for sec_id in section_ids_to_delete:
            for kb_key, ds_id in datasets_to_try.items():
                try:
                    resp = await client.delete(
                        f"{DIFY_BASE_URL}/datasets/{ds_id}/documents/{sec_id}",
                        headers={"Authorization": f"Bearer {DIFY_DATASET_API_KEY}"},
                    )
                    if resp.status_code in (200, 204):
                        deleted_count += 1
                        break
                except Exception as e:
                    continue

    # Remove from registry
    if doc_id in FILE_REGISTRY:
        del FILE_REGISTRY[doc_id]
        save_registry()

    if deleted_count > 0 or parent:
        return {
            "status": "ok",
            "message": f"Đã xoá tài liệu và {deleted_count} sections",
            "deleted_sections": deleted_count,
        }

    raise HTTPException(404, "Không tìm thấy tài liệu trong hệ thống")


# --- Edit Knowledge Document ---

class EditDocumentRequest(BaseModel):
    title: Optional[str] = None
    description: Optional[str] = None
    tags: Optional[List[str]] = None


@app.put("/api/v1/admin/knowledge/{doc_id}")
async def edit_knowledge(doc_id: str, req: EditDocumentRequest):
    """Edit document title, description, and tags. If description changes, push summary to Dify."""
    if doc_id not in FILE_REGISTRY:
        raise HTTPException(404, "Document not found")
    entry = FILE_REGISTRY[doc_id]
    old_description = entry.get("description", "")

    if req.title is not None:
        entry["title"] = req.title
    if req.description is not None:
        entry["description"] = req.description
    if req.tags is not None:
        entry["tags"] = req.tags
    save_registry()

    # If description changed and non-empty, push summary as a Dify metadata section
    summary_section_id = None
    new_desc = req.description or ""
    if new_desc and new_desc != old_description and DIFY_DATASET_API_KEY:
        kb = entry.get("knowledge_base", "general")
        ds_id = DIFY_DATASET_IDS.get(kb)
        if ds_id:
            doc_title = entry.get("title", entry.get("file_name", doc_id))
            summary_text = f"# Tổng quan: {doc_title}\n\n{new_desc}\n\nFile: {entry.get('file_name', '')}\nLoại: {entry.get('file_type', '')}\nSản phẩm: {', '.join(entry.get('tags', []))}"

            # Delete old summary section if exists
            old_summary_id = entry.get("summary_section_id")
            if old_summary_id:
                try:
                    async with httpx.AsyncClient(timeout=30.0) as client:
                        await client.delete(
                            f"{DIFY_BASE_URL}/datasets/{ds_id}/documents/{old_summary_id}",
                            headers={"Authorization": f"Bearer {DIFY_DATASET_API_KEY}"},
                        )
                except Exception:
                    pass

            # Create new summary section in Dify
            try:
                async with httpx.AsyncClient(timeout=30.0) as client:
                    resp = await client.post(
                        f"{DIFY_BASE_URL}/datasets/{ds_id}/document/create_by_text",
                        headers={
                            "Authorization": f"Bearer {DIFY_DATASET_API_KEY}",
                            "Content-Type": "application/json",
                        },
                        json={
                            "name": f"[Summary] {doc_title}",
                            "text": summary_text,
                            "indexing_technique": "high_quality",
                            "process_rule": {"mode": "automatic"},
                        },
                    )
                    if resp.status_code == 200:
                        did = resp.json().get("document", {}).get("id")
                        if did:
                            summary_section_id = did
                            entry["summary_section_id"] = did
                            # Also add to section_doc_ids for tracking
                            if did not in entry.get("section_doc_ids", []):
                                entry.setdefault("section_doc_ids", []).append(did)
                            save_registry()
                            print(f"[Edit] Pushed summary to Dify: {did}")
            except Exception as e:
                print(f"[Edit] Failed to push summary to Dify: {e}")

    return {"status": "ok", "id": doc_id, "summary_section_id": summary_section_id}


@app.post("/api/v1/admin/knowledge/{doc_id}/auto-summary")
async def auto_summary(doc_id: str):
    """Read document content and generate an AI summary using Claude."""
    if doc_id not in FILE_REGISTRY:
        raise HTTPException(404, "Document not found")

    entry = FILE_REGISTRY[doc_id]
    file_name = entry.get("file_name", "")
    file_path = entry.get("file_path", "")
    source_type = entry.get("source_type", "")
    file_type = entry.get("file_type", "")
    drive_file_id = entry.get("drive_file_id", "")

    content_text = ""

    # Strategy 1: Read from local file if available
    if file_path and Path(file_path).exists():
        try:
            file_content = Path(file_path).read_bytes()
            ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""
            if ext in ("xlsx", "xls"):
                content_text = await asyncio.to_thread(preprocess_excel, file_content, file_name) or ""
            elif ext in ("docx", "doc"):
                content_text = await asyncio.to_thread(preprocess_docx, file_content, file_name) or ""
            elif ext == "pdf":
                def _read_pdf():
                    import pdfplumber, tempfile as tf
                    with tf.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                        tmp.write(file_content)
                        tmp_path = tmp.name
                    parts = []
                    try:
                        with pdfplumber.open(tmp_path) as pdf:
                            for page in pdf.pages:
                                text = page.extract_text()
                                if text:
                                    parts.append(text)
                    finally:
                        import os; os.unlink(tmp_path)
                    return "\n\n".join(parts)
                content_text = await asyncio.to_thread(_read_pdf)
            elif ext in ("txt", "md", "csv"):
                content_text = file_content.decode("utf-8", errors="replace")
        except Exception as e:
            print(f"[AutoSummary] Local file read error: {e}")

    # Strategy 2: Fetch from Google Drive if no local content
    if not content_text and drive_file_id and source_type in ("google_drive", "google_sheet", "google_doc"):
        try:
            import sys
            sys.path.insert(0, str(Path(__file__).parent / "api"))
            creds_path = str(Path(__file__).parent / "config" / "google-credentials.json")

            if file_type == "GSHEET" or source_type == "google_sheet":
                from services.google_sheets_sync import GoogleSheetsSync
                syncer = GoogleSheetsSync(
                    credentials_path=creds_path,
                    dify_base_url=DIFY_BASE_URL,
                    dify_dataset_api_key=DIFY_DATASET_API_KEY,
                )
                content_text = await asyncio.to_thread(syncer.sheet_to_markdown, drive_file_id)
            else:
                from services.google_drive_sync import GoogleDriveSync
                syncer = GoogleDriveSync(
                    credentials_path=creds_path,
                    dify_base_url=DIFY_BASE_URL,
                    dify_dataset_api_key=DIFY_DATASET_API_KEY,
                )
                mime_map = {
                    "GDOC": "application/vnd.google-apps.document",
                    "GSLIDES": "application/vnd.google-apps.presentation",
                    "PDF": "application/pdf",
                    "DOCX": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                }
                mime = mime_map.get(file_type, "application/vnd.google-apps.document")
                ext, raw = await asyncio.to_thread(syncer._download_file, drive_file_id, mime)
                content_text = await asyncio.to_thread(syncer._preprocess_file, ext, raw, file_name) or ""
        except Exception as e:
            print(f"[AutoSummary] Google Drive fetch error: {e}")

    if not content_text:
        raise HTTPException(400, "Không thể đọc nội dung tài liệu. File có thể không tồn tại hoặc không hỗ trợ.")

    # Truncate to avoid token limits (keep first ~8000 chars)
    truncated = content_text[:8000]

    # Call Claude to generate summary
    if not ANTHROPIC_API_KEY:
        raise HTTPException(500, "ANTHROPIC_API_KEY chưa được cấu hình")

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            resp = await client.post(
                "https://api.anthropic.com/v1/messages",
                headers={
                    "x-api-key": ANTHROPIC_API_KEY,
                    "anthropic-version": "2023-06-01",
                    "content-type": "application/json",
                },
                json={
                    "model": "claude-sonnet-4-20250514",
                    "max_tokens": 500,
                    "system": (
                        "Bạn là trợ lý AI tóm tắt nội dung tài liệu nội bộ công ty. "
                        "Hãy viết MÔ TẢ NGẮN GỌN (3-5 câu) bằng tiếng Việt về nội dung chính của tài liệu. "
                        "Tập trung vào: tài liệu này nói về gì, dùng cho mục đích gì, thông tin chính là gì. "
                        "KHÔNG dùng bullet points. Viết thành đoạn văn liền mạch. "
                        "KHÔNG lặp lại tên file. Chỉ trả về nội dung mô tả, không thêm gì khác."
                    ),
                    "messages": [{"role": "user", "content": f"Tên file: {file_name}\n\nNội dung:\n{truncated}"}],
                },
            )
            if resp.status_code == 200:
                data = resp.json()
                summary = data["content"][0]["text"].strip()
                print(f"[AutoSummary] Generated for '{file_name}': {summary[:100]}...")
                return {"status": "ok", "summary": summary, "content_length": len(content_text)}
            else:
                print(f"[AutoSummary] Claude API error {resp.status_code}: {resp.text[:200]}")
                raise HTTPException(500, f"Lỗi Claude API: {resp.status_code}")
    except HTTPException:
        raise
    except Exception as e:
        print(f"[AutoSummary] Error: {e}")
        raise HTTPException(500, f"Lỗi tạo tóm tắt: {str(e)}")


# --- Knowledge Sources CRUD ---

@app.get("/api/v1/admin/knowledge/sources")
async def list_sources():
    """List all imported Drive sources (folders, sheets, docs)."""
    sources = []
    for sid, src in DRIVE_SOURCES.items():
        # Count documents that still exist in FILE_REGISTRY
        valid_doc_ids = [d for d in src.get("document_ids", []) if d in FILE_REGISTRY]
        sources.append({
            **src,
            "document_count": len(valid_doc_ids),
            "uploaded_by_name": _get_user_name(src.get("uploaded_by")),
        })
    sources.sort(key=lambda s: s.get("created_at", ""), reverse=True)
    return sources


@app.delete("/api/v1/admin/knowledge/sources/{source_id}")
async def delete_source(source_id: str):
    """Delete a source and ALL its documents from the system (NOT from Google Drive)."""
    if source_id not in DRIVE_SOURCES:
        raise HTTPException(404, "Source not found")
    source = DRIVE_SOURCES[source_id]
    doc_ids = list(source.get("document_ids", []))

    # Cascade delete all documents from this source
    deleted_count = 0
    for doc_id in doc_ids:
        if doc_id in FILE_REGISTRY:
            try:
                await delete_knowledge(doc_id)
                deleted_count += 1
            except Exception:
                # If Dify delete fails, still remove from registry
                if doc_id in FILE_REGISTRY:
                    del FILE_REGISTRY[doc_id]
                    deleted_count += 1

    save_registry()
    del DRIVE_SOURCES[source_id]
    save_sources()
    return {"status": "ok", "deleted_documents": deleted_count, "message": f"Đã xóa nguồn và {deleted_count} tài liệu khỏi hệ thống"}


# --- Re-index by URL ---

class ReindexByUrlRequest(BaseModel):
    url: str

@app.post("/api/v1/admin/knowledge/search-by-url")
async def search_knowledge_by_url(req: ReindexByUrlRequest):
    """Search FILE_REGISTRY for documents matching a Google Drive/Sheet/Doc URL."""
    url = req.url.strip()

    # Extract file ID from various Google URL formats
    file_id = None
    id_match = re.search(r"/d/([a-zA-Z0-9_-]+)", url)
    if id_match:
        file_id = id_match.group(1)
    else:
        # Try /file/d/ format
        id_match = re.search(r"/file/d/([a-zA-Z0-9_-]+)", url)
        if id_match:
            file_id = id_match.group(1)

    if not file_id:
        raise HTTPException(400, "Không thể trích xuất ID từ URL. Vui lòng nhập đúng link Google Sheet/Doc/Drive.")

    # Search FILE_REGISTRY by drive_file_id or drive_url containing the file_id
    matches = []
    for parent_id, entry in FILE_REGISTRY.items():
        entry_fid = entry.get("drive_file_id", "")
        entry_url = entry.get("drive_url", "")
        if entry_fid == file_id or file_id in entry_url:
            matches.append({
                "id": parent_id,
                "title": entry.get("title", ""),
                "file_name": entry.get("file_name", ""),
                "file_type": entry.get("file_type", ""),
                "knowledge_base": entry.get("knowledge_base", ""),
                "drive_url": entry_url,
                "drive_file_id": entry_fid,
                "source_type": entry.get("source_type", ""),
                "tags": entry.get("tags", []),
                "sections_count": len(entry.get("section_doc_ids", [])),
                "created_at": entry.get("uploaded_at", ""),
            })

    return {"file_id": file_id, "matches": matches}


@app.post("/api/v1/admin/knowledge/reindex-by-url")
async def reindex_by_url(req: ReindexByUrlRequest, request: Request):
    """Delete existing index for a Google Sheet/Doc URL, then re-import it."""
    url = req.url.strip()

    # Extract file ID
    file_id = None
    id_match = re.search(r"/d/([a-zA-Z0-9_-]+)", url)
    if id_match:
        file_id = id_match.group(1)
    if not file_id:
        raise HTTPException(400, "Không thể trích xuất ID từ URL.")

    # Find matching document in FILE_REGISTRY
    found_parent_id = None
    found_entry = None
    for parent_id, entry in FILE_REGISTRY.items():
        entry_fid = entry.get("drive_file_id", "")
        entry_url = entry.get("drive_url", "")
        if entry_fid == file_id or file_id in entry_url:
            found_parent_id = parent_id
            found_entry = dict(entry)
            break

    if not found_parent_id or not found_entry:
        raise HTTPException(404, "Không tìm thấy tài liệu với URL này trong hệ thống.")

    # Save info before deleting
    knowledge_base = found_entry.get("knowledge_base", "general")
    tags = found_entry.get("tags", [])
    title = found_entry.get("title", "")
    source_id = found_entry.get("source_id", "")
    source_type = found_entry.get("source_type", "")
    drive_url = found_entry.get("drive_url", url)

    # Step 1: Delete existing index (sections from Dify + FILE_REGISTRY)
    try:
        await delete_knowledge(found_parent_id)
    except HTTPException:
        # Even if Dify delete fails, continue with re-import
        if found_parent_id in FILE_REGISTRY:
            del FILE_REGISTRY[found_parent_id]
            save_registry()

    # Also remove from DRIVE_SOURCES document_ids if applicable
    if source_id and source_id in DRIVE_SOURCES:
        doc_ids = DRIVE_SOURCES[source_id].get("document_ids", [])
        if found_parent_id in doc_ids:
            doc_ids.remove(found_parent_id)
            save_sources()

    # Step 2: Re-import
    dataset_id = DIFY_DATASET_IDS.get(knowledge_base)
    if not dataset_id:
        raise HTTPException(400, f"Không tìm thấy dataset cho knowledge_base '{knowledge_base}'")

    if not DIFY_DATASET_API_KEY:
        raise HTTPException(500, "DIFY_DATASET_API_KEY not configured")

    # Get current user
    user = get_current_user(request)
    uploaded_by = user.get("sub", "")

    is_sheet = "/spreadsheets/d/" in drive_url or source_type in ("google_sheet",)
    is_doc = "/document/d/" in drive_url or source_type in ("google_doc",)

    # Default to sheet if we can't determine type from URL
    if not is_sheet and not is_doc:
        file_type = found_entry.get("file_type", "")
        if file_type == "GSHEET":
            is_sheet = True
        elif file_type == "GDOC":
            is_doc = True
        else:
            # For drive files, try re-import as sheet since that's the common case
            is_sheet = True

    task_id = str(uuid.uuid4())[:8]
    link_type = "sheet" if is_sheet else "doc"
    IMPORT_TASKS[task_id] = {
        "task_id": task_id,
        "type": f"reindex-{link_type}",
        "status": "importing",
        "url": drive_url,
        "title": title,
        "knowledge_base": knowledge_base,
        "source_id": source_id,
        "sections_count": 0,
        "started_at": datetime.now().isoformat(),
        "completed_at": None,
    }

    asyncio.create_task(_run_link_import(
        task_id, file_id, dataset_id, knowledge_base,
        drive_url, title, is_sheet, source_id, tags, uploaded_by,
    ))

    return {
        "task_id": task_id,
        "status": "importing",
        "deleted_doc_id": found_parent_id,
        "message": f"Đã xóa index cũ và đang re-index '{title}'...",
    }


# --- Logs ---

@app.get("/api/v1/admin/logs")
async def list_logs(page: int = 1, limit: int = 20, channel: Optional[str] = None, department: Optional[str] = None):
    filtered = MOCK_LOGS
    if channel:
        filtered = [l for l in filtered if l["channel"] == channel]
    if department:
        filtered = [l for l in filtered if l["department_routed"] == department]

    total = len(filtered)
    pages = (total + limit - 1) // limit if total > 0 else 1
    start = (page - 1) * limit
    items = filtered[start:start + limit]

    return {"items": items, "total": total, "page": page, "limit": limit, "pages": pages}


# --- AI-Powered Query Rewriting ---

import re
import unicodedata

ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")

QUERY_REWRITE_PROMPT = """Bạn là module rewrite query tiếng Việt cho hệ thống RAG nội bộ của ViHAT Group (công ty CPaaS/Contact Center).

Nhiệm vụ: Nhận câu hỏi gốc từ user (có thể ngắn, viết tắt, mơ hồ) → rewrite thành 2-3 câu hỏi rõ ràng, cụ thể, tối ưu cho vector search.

Quy tắc:
1. Giữ nguyên ý nghĩa gốc
2. Mở rộng viết tắt: "omi" → "OmiCall", "zns" → "Zalo ZNS", "hđ" → "hợp đồng"
3. Nếu hỏi về 1 chủ đề → liệt kê các khía cạnh liên quan. VD: "ốm đau" → hỏi cả chế độ nghỉ bệnh + thăm hỏi nhập viện + mức hỗ trợ
4. Thêm context "tại ViHAT Group" nếu câu hỏi về nội bộ
5. Output CHỈ là các câu hỏi đã rewrite, cách nhau bằng dấu xuống dòng, KHÔNG giải thích

Ví dụ:
Input: "ốm đau thì sao"
Output: Chính sách ốm đau tại ViHAT Group là gì? Chế độ nghỉ bệnh BHXH và thăm hỏi ốm đau tai nạn nhập viện? Mức hỗ trợ thăm hỏi theo từng cấp bậc nhân viên?

Input: "giá omi"
Output: Bảng giá dịch vụ OmiCall tổng đài ảo? Giá các gói PBX 1, PBX 2, Call Center, OMNI Channel? Đơn giá theo user và thời hạn hợp đồng?

Input: "nghỉ phép mấy ngày"
Output: Chính sách nghỉ phép năm tại ViHAT Group được bao nhiêu ngày? Quy trình xin nghỉ phép dài ngày và ngắn ngày? Thưởng phép thâm niên từ năm thứ mấy?"""


async def ai_rewrite_query(query: str) -> str:
    """Use Claude to rewrite query for better RAG retrieval."""
    if not ANTHROPIC_API_KEY:
        print("[QueryRewrite] No ANTHROPIC_API_KEY, skipping AI rewrite")
        return query

    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            resp = await client.post(
                "https://api.anthropic.com/v1/messages",
                headers={
                    "x-api-key": ANTHROPIC_API_KEY,
                    "anthropic-version": "2023-06-01",
                    "content-type": "application/json",
                },
                json={
                    "model": "claude-sonnet-4-20250514",
                    "max_tokens": 200,
                    "system": QUERY_REWRITE_PROMPT,
                    "messages": [{"role": "user", "content": query}],
                },
            )
            if resp.status_code == 200:
                data = resp.json()
                rewritten = data["content"][0]["text"].strip()
                print(f"[QueryRewrite] '{query}' => '{rewritten[:120]}...'")
                return rewritten
            else:
                print(f"[QueryRewrite] API error {resp.status_code}: {resp.text[:100]}")
                return query
    except Exception as e:
        print(f"[QueryRewrite] Error: {e}")
        return query


def normalize_query(query: str) -> str:
    """Basic Vietnamese normalization."""
    query = unicodedata.normalize("NFC", query)
    # Simple product alias expansion (fast, no API call needed)
    aliases = {
        "omi": "OmiCall", "esms": "eSMS", "zns": "Zalo ZNS",
        "zalo oa": "Zalo OA", "omiflow": "OmiFlow",
    }
    lower = query.lower()
    for alias, full in aliases.items():
        if alias in lower:
            query = re.sub(re.escape(alias), full, query, flags=re.IGNORECASE)
    return query


# --- Query (Chat) ---

class QueryRequest(BaseModel):
    user_id: str
    query: str
    department: Optional[str] = None
    conversation_id: Optional[str] = None


FILE_REQUEST_PATTERNS = [
    # "gửi" variants
    "gửi file", "gửi cho tôi", "gửi cho anh", "gửi cho em",
    "gửi anh", "gửi tôi", "gửi em",
    "gửi tài liệu", "gửi bảng giá",
    # "cho" variants
    "cho tôi file", "cho anh file", "cho em file",
    "cho tôi tài liệu", "cho anh tài liệu",
    "cho bảng giá", "cho anh bảng giá", "cho tôi bảng giá",
    "cho anh link", "cho tôi link", "cho em link",
    "cho link", "cho đường link", "cho url",
    # "tải/lấy/xem/download"
    "tải file", "lấy file", "xem file", "download file",
    "tải tài liệu", "lấy tài liệu", "tải bảng giá", "lấy bảng giá",
    # "link" variants
    "link file", "link tài liệu", "link bảng giá", "link báo giá",
    "đường link", "share link", "chia sẻ link",
    # file-centric
    "file gốc", "file bảng giá",
]

# These words confirm it's a file request (not just "gửi anh lời chào")
FILE_CONFIRM_WORDS = {"file", "tài liệu", "bảng giá", "báo giá", "document", "tải", "download", "link", "đường link", "url"}


def detect_file_request(query: str) -> Optional[str]:
    """Detect if user is asking for a file download."""
    query_lower = query.lower()
    for kw in FILE_REQUEST_PATTERNS:
        if kw in query_lower:
            # Must also contain a file-related confirm word
            if any(w in query_lower for w in FILE_CONFIRM_WORDS):
                return query_lower
    return None


# Product aliases — user mentions a specific product/service
PRODUCT_ALIASES = {
    "omicall": ["omi", "omicall", "tổng đài", "contact center", "dịch vụ omi"],
    "esms": ["esms", "sms", "brandname", "tin nhắn"],
    "zns": ["zns", "zalo zns", "zalo notification"],
    "becare": ["becare", "be care"],
    "omiflow": ["omiflow", "omi flow", "automation"],
    "sdk": ["sdk", "omi sdk"],
    "zcc": ["zcc", "zalo cloud connect"],
    "voice_brandname": ["voice brandname", "voice bn"],
}
# Topic aliases — what kind of document
TOPIC_ALIASES = {
    "bảng giá": ["bảng giá", "giá", "pricing", "báo giá", "price"],
    "phúc lợi": ["phúc lợi", "phuc loi", "welfare", "chế độ"],
    "nội quy": ["nội quy", "noi quy", "quy định", "lao động"],
    "hướng dẫn": ["hướng dẫn", "tài liệu hướng dẫn", "user guide", "manual"],
    "proposal": ["proposal", "đề xuất"],
    "profile": ["profile", "giới thiệu công ty"],
    "hợp đồng": ["hợp đồng", "hop dong", "contract"],
    "chính sách": ["chính sách", "chinh sach", "policy"],
}

STOPWORDS = {"cho", "anh", "tôi", "em", "gửi", "link", "đường", "file", "tài", "liệu",
             "mới", "nhất", "nhé", "ạ", "à", "nha", "xin", "hãy", "giúp", "tải", "lấy",
             "xem", "download", "share", "chia", "sẻ", "về", "của", "và", "các", "những",
             "đi", "dùm", "giùm", "với", "được", "không", "nào"}

# Combined for backward compat in detect logic
SEARCH_ALIASES = {**PRODUCT_ALIASES, **TOPIC_ALIASES}


def find_matching_files(search_term: str) -> list:
    """Find files matching search term with product/topic separation + relevance scoring.

    Logic:
    - If user mentions a specific product (zns, omicall...) → files MUST match that product
    - Topic (bảng giá, hướng dẫn...) further filters/boosts relevance
    - Files matching both product + topic score highest
    """
    results = []
    search_lower = search_term.lower()

    # Detect which product and topic groups the user mentioned
    matched_products = []
    for key, aliases in PRODUCT_ALIASES.items():
        if any(a in search_lower for a in aliases):
            matched_products.append(key)

    matched_topics = []
    for key, aliases in TOPIC_ALIASES.items():
        if any(a in search_lower for a in aliases):
            matched_topics.append(key)

    # Build expanded search terms for each group
    product_terms = set()
    for p in matched_products:
        product_terms.update(PRODUCT_ALIASES[p])
        product_terms.add(p)

    topic_terms = set()
    for t in matched_topics:
        topic_terms.update(TOPIC_ALIASES[t])
        topic_terms.add(t)

    # Extra content words not covered by aliases
    content_words = set()
    for word in search_lower.split():
        if len(word) > 1 and word not in STOPWORDS:
            content_words.add(word)

    print(f"[FileSearch] query='{search_term}', products={matched_products}, "
          f"topics={matched_topics}, extra_words={content_words - product_terms - topic_terms}")

    seen_files = set()
    for doc_id, entry in FILE_REGISTRY.items():
        fname = entry["file_name"].lower()
        title = entry.get("title", "").lower()
        kb = entry.get("knowledge_base", "").lower()
        tags = " ".join(entry.get("tags", [])).lower()
        searchable = f"{fname} {title} {kb} {tags}"

        if fname in seen_files:
            continue

        # Score: product match is most important
        product_score = 0
        for p in matched_products:
            if any(a in searchable for a in PRODUCT_ALIASES[p]):
                product_score += 20

        topic_score = 0
        for t in matched_topics:
            if any(a in searchable for a in TOPIC_ALIASES[t]):
                topic_score += 10

        extra_score = sum(1 for w in content_words if len(w) > 2 and w in searchable and w not in STOPWORDS)

        total_score = product_score + topic_score + extra_score

        # FILTER: if user asked for a specific product, file MUST match that product
        if matched_products and product_score == 0:
            continue

        # FILTER: must match at least something meaningful
        if total_score == 0:
            continue

        # If only topic matched (no product specified), require topic match
        if not matched_products and matched_topics and topic_score == 0:
            continue

        seen_files.add(fname)
        drive_url = entry.get("drive_url", "")
        download_url = f"{PUBLIC_API_URL}/api/v1/files/{doc_id}/download"
        results.append({
            "file_name": entry["file_name"],
            "knowledge_base": entry.get("knowledge_base", ""),
            "download_url": download_url,
            "drive_url": drive_url,
            "source_type": entry.get("source_type", "upload"),
            "_score": total_score,
        })

    # Sort by relevance, limit to 10
    results.sort(key=lambda x: x.get("_score", 0), reverse=True)
    for r in results:
        r.pop("_score", None)
    return results[:10]


@app.post("/api/v1/query")
async def query(req: QueryRequest):
    """Chat via Dify RAG Engine with AI-powered query rewriting + file download detection."""
    import time
    start = time.time()

    if not DIFY_API_KEY:
        raise HTTPException(500, "DIFY_API_KEY not configured")

    # Check if user is requesting a file download
    file_search = detect_file_request(req.query)
    if file_search:
        matching_files = find_matching_files(file_search)
        if matching_files:
            file_lines = []
            for f in matching_files:
                # Use drive_url for Google Drive files (shareable), download_url for uploaded files
                if f.get("drive_url"):
                    file_lines.append(f"- 📎 **{f['file_name']}** — [Xem trên Drive]({f['drive_url']})")
                else:
                    file_lines.append(f"- 📎 **{f['file_name']}** — [Tải xuống]({f['download_url']})")
            file_list = "\n".join(file_lines)
            elapsed = int((time.time() - start) * 1000)
            return {
                "status": "success",
                "answer": f"Đây là các tài liệu bạn yêu cầu:\n\n{file_list}\n\nBạn có thể click vào link để xem hoặc tải xuống.",
                "sources": [],
                "conversation_id": str(uuid.uuid4()),
                "tokens_used": {"prompt": 0, "completion": 0},
                "processing_time_ms": elapsed,
                "files": matching_files,
            }
        # No matching files found — fall through to Dify RAG for a knowledge-based answer
        print(f"[FileSearch] No matching files, falling through to Dify RAG")

    # Step 1: Basic normalization (fast)
    normalized = normalize_query(req.query)

    # Step 2: AI-powered rewrite (only for first message, not follow-ups)
    if not req.conversation_id:
        rewritten = await ai_rewrite_query(normalized)
    else:
        rewritten = normalized

    payload = {
        "inputs": {},
        "query": rewritten,
        "response_mode": "blocking",
        "user": req.user_id,
    }
    if req.conversation_id:
        payload["conversation_id"] = req.conversation_id

    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            resp = await client.post(
                f"{DIFY_BASE_URL}/chat-messages",
                headers={"Authorization": f"Bearer {DIFY_API_KEY}", "Content-Type": "application/json"},
                json=payload,
            )
            # If 404 (conversation expired), retry without conversation_id
            if resp.status_code == 404 and "conversation_id" in payload:
                print(f"[Dify] Conversation expired, starting new one")
                payload.pop("conversation_id")
                resp = await client.post(
                    f"{DIFY_BASE_URL}/chat-messages",
                    headers={"Authorization": f"Bearer {DIFY_API_KEY}", "Content-Type": "application/json"},
                    json=payload,
                )
            resp.raise_for_status()
            data = resp.json()
    except httpx.HTTPStatusError as e:
        raise HTTPException(500, f"Dify error: {e}")
    except Exception as e:
        raise HTTPException(500, f"Dify error: {e}")

    # Extract sources
    sources = []
    for r in data.get("metadata", {}).get("retriever_resources", []):
        sources.append({
            "document": r.get("document_name", ""),
            "chunk": r.get("content", "")[:200],
            "score": r.get("score", 0),
        })

    elapsed = int((time.time() - start) * 1000)

    answer = data.get("answer", "")

    # Post-processing: attach real document links based on RAG sources
    answer = _enrich_answer_with_links(answer, sources)

    return {
        "status": "success",
        "answer": answer,
        "sources": sources,
        "conversation_id": data.get("conversation_id", str(uuid.uuid4())),
        "tokens_used": {
            "prompt": data.get("metadata", {}).get("usage", {}).get("prompt_tokens", 0),
            "completion": data.get("metadata", {}).get("usage", {}).get("completion_tokens", 0),
        },
        "processing_time_ms": elapsed,
    }


def _enrich_answer_with_links(answer: str, sources: list) -> str:
    """Post-process Dify answer: remove hallucinated links & append real document links.

    Dify LLM sometimes creates fake links like 'Xem chi tiết tại: [link](...)'.
    We strip those and append real links from FILE_REGISTRY based on source docs.
    """
    import re

    # 1. Remove ALL markdown links from Dify answer (Dify LLM hallucinates wrong URLs)
    # We will append correct links at the end based on FILE_REGISTRY
    def _replace_link(m):
        text = m.group(1)
        # Keep only display text, strip the URL
        return text

    answer = re.sub(r'\[([^\]]+)\]\(([^)]+)\)', _replace_link, answer)
    # Also strip bare URLs that Dify might include
    answer = re.sub(r'https?://docs\.google\.com/\S+', '', answer)
    answer = re.sub(r'https?://drive\.google\.com/\S+', '', answer)

    # 2. Find matching real documents from registry based on sources
    if not sources:
        return answer

    seen_links = set()
    real_links = []
    for src in sources[:3]:
        doc_name = src.get("document", "")
        score = src.get("score", 0)
        if score < 0.55 or not doc_name:
            continue

        # Match source doc name to FILE_REGISTRY entries
        doc_lower = doc_name.lower()
        for pid, entry in FILE_REGISTRY.items():
            fname = entry.get("file_name", "").lower()
            title = entry.get("title", "").lower()
            # Match by partial name (Dify section names are like "Phúc lợi ViHAT — 2. Mừng Thâm niên")
            base_name = doc_lower.split("—")[0].strip() if "—" in doc_lower else doc_lower
            if base_name and (base_name in fname or base_name in title or fname in base_name or title in base_name):
                drive_url = entry.get("drive_url", "")
                if drive_url and drive_url not in seen_links:
                    seen_links.add(drive_url)
                    display_name = entry.get("title", entry.get("file_name", ""))
                    real_links.append(f"🔗 [{display_name}]({drive_url})")
                elif not drive_url:
                    dl_url = f"{PUBLIC_API_URL}/api/v1/files/{pid}/download"
                    if dl_url not in seen_links:
                        seen_links.add(dl_url)
                        display_name = entry.get("title", entry.get("file_name", ""))
                        real_links.append(f"📎 [{display_name}]({dl_url})")
                break

    # 3. Append real links if found
    if real_links:
        links_section = "\n\n📂 **Tài liệu tham khảo:**\n" + "\n".join(f"- {l}" for l in real_links)
        answer += links_section

    return answer


# --- Health ---

# === Google Sheets Sync endpoints ===

@app.post("/api/v1/sheets/sync")
async def sheets_sync(req: dict):
    """Sync a Google Sheet to Dify KB."""
    try:
        import sys
        sys.path.insert(0, str(Path(__file__).parent / "api"))
        from services.google_sheets_sync import GoogleSheetsSync

        syncer = GoogleSheetsSync(
            credentials_path=str(Path(__file__).parent / "config" / "google-credentials.json"),
            dify_base_url=DIFY_BASE_URL,
            dify_dataset_api_key=DIFY_DATASET_API_KEY,
        )
        result = await syncer.sync_sheet(
            spreadsheet_id=req.get("spreadsheet_id", ""),
            dataset_id=req.get("dataset_id", ""),
            title=req.get("title", ""),
            force=req.get("force", False),
        )
        return result
    except FileNotFoundError:
        raise HTTPException(400, "Google credentials not found. Place google-credentials.json in config/")
    except Exception as e:
        raise HTTPException(500, str(e))


@app.get("/api/v1/sheets/status")
async def sheets_status():
    """Get sync status of all configured sheets."""
    try:
        import sys
        sys.path.insert(0, str(Path(__file__).parent / "api"))
        from services.google_sheets_sync import GoogleSheetsSync

        syncer = GoogleSheetsSync(
            credentials_path=str(Path(__file__).parent / "config" / "google-credentials.json"),
            dify_base_url=DIFY_BASE_URL,
            dify_dataset_api_key=DIFY_DATASET_API_KEY,
        )
        return {"sheets": syncer.get_sync_status()}
    except Exception:
        return {"sheets": []}


# ---- Conversation Persistence ----
CONVERSATIONS_FILE = Path(__file__).parent / "data" / "conversations.json"


def _load_conversations() -> dict:
    if CONVERSATIONS_FILE.exists():
        try:
            return json.loads(CONVERSATIONS_FILE.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            return {}
    return {}


def _save_conversations(data: dict):
    CONVERSATIONS_FILE.parent.mkdir(parents=True, exist_ok=True)
    CONVERSATIONS_FILE.write_text(json.dumps(data, ensure_ascii=False, indent=2, default=str), encoding="utf-8")


class SaveMessageRequest(BaseModel):
    role: str
    content: str
    sources: Optional[list] = None


@app.get("/api/v1/chat/conversations")
async def list_conversations():
    """List all conversations, newest first."""
    convs = _load_conversations()
    result = []
    for conv in convs.values():
        result.append({
            "id": conv["id"],
            "title": conv.get("title", "Cuộc hội thoại mới"),
            "created_at": conv["created_at"],
            "message_count": len(conv.get("messages", [])),
        })
    result.sort(key=lambda c: c["created_at"], reverse=True)
    return result


@app.get("/api/v1/chat/conversations/{conv_id}/messages")
async def get_conversation_messages(conv_id: str):
    """Get all messages for a conversation."""
    convs = _load_conversations()
    if conv_id not in convs:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return convs[conv_id].get("messages", [])


@app.post("/api/v1/chat/conversations")
async def create_conversation():
    """Create a new empty conversation."""
    convs = _load_conversations()
    new_id = str(uuid.uuid4())
    conv = {
        "id": new_id,
        "title": "Cuộc hội thoại mới",
        "created_at": datetime.utcnow().isoformat(),
        "messages": [],
    }
    convs[new_id] = conv
    _save_conversations(convs)
    return {"id": new_id, "title": conv["title"], "created_at": conv["created_at"]}


@app.post("/api/v1/chat/conversations/{conv_id}/messages")
async def save_message(conv_id: str, msg: SaveMessageRequest):
    """Append a message to a conversation."""
    convs = _load_conversations()
    if conv_id not in convs:
        raise HTTPException(status_code=404, detail="Conversation not found")

    message = {
        "role": msg.role,
        "content": msg.content,
        "sources": msg.sources,
        "timestamp": datetime.utcnow().isoformat(),
    }
    convs[conv_id]["messages"].append(message)

    # Update title from first user message
    if msg.role == "user" and convs[conv_id]["title"] == "Cuộc hội thoại mới":
        convs[conv_id]["title"] = msg.content[:50] + ("..." if len(msg.content) > 50 else "")

    _save_conversations(convs)
    return message


@app.delete("/api/v1/chat/conversations/{conv_id}")
async def delete_conversation(conv_id: str):
    """Delete a conversation."""
    convs = _load_conversations()
    if conv_id not in convs:
        raise HTTPException(status_code=404, detail="Conversation not found")
    del convs[conv_id]
    _save_conversations(convs)
    return {"status": "deleted"}


# ---- Feedback ----

FEEDBACK_FILE = Path(__file__).parent / "data" / "feedback.json"

def _load_feedback() -> dict:
    if FEEDBACK_FILE.exists():
        return json.loads(FEEDBACK_FILE.read_text(encoding="utf-8"))
    return {}

def _save_feedback(data: dict):
    FEEDBACK_FILE.parent.mkdir(parents=True, exist_ok=True)
    FEEDBACK_FILE.write_text(json.dumps(data, ensure_ascii=False, default=str, indent=2), encoding="utf-8")

class FeedbackSubmit(BaseModel):
    query_text: str
    answer_text: str
    sources: list = []
    category: str  # "wrong_answer" | "no_answer" | "outdated"
    user_comment: str
    conversation_id: Optional[str] = None

class FeedbackStatusUpdate(BaseModel):
    status: str  # "reviewing" | "resolved"
    admin_note: str = ""

@app.post("/api/v1/feedback")
async def submit_feedback(body: FeedbackSubmit):
    """Submit feedback on a chat response."""
    feedback_id = str(uuid.uuid4())
    feedback = {
        "id": feedback_id,
        "query_text": body.query_text,
        "answer_text": body.answer_text,
        "sources": body.sources,
        "category": body.category,
        "user_comment": body.user_comment,
        "conversation_id": body.conversation_id,
        "status": "new",
        "admin_note": "",
        "created_at": datetime.utcnow().isoformat(),
    }
    data = _load_feedback()
    data[feedback_id] = feedback
    _save_feedback(data)
    return {"id": feedback_id, "status": "new"}

@app.get("/api/v1/admin/feedback")
async def list_feedback(status: Optional[str] = None, page: int = 1, limit: int = 20):
    """List all feedback (admin)."""
    data = _load_feedback()
    items = list(data.values())

    # Filter by status
    if status:
        items = [f for f in items if f["status"] == status]

    # Sort by created_at descending
    items.sort(key=lambda x: x.get("created_at", ""), reverse=True)

    total = len(items)
    pages = max(1, (total + limit - 1) // limit)
    start = (page - 1) * limit
    end = start + limit

    return {
        "items": items[start:end],
        "total": total,
        "page": page,
        "limit": limit,
        "pages": pages,
    }

@app.put("/api/v1/admin/feedback/{feedback_id}/status")
async def update_feedback_status(feedback_id: str, body: FeedbackStatusUpdate):
    """Update feedback status."""
    data = _load_feedback()
    if feedback_id not in data:
        raise HTTPException(status_code=404, detail="Feedback not found")

    data[feedback_id]["status"] = body.status
    data[feedback_id]["admin_note"] = body.admin_note
    _save_feedback(data)
    return data[feedback_id]


# --- Google Drive / Sheet Import Endpoints (Async Background Tasks) ---

import asyncio

# In-memory import task tracker
IMPORT_TASKS: dict = {}  # task_id -> {status, type, url, knowledge_base, ...}

class ImportDriveRequest(BaseModel):
    folder_url: str
    knowledge_base: str
    name: str = ""
    description: str = ""
    product_tags: List[str] = []
    note: str = ""

class ImportLinkRequest(BaseModel):
    url: str
    knowledge_base: str
    title: str = ""
    description: str = ""
    product_tags: List[str] = []
    note: str = ""


async def _run_drive_import(task_id: str, folder_id: str, dataset_id: str, knowledge_base: str, folder_url: str, source_id: str = "", product_tags: List[str] = [], uploaded_by: str = ""):
    """Background worker: import Drive folder."""
    task = IMPORT_TASKS[task_id]
    try:
        import sys
        sys.path.insert(0, str(Path(__file__).parent / "api"))
        from services.google_drive_sync import GoogleDriveSync

        creds_path = str(Path(__file__).parent / "config" / "google-credentials.json")
        if not Path(creds_path).exists():
            raise Exception("File google-credentials.json không tồn tại trong config/")

        syncer = GoogleDriveSync(
            credentials_path=creds_path,
            dify_base_url=DIFY_BASE_URL,
            dify_dataset_api_key=DIFY_DATASET_API_KEY,
        )

        # List files first — offload blocking Google API call to thread
        task["status_detail"] = "Đang quét thư mục..."
        try:
            files = await asyncio.to_thread(syncer.list_folder, folder_id)
        except Exception as list_err:
            err_msg = str(list_err)
            if "404" in err_msg or "not found" in err_msg.lower():
                raise Exception(f"Không tìm thấy thư mục. Hãy kiểm tra URL và đảm bảo đã chia sẻ thư mục với Service Account email.")
            elif "403" in err_msg or "permission" in err_msg.lower() or "forbidden" in err_msg.lower():
                raise Exception(f"Không có quyền truy cập thư mục. Hãy chia sẻ thư mục với Service Account email (Viewer).")
            raise Exception(f"Lỗi quét thư mục: {err_msg}")

        task["files_found"] = len(files)

        if len(files) == 0:
            task["status"] = "completed"
            task["synced"] = 0
            task["errors"] = 0
            task["details"] = []
            task["status_detail"] = "Thư mục trống hoặc không có file được hỗ trợ"
            task["completed_at"] = datetime.now().isoformat()
            return

        # Sync folder with progress updates
        task["status_detail"] = f"Đang sync {len(files)} file..."
        task["files_processed"] = 0
        result = await syncer.sync_folder(folder_id, dataset_id)

        # Register synced files in FILE_REGISTRY so they appear in the knowledge list
        sync_state = syncer._sync_state
        state_key = f"folder:{folder_id}:{dataset_id}"
        folder_state = sync_state.get(state_key, {}).get("files", {})
        registered_count = 0

        for file_info in files:
            fid = file_info["id"]
            fname = file_info["name"]
            fs = folder_state.get(fid, {})
            doc_ids = fs.get("doc_ids", [])

            # Skip if no doc_ids (not synced) or already in registry
            if not doc_ids:
                continue
            # Check if already registered (by drive file id)
            already = any(
                e.get("drive_file_id") == fid
                for e in FILE_REGISTRY.values()
            )
            if already:
                continue

            parent_id = f"drive-{fid[:12]}"
            mime = file_info.get("mimeType", "")
            ext_map = {
                "application/vnd.google-apps.spreadsheet": "GSHEET",
                "application/vnd.google-apps.document": "GDOC",
                "application/vnd.google-apps.presentation": "GSLIDES",
                "application/pdf": "PDF",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "DOCX",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "XLSX",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation": "PPTX",
                "text/plain": "TXT",
            }
            file_type = ext_map.get(mime, fname.rsplit(".", 1)[-1].upper() if "." in fname else "FILE")

            # Build proper Google URL based on mime type
            if mime == "application/vnd.google-apps.spreadsheet":
                drive_url = f"https://docs.google.com/spreadsheets/d/{fid}"
            elif mime == "application/vnd.google-apps.document":
                drive_url = f"https://docs.google.com/document/d/{fid}"
            elif mime == "application/vnd.google-apps.presentation":
                drive_url = f"https://docs.google.com/presentation/d/{fid}"
            else:
                drive_url = f"https://drive.google.com/file/d/{fid}/view"

            FILE_REGISTRY[parent_id] = {
                "file_name": fname,
                "title": fname,
                "file_path": "",  # No local file for Drive imports
                "knowledge_base": knowledge_base,
                "uploaded_at": datetime.now().isoformat(),
                "section_doc_ids": doc_ids,
                "source_type": "google_drive",
                "drive_url": drive_url,
                "drive_file_id": fid,
                "file_type": file_type,
                "file_size_bytes": int(file_info.get("size", 0)) if file_info.get("size") else 0,
                "tags": product_tags,
                "source_id": source_id or None,
                "uploaded_by": uploaded_by or None,
                "description": "",
            }
            # Track document in source
            if source_id and source_id in DRIVE_SOURCES:
                DRIVE_SOURCES[source_id]["document_ids"].append(parent_id)
            registered_count += 1

        if registered_count > 0:
            save_registry()
        if source_id and source_id in DRIVE_SOURCES:
            save_sources()

        task["status"] = "completed"
        task["synced"] = result.get("synced", 0)
        task["skipped"] = result.get("skipped", 0)
        task["errors"] = result.get("errors", 0)
        task["details"] = result.get("details", [])
        task["files_processed"] = len(files)
        task["registered"] = registered_count
        task["status_detail"] = f"Hoàn tất: {result.get('synced', 0)} synced, {result.get('skipped', 0)} bỏ qua, {result.get('errors', 0)} lỗi"
        task["completed_at"] = datetime.now().isoformat()
    except Exception as e:
        import traceback
        task["status"] = "error"
        task["error"] = str(e)
        task["error_trace"] = traceback.format_exc()
        task["completed_at"] = datetime.now().isoformat()


async def _run_link_import(task_id: str, doc_id: str, dataset_id: str, knowledge_base: str, url: str, title: str, is_sheet: bool, source_id: str = "", product_tags: List[str] = [], uploaded_by: str = ""):
    """Background worker: import single Sheet or Doc."""
    task = IMPORT_TASKS[task_id]
    try:
        import sys
        sys.path.insert(0, str(Path(__file__).parent / "api"))
        task["status_detail"] = "Đang tải và xử lý tài liệu..."

        if is_sheet:
            from services.google_sheets_sync import GoogleSheetsSync
            syncer = GoogleSheetsSync(
                credentials_path=str(Path(__file__).parent / "config" / "google-credentials.json"),
                dify_base_url=DIFY_BASE_URL,
                dify_dataset_api_key=DIFY_DATASET_API_KEY,
            )
            result = await syncer.sync_sheet(
                spreadsheet_id=doc_id,
                dataset_id=dataset_id,
                title=title,
                force=True,
            )
            sections_count = result.get("sections", 0)
            doc_ids = result.get("doc_ids", [])

            # Register in FILE_REGISTRY
            parent_id = f"drive-sheet-{doc_id[:12]}"
            FILE_REGISTRY[parent_id] = {
                "file_name": title or doc_id,
                "title": title or doc_id,
                "file_path": "",
                "knowledge_base": knowledge_base,
                "uploaded_at": datetime.now().isoformat(),
                "section_doc_ids": doc_ids,
                "source_type": "google_sheet",
                "drive_url": url,
                "drive_file_id": doc_id,
                "file_type": "GSHEET",
                "file_size_bytes": 0,
                "tags": product_tags,
                "source_id": source_id or None,
                "uploaded_by": uploaded_by or None,
                "description": "",
            }
            save_registry()
            if source_id and source_id in DRIVE_SOURCES:
                DRIVE_SOURCES[source_id]["document_ids"].append(parent_id)
                save_sources()

            task["status"] = "completed"
            task["sections_count"] = sections_count
            task["completed_at"] = datetime.now().isoformat()
        else:
            from services.google_drive_sync import GoogleDriveSync
            syncer = GoogleDriveSync(
                credentials_path=str(Path(__file__).parent / "config" / "google-credentials.json"),
                dify_base_url=DIFY_BASE_URL,
                dify_dataset_api_key=DIFY_DATASET_API_KEY,
            )
            ext, content = await asyncio.to_thread(syncer._download_file, doc_id, "application/vnd.google-apps.document")
            doc_title = title or doc_id
            markdown = await asyncio.to_thread(syncer._preprocess_file, ext, content, doc_title)
            if not markdown:
                raise Exception("Failed to process Google Doc content")

            sections = await asyncio.to_thread(syncer._split_sections, markdown, doc_title)
            uploaded = 0
            uploaded_doc_ids = []
            async with httpx.AsyncClient(timeout=120.0) as client:
                for section in sections:
                    resp = await client.post(
                        f"{DIFY_BASE_URL}/datasets/{dataset_id}/document/create_by_text",
                        headers={
                            "Authorization": f"Bearer {DIFY_DATASET_API_KEY}",
                            "Content-Type": "application/json",
                        },
                        json={
                            "name": section["name"],
                            "text": section["text"],
                            "indexing_technique": "high_quality",
                            "process_rule": {"mode": "automatic"},
                        },
                    )
                    if resp.status_code == 200:
                        uploaded += 1
                        did = resp.json().get("document", {}).get("id")
                        if did:
                            uploaded_doc_ids.append(did)

            # Register in FILE_REGISTRY
            parent_id = f"drive-doc-{doc_id[:12]}"
            FILE_REGISTRY[parent_id] = {
                "file_name": doc_title,
                "title": doc_title,
                "file_path": "",
                "knowledge_base": knowledge_base,
                "uploaded_at": datetime.now().isoformat(),
                "section_doc_ids": uploaded_doc_ids,
                "source_type": "google_doc",
                "drive_url": url,
                "drive_file_id": doc_id,
                "file_type": "GDOC",
                "file_size_bytes": 0,
                "tags": product_tags,
                "source_id": source_id or None,
                "uploaded_by": uploaded_by or None,
                "description": "",
            }
            save_registry()
            if source_id and source_id in DRIVE_SOURCES:
                DRIVE_SOURCES[source_id]["document_ids"].append(parent_id)
                save_sources()

            task["status"] = "completed"
            task["title"] = doc_title
            task["sections_count"] = uploaded
            task["completed_at"] = datetime.now().isoformat()

    except Exception as e:
        task["status"] = "error"
        task["error"] = str(e)
        task["completed_at"] = datetime.now().isoformat()


@app.post("/api/v1/admin/knowledge/import-drive")
async def import_drive(req: ImportDriveRequest, request: Request):
    """Start async import of Google Drive folder. Returns task_id immediately."""
    m = re.search(r"/folders/([a-zA-Z0-9_-]+)", req.folder_url)
    if not m:
        raise HTTPException(400, "Invalid Google Drive folder URL. Expected a URL containing /folders/<id>")
    folder_id = m.group(1)

    dataset_id = DIFY_DATASET_IDS.get(req.knowledge_base)
    if not dataset_id:
        raise HTTPException(400, f"Unknown knowledge base '{req.knowledge_base}'. Valid: {list(DIFY_DATASET_IDS.keys())}")

    if not DIFY_DATASET_API_KEY:
        raise HTTPException(500, "DIFY_DATASET_API_KEY not configured")

    # Get current user for uploaded_by
    user = get_current_user(request)
    uploaded_by = user.get("sub", "")

    # Create source entry
    source_id = f"src-{uuid.uuid4().hex[:8]}"
    DRIVE_SOURCES[source_id] = {
        "id": source_id,
        "name": req.name or f"Drive folder {folder_id[:8]}",
        "description": req.description,
        "type": "folder",
        "url": req.folder_url,
        "folder_id": folder_id,
        "knowledge_base": req.knowledge_base,
        "product_tags": req.product_tags,
        "document_ids": [],
        "uploaded_by": uploaded_by,
        "created_at": datetime.now().isoformat(),
        "updated_at": datetime.now().isoformat(),
    }
    save_sources()

    task_id = str(uuid.uuid4())[:8]
    IMPORT_TASKS[task_id] = {
        "task_id": task_id,
        "type": "folder",
        "status": "importing",
        "status_detail": "Đang khởi tạo...",
        "url": req.folder_url,
        "title": req.name,
        "knowledge_base": req.knowledge_base,
        "source_id": source_id,
        "files_found": 0,
        "files_processed": 0,
        "synced": 0,
        "skipped": 0,
        "errors": 0,
        "details": [],
        "started_at": datetime.now().isoformat(),
        "completed_at": None,
    }

    asyncio.create_task(_run_drive_import(task_id, folder_id, dataset_id, req.knowledge_base, req.folder_url, source_id, req.product_tags, uploaded_by))

    return {"task_id": task_id, "status": "importing", "source_id": source_id, "message": "Import đang chạy nền..."}


@app.post("/api/v1/admin/knowledge/import-link")
async def import_link(req: ImportLinkRequest, request: Request):
    """Start async import of a single Google Sheet or Doc. Returns task_id immediately."""
    dataset_id = DIFY_DATASET_IDS.get(req.knowledge_base)
    if not dataset_id:
        raise HTTPException(400, f"Unknown knowledge base '{req.knowledge_base}'. Valid: {list(DIFY_DATASET_IDS.keys())}")

    if not DIFY_DATASET_API_KEY:
        raise HTTPException(500, "DIFY_DATASET_API_KEY not configured")

    is_sheet = "/spreadsheets/d/" in req.url
    is_doc = "/document/d/" in req.url

    if not is_sheet and not is_doc:
        raise HTTPException(400, "URL must be a Google Sheet (/spreadsheets/d/...) or Google Doc (/document/d/...)")

    id_match = re.search(r"/d/([a-zA-Z0-9_-]+)", req.url)
    if not id_match:
        raise HTTPException(400, "Could not extract document ID from URL")
    doc_id = id_match.group(1)

    # Get current user for uploaded_by
    user = get_current_user(request)
    uploaded_by = user.get("sub", "")

    # Create source entry
    source_id = f"src-{uuid.uuid4().hex[:8]}"
    link_type = "sheet" if is_sheet else "doc"
    DRIVE_SOURCES[source_id] = {
        "id": source_id,
        "name": req.title or doc_id,
        "description": req.description,
        "type": link_type,
        "url": req.url,
        "folder_id": doc_id,
        "knowledge_base": req.knowledge_base,
        "product_tags": req.product_tags,
        "document_ids": [],
        "uploaded_by": uploaded_by,
        "created_at": datetime.now().isoformat(),
        "updated_at": datetime.now().isoformat(),
    }
    save_sources()

    task_id = str(uuid.uuid4())[:8]
    IMPORT_TASKS[task_id] = {
        "task_id": task_id,
        "type": link_type,
        "status": "importing",
        "url": req.url,
        "title": req.title or doc_id,
        "knowledge_base": req.knowledge_base,
        "source_id": source_id,
        "sections_count": 0,
        "started_at": datetime.now().isoformat(),
        "completed_at": None,
    }

    asyncio.create_task(_run_link_import(task_id, doc_id, dataset_id, req.knowledge_base, req.url, req.title, is_sheet, source_id, req.product_tags, uploaded_by))

    return {"task_id": task_id, "status": "importing", "source_id": source_id, "message": "Import đang chạy nền..."}


class ImportWebUrlRequest(BaseModel):
    url: str
    knowledge_base: str
    title: str = ""
    description: str = ""
    product_tags: List[str] = []


async def _run_web_import(task_id: str, url: str, dataset_id: str, knowledge_base: str, title: str, source_id: str, product_tags: List[str], uploaded_by: str):
    """Background worker: fetch web page content and send to Dify."""
    task = IMPORT_TASKS[task_id]
    try:
        # Fetch the web page
        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
            resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0 (compatible; ViHAT-Bot/1.0)"})
            resp.raise_for_status()
            html = resp.text

        # Extract text from HTML
        from html.parser import HTMLParser
        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.result = []
                self.skip = False
            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style", "nav", "footer", "header"):
                    self.skip = True
            def handle_endtag(self, tag):
                if tag in ("script", "style", "nav", "footer", "header"):
                    self.skip = False
            def handle_data(self, data):
                if not self.skip:
                    text = data.strip()
                    if text:
                        self.result.append(text)

        extractor = TextExtractor()
        extractor.feed(html)
        page_text = "\n".join(extractor.result)

        # If text extraction failed (SPA/JS-rendered page), try meta tags
        if not page_text.strip() or len(page_text) < 50:
            meta_parts = []
            # Extract meta title
            meta_title = re.search(r'<meta\s+name="title"\s+content="([^"]+)"', html)
            if meta_title:
                meta_parts.append(f"Title: {meta_title.group(1)}")
            # Extract meta description
            meta_desc = re.search(r'<meta\s+name="description"\s+content="([^"]+)"', html)
            if meta_desc:
                meta_parts.append(f"Description: {meta_desc.group(1)}")
            # Extract og:title and og:description
            og_title = re.search(r'<meta\s+property="og:title"\s+content="([^"]+)"', html)
            if og_title and og_title.group(1) not in str(meta_parts):
                meta_parts.append(f"OG Title: {og_title.group(1)}")
            og_desc = re.search(r'<meta\s+property="og:description"\s+content="([^"]+)"', html)
            if og_desc and og_desc.group(1) not in str(meta_parts):
                meta_parts.append(f"OG Description: {og_desc.group(1)}")
            # Try <title> tag
            title_match = re.search(r'<title>([^<]+)</title>', html)
            if title_match:
                meta_parts.insert(0, f"Page Title: {title_match.group(1)}")
            # Try JSON-LD structured data
            json_ld = re.findall(r'<script[^>]*type="application/ld\+json"[^>]*>(.*?)</script>', html, re.DOTALL)
            for jld in json_ld:
                try:
                    ld_data = json.loads(jld)
                    if isinstance(ld_data, dict):
                        for key in ("description", "articleBody", "text", "name"):
                            if key in ld_data and ld_data[key]:
                                meta_parts.append(f"{key}: {ld_data[key]}")
                except Exception:
                    pass
            if meta_parts:
                page_text = "\n\n".join(meta_parts)
                print(f"[WebImport] SPA detected, using meta tags ({len(page_text)} chars)")

        if not page_text.strip():
            task["status"] = "error"
            task["error"] = "Không thể trích xuất nội dung từ URL (trang SPA/JS). Thử dùng Import từ Google Drive."
            task["completed_at"] = datetime.now().isoformat()
            return

        # Save as text file
        safe_name = re.sub(r'[^\w\s-]', '', title or "webpage").strip().replace(' ', '_')
        file_name = f"{safe_name}_{uuid.uuid4().hex[:6]}.txt"
        file_path = UPLOADS_DIR / file_name
        file_path.write_text(page_text, encoding="utf-8")

        # Send to Dify
        file_size = file_path.stat().st_size
        DIFY_FILE_LIMIT = 15 * 1024 * 1024
        section_doc_ids = []

        if file_size <= DIFY_FILE_LIMIT and DIFY_DATASET_API_KEY:
            try:
                async with httpx.AsyncClient(timeout=120.0) as client2:
                    with open(file_path, "rb") as f:
                        files = {"file": (file_name, f, "text/plain")}
                        data = {"indexing_technique": "high_quality", "process_rule": json.dumps({"mode": "automatic"})}
                        headers = {"Authorization": f"Bearer {DIFY_DATASET_API_KEY}"}
                        resp2 = await client2.post(
                            f"https://api.dify.ai/v1/datasets/{dataset_id}/document/create_by_file",
                            files=files, data=data, headers=headers,
                        )
                    if resp2.status_code in (200, 201):
                        doc_data = resp2.json()
                        doc = doc_data.get("document", {})
                        section_doc_ids.append(doc.get("id", ""))
            except Exception as e:
                print(f"[WebImport] Dify upload error: {e}")

        # Register in FILE_REGISTRY
        parent_id = f"web-{uuid.uuid4().hex[:8]}"
        FILE_REGISTRY[parent_id] = {
            "id": parent_id,
            "knowledge_base": knowledge_base,
            "title": title or url,
            "description": "",
            "file_name": file_name,
            "file_path": str(file_path),
            "file_type": "web",
            "tags": product_tags,
            "section_doc_ids": section_doc_ids,
            "status": "active",
            "uploaded_at": datetime.now().isoformat(),
            "uploaded_by": uploaded_by,
            "source_id": source_id,
            "drive_url": url,
        }
        save_registry()

        # Update source
        if source_id and source_id in DRIVE_SOURCES:
            DRIVE_SOURCES[source_id]["document_ids"].append(parent_id)
            save_sources()

        task["status"] = "completed"
        task["sections_count"] = 1
        task["completed_at"] = datetime.now().isoformat()
        print(f"[WebImport] Completed: {title or url} ({len(page_text)} chars)")

    except Exception as e:
        import traceback
        task["status"] = "error"
        task["error"] = str(e)
        task["completed_at"] = datetime.now().isoformat()
        print(f"[WebImport] Error: {e}\n{traceback.format_exc()}")


@app.post("/api/v1/admin/knowledge/import-web")
async def import_web_url(req: ImportWebUrlRequest, request: Request):
    """Import knowledge from a web URL (fetch content, extract text, send to Dify)."""
    dataset_id = DIFY_DATASET_IDS.get(req.knowledge_base)
    if not dataset_id:
        raise HTTPException(400, f"Unknown knowledge base '{req.knowledge_base}'")

    if not req.url.startswith(("http://", "https://")):
        raise HTTPException(400, "URL must start with http:// or https://")

    # Don't allow Google Drive URLs — use import-link for those
    if "drive.google.com" in req.url or "docs.google.com" in req.url:
        raise HTTPException(400, "Dùng 'Import Link' cho Google Drive/Docs/Sheets")

    user = get_current_user(request)
    uploaded_by = user.get("sub", "")

    source_id = f"src-{uuid.uuid4().hex[:8]}"
    DRIVE_SOURCES[source_id] = {
        "id": source_id,
        "name": req.title or req.url,
        "description": req.description,
        "type": "web",
        "url": req.url,
        "folder_id": "",
        "knowledge_base": req.knowledge_base,
        "product_tags": req.product_tags,
        "document_ids": [],
        "uploaded_by": uploaded_by,
        "created_at": datetime.now().isoformat(),
        "updated_at": datetime.now().isoformat(),
    }
    save_sources()

    task_id = str(uuid.uuid4())[:8]
    IMPORT_TASKS[task_id] = {
        "task_id": task_id,
        "type": "web",
        "status": "importing",
        "url": req.url,
        "title": req.title or req.url,
        "knowledge_base": req.knowledge_base,
        "source_id": source_id,
        "sections_count": 0,
        "started_at": datetime.now().isoformat(),
        "completed_at": None,
    }

    asyncio.create_task(_run_web_import(task_id, req.url, dataset_id, req.knowledge_base, req.title, source_id, req.product_tags, uploaded_by))

    return {"task_id": task_id, "status": "importing", "source_id": source_id, "message": "Đang import từ web..."}


@app.get("/api/v1/admin/knowledge/import-tasks")
async def list_import_tasks():
    """List all import tasks (active and recent)."""
    tasks = sorted(IMPORT_TASKS.values(), key=lambda t: t.get("started_at", ""), reverse=True)
    return tasks


@app.get("/api/v1/admin/knowledge/drive-status")
async def drive_status():
    """Get sync status of all imported Drive sources (folders and sheets)."""
    result = {"folders": [], "sheets": []}

    try:
        import sys
        sys.path.insert(0, str(Path(__file__).parent / "api"))

        # Drive folders status
        try:
            from services.google_drive_sync import GoogleDriveSync
            drive_syncer = GoogleDriveSync(
                credentials_path=str(Path(__file__).parent / "config" / "google-credentials.json"),
                dify_base_url=DIFY_BASE_URL,
                dify_dataset_api_key=DIFY_DATASET_API_KEY,
            )
            result["folders"] = drive_syncer.get_sync_status()
        except Exception:
            pass

        # Sheets status
        try:
            from services.google_sheets_sync import GoogleSheetsSync
            sheets_syncer = GoogleSheetsSync(
                credentials_path=str(Path(__file__).parent / "config" / "google-credentials.json"),
                dify_base_url=DIFY_BASE_URL,
                dify_dataset_api_key=DIFY_DATASET_API_KEY,
            )
            result["sheets"] = sheets_syncer.get_sync_status()
        except Exception:
            pass

    except Exception:
        pass

    return result


# ═══════════════════════════════════════════════════════════════════════════════
# PROPOSAL BUILDER (Sales Proposal Generator)
# ═══════════════════════════════════════════════════════════════════════════════

PROPOSALS_DIR = Path(__file__).parent / "data" / "proposals"
PROPOSALS_DIR.mkdir(parents=True, exist_ok=True)
TEMPLATES_DIR = Path(__file__).parent / "data" / "templates"
TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)
PRODUCTS_CONFIG_PATH = Path(__file__).parent / "data" / "products_config.json"
RFI_TEMPLATES_PATH = Path(__file__).parent / "data" / "rfi_templates.json"

PROPOSAL_TASKS: dict = {}  # task_id -> task info
PROPOSAL_TASKS_PATH = Path(__file__).parent / "data" / "proposal_tasks.json"

def _load_proposal_tasks():
    """Load proposal tasks from JSON file."""
    global PROPOSAL_TASKS
    if PROPOSAL_TASKS_PATH.exists():
        try:
            PROPOSAL_TASKS = json.loads(PROPOSAL_TASKS_PATH.read_text(encoding="utf-8"))
        except Exception:
            PROPOSAL_TASKS = {}

def _save_proposal_tasks():
    """Persist proposal tasks to JSON file."""
    try:
        # Only save a safe subset (no rfi_answers/company_info to save space)
        safe = {}
        for tid, t in PROPOSAL_TASKS.items():
            safe[tid] = {k: v for k, v in t.items() if k not in ("rfi_answers", "company_info")}
        PROPOSAL_TASKS_PATH.write_text(json.dumps(safe, ensure_ascii=False, indent=2, default=str), encoding="utf-8")
    except Exception as e:
        print(f"[Proposal] Save tasks error: {e}")

_load_proposal_tasks()

def _get_legal_entities() -> list:
    """Get legal entities from current tenant config."""
    tenant = TENANTS.get(DEFAULT_TENANT_ID, {})
    return tenant.get("config", {}).get("legal_entities", [])

LEGAL_ENTITIES = _get_legal_entities()  # backward compat

def _load_products_config() -> list:
    """Legacy: return solutions as selectable items, grouped under products.
    Solutions ARE what customers search for. Products with no solutions show directly."""
    # Collect all solution slugs to know which products are covered
    solution_slugs = set()
    result = []
    for s in sorted(SOLUTIONS.values(), key=lambda x: x.get("sort_order", 0)):
        if s.get("status") != "active":
            continue
        parent = PRODUCTS.get(s.get("product_id", ""))
        parent_name = parent["name"] if parent else ""
        result.append({
            "id": s["slug"],
            "label": f"{s['name']}" + (f" ({parent_name})" if parent_name and parent_name != s["name"] else ""),
            "type": "solution",
            "product_slug": parent["slug"] if parent else "",
        })
        solution_slugs.add(s["slug"])
    # Add products that have NO solutions (standalone)
    for p in PRODUCTS.values():
        if p.get("status") == "active" and p["slug"] not in solution_slugs:
            # Check if any solution references this product
            has_solutions = any(
                s.get("product_id") == p["id"] and s.get("status") == "active"
                for s in SOLUTIONS.values()
            )
            if not has_solutions:
                result.append({"id": p["slug"], "label": p["name"], "type": "product"})
    return result

def _save_products_config(data: list):
    """Legacy: no-op, products now stored in-memory."""
    pass

def _load_rfi_templates() -> dict:
    try:
        return json.loads(RFI_TEMPLATES_PATH.read_text(encoding="utf-8"))
    except Exception:
        return {}

def _save_rfi_templates(data: dict):
    RFI_TEMPLATES_PATH.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


# --- Legacy Products Config API (backward compat) ---

@app.get("/api/v1/proposals/products")
async def get_products_legacy():
    """Legacy endpoint: returns simple {id, label} list."""
    return _load_products_config()

@app.put("/api/v1/proposals/products")
async def update_products_legacy(products: list):
    """Legacy: no-op for backward compat."""
    return {"status": "ok", "count": len(products)}


# ═══════════════════════════════════════════════════════════════════════════════
# TENANT API
# ═══════════════════════════════════════════════════════════════════════════════

def _resolve_tenant(request: Request):
    """Resolve tenant from X-Tenant-Slug header, falling back to default."""
    slug = request.headers.get("X-Tenant-Slug", "")
    if slug:
        for t in TENANTS.values():
            if t["slug"] == slug:
                return t
    return TENANTS.get(DEFAULT_TENANT_ID)

@app.get("/api/v1/tenant/by-slug/{slug}")
async def get_tenant_by_slug(slug: str):
    """Get tenant info by slug (public — used by frontend to resolve subdomain)."""
    for t in TENANTS.values():
        if t["slug"] == slug:
            # Return safe public info (no sensitive config)
            return {
                "id": t["id"],
                "slug": t["slug"],
                "name": t["name"],
                "logo_url": t.get("logo_url"),
                "primary_color": t.get("primary_color"),
                "config": {
                    "departments": t.get("config", {}).get("departments", []),
                    "features": t.get("config", {}).get("features", {}),
                },
                "is_active": t.get("is_active", True),
            }
    raise HTTPException(404, "Tenant not found")

@app.get("/api/v1/tenant")
async def get_current_tenant(request: Request):
    """Get current tenant info (branding, config). Resolves from X-Tenant-Slug header."""
    tenant = _resolve_tenant(request)
    if not tenant:
        raise HTTPException(404, "Tenant not found")
    return tenant

@app.put("/api/v1/tenant")
async def update_tenant(body: dict, request: Request):
    """Update current tenant settings (super_admin only)."""
    tenant = _resolve_tenant(request)
    if not tenant:
        raise HTTPException(404, "Tenant not found")
    for key in ["name", "logo_url", "primary_color", "config"]:
        if key in body:
            tenant[key] = body[key]
    tenant["updated_at"] = datetime.utcnow().isoformat()
    return tenant


# ═══════════════════════════════════════════════════════════════════════════════
# RICH PRODUCTS API (with versioning)
# ═══════════════════════════════════════════════════════════════════════════════

class ProductCreateRequest(BaseModel):
    slug: str
    name: str
    short_description: str = ""
    full_description: str = ""
    features: list = []
    use_cases: list = []
    target_industries: list = []
    pricing_model: str = ""
    competitive_advantages: list = []
    integration_options: list = []
    sort_order: int = 0

class ProductUpdateRequest(BaseModel):
    name: str = ""
    short_description: str = ""
    full_description: str = ""
    features: list = []
    use_cases: list = []
    target_industries: list = []
    pricing_model: str = ""
    competitive_advantages: list = []
    integration_options: list = []
    sort_order: int = 0
    version_label: str = ""
    change_summary: str = ""

@app.get("/api/v1/products")
async def list_products():
    """List all active products for current tenant."""
    products = [p for p in PRODUCTS.values()
                if p.get("tenant_id") == DEFAULT_TENANT_ID and p.get("status") != "deprecated"]
    products.sort(key=lambda x: x.get("sort_order", 0))
    # Add version count, solutions count, and related docs count
    for p in products:
        p["version_count"] = len(PRODUCT_VERSIONS.get(p["id"], []))
        # Get solutions for this product
        product_solutions = [s for s in SOLUTIONS.values()
                             if s.get("product_id") == p["id"] and s.get("status") == "active"]
        p["solutions_count"] = len(product_solutions)
        p["solutions"] = sorted(product_solutions, key=lambda x: x.get("sort_order", 0))
        # Count knowledge docs matching product slug OR any solution slug
        all_slugs = {p["slug"]}
        all_slugs.update(s["slug"] for s in product_solutions)
        p["related_docs_count"] = sum(
            1 for d in FILE_REGISTRY.values()
            if all_slugs.intersection(set(d.get("tags") or []))
        )
    return products

@app.get("/api/v1/products/{product_id}")
async def get_product(product_id: str):
    """Get product detail by ID or slug."""
    product = PRODUCTS.get(product_id)
    if not product:
        # Try by slug
        product = next((p for p in PRODUCTS.values() if p["slug"] == product_id), None)
    if not product:
        raise HTTPException(404, "Product not found")
    result = {**product}
    result["version_count"] = len(PRODUCT_VERSIONS.get(product["id"], []))
    product_solutions = [s for s in SOLUTIONS.values()
                         if s.get("product_id") == product["id"] and s.get("status") == "active"]
    result["solutions_count"] = len(product_solutions)
    result["solutions"] = sorted(product_solutions, key=lambda x: x.get("sort_order", 0))
    all_slugs = {product["slug"]}
    all_slugs.update(s["slug"] for s in product_solutions)
    result["related_docs_count"] = sum(
        1 for d in FILE_REGISTRY.values()
        if all_slugs.intersection(set(d.get("tags") or []))
    )
    return result

@app.post("/api/v1/products")
async def create_product(req: ProductCreateRequest):
    """Create a new product."""
    if not req.name.strip():
        raise HTTPException(400, "Tên sản phẩm không được để trống")
    # Auto-generate slug from name if empty
    if not req.slug.strip():
        import unicodedata
        slug = unicodedata.normalize("NFD", req.name.lower())
        slug = "".join(c for c in slug if not unicodedata.combining(c))
        slug = slug.replace("đ", "d").replace(" ", "_")
        slug = re.sub(r"[^a-z0-9_]", "", slug)
        req.slug = slug or f"product_{uuid.uuid4().hex[:6]}"
    # Check slug uniqueness within tenant
    if any(p["slug"] == req.slug and p["tenant_id"] == DEFAULT_TENANT_ID for p in PRODUCTS.values()):
        raise HTTPException(400, f"Slug '{req.slug}' đã tồn tại. Vui lòng chọn slug khác.")
    now = datetime.utcnow().isoformat()
    pid = str(uuid.uuid4())[:8]
    PRODUCTS[pid] = {
        "id": pid,
        "tenant_id": DEFAULT_TENANT_ID,
        "slug": req.slug,
        "name": req.name,
        "short_description": req.short_description,
        "full_description": req.full_description,
        "features": req.features,
        "use_cases": req.use_cases,
        "target_industries": req.target_industries,
        "pricing_model": req.pricing_model,
        "competitive_advantages": req.competitive_advantages,
        "integration_options": req.integration_options,
        "status": "active",
        "sort_order": req.sort_order,
        "created_at": now,
        "updated_at": now,
    }
    PRODUCT_VERSIONS[pid] = []
    _save_products()
    return PRODUCTS[pid]

@app.put("/api/v1/products/{product_id}")
async def update_product(product_id: str, req: ProductUpdateRequest):
    """Update product and auto-save version history."""
    product = PRODUCTS.get(product_id)
    if not product:
        product = next((p for p in PRODUCTS.values() if p["slug"] == product_id), None)
    if not product:
        raise HTTPException(404, "Product not found")

    pid = product["id"]
    # Save current state as version before updating
    save_version(pid, PRODUCT_VERSIONS, product,
                 version_label=req.version_label, change_summary=req.change_summary)

    # Update fields (only non-empty values)
    if req.name:
        product["name"] = req.name
    product["short_description"] = req.short_description
    product["full_description"] = req.full_description
    product["features"] = req.features
    product["use_cases"] = req.use_cases
    product["target_industries"] = req.target_industries
    product["pricing_model"] = req.pricing_model
    product["competitive_advantages"] = req.competitive_advantages
    product["integration_options"] = req.integration_options
    if req.sort_order is not None:
        product["sort_order"] = req.sort_order
    product["updated_at"] = datetime.utcnow().isoformat()

    _save_products()
    _save_product_versions()
    return product

@app.delete("/api/v1/products/{product_id}")
async def delete_product(product_id: str):
    """Soft-delete product (set status=deprecated)."""
    product = PRODUCTS.get(product_id)
    if not product:
        product = next((p for p in PRODUCTS.values() if p["slug"] == product_id), None)
    if not product:
        raise HTTPException(404, "Product not found")
    product["status"] = "deprecated"
    product["updated_at"] = datetime.utcnow().isoformat()
    _save_products()
    return {"status": "ok", "id": product["id"]}

@app.get("/api/v1/products/{product_id}/versions")
async def list_product_versions(product_id: str):
    """List version history for a product."""
    product = PRODUCTS.get(product_id) or next((p for p in PRODUCTS.values() if p["slug"] == product_id), None)
    if not product:
        raise HTTPException(404, "Product not found")
    versions = PRODUCT_VERSIONS.get(product["id"], [])
    return sorted(versions, key=lambda v: v["version_number"], reverse=True)

@app.get("/api/v1/products/{product_id}/versions/{version_number}")
async def get_product_version(product_id: str, version_number: int):
    """Get a specific version of a product."""
    product = PRODUCTS.get(product_id) or next((p for p in PRODUCTS.values() if p["slug"] == product_id), None)
    if not product:
        raise HTTPException(404, "Product not found")
    versions = PRODUCT_VERSIONS.get(product["id"], [])
    version = next((v for v in versions if v["version_number"] == version_number), None)
    if not version:
        raise HTTPException(404, f"Version {version_number} not found")
    return version

@app.post("/api/v1/products/{product_id}/versions/{version_number}/restore")
async def restore_product_version(product_id: str, version_number: int):
    """Restore product to an older version."""
    product = PRODUCTS.get(product_id) or next((p for p in PRODUCTS.values() if p["slug"] == product_id), None)
    if not product:
        raise HTTPException(404, "Product not found")
    pid = product["id"]
    versions = PRODUCT_VERSIONS.get(pid, [])
    version = next((v for v in versions if v["version_number"] == version_number), None)
    if not version:
        raise HTTPException(404, f"Version {version_number} not found")

    # Save current state as new version before restoring
    save_version(pid, PRODUCT_VERSIONS, product,
                 version_label=f"Before restore to v{version_number}",
                 change_summary=f"Auto-saved before restoring to version {version_number}")

    # Restore from snapshot
    snapshot = version["snapshot"]
    for key, val in snapshot.items():
        if key not in ("id", "tenant_id", "created_at"):
            product[key] = val
    product["updated_at"] = datetime.utcnow().isoformat()

    return product

@app.get("/api/v1/products/{product_id}/documents")
async def get_product_documents(product_id: str):
    """Get knowledge documents related to this product (by tag)."""
    product = PRODUCTS.get(product_id) or next((p for p in PRODUCTS.values() if p["slug"] == product_id), None)
    if not product:
        raise HTTPException(404, "Product not found")
    # Collect all slugs: product slug + all its solution slugs
    product_slugs = {product["slug"]}
    product_slugs.update(
        s["slug"] for s in SOLUTIONS.values()
        if s.get("product_id") == product["id"] and s.get("status") == "active"
    )
    docs = [d for d in FILE_REGISTRY.values()
            if product_slugs.intersection(set(d.get("tags") or []))]
    return docs


# ═══════════════════════════════════════════════════════════════════════════════
# SOLUTIONS API
# ═══════════════════════════════════════════════════════════════════════════════

class SolutionCreateRequest(BaseModel):
    name: str
    slug: str = ""
    description: str = ""
    product_id: str
    aliases: list = []
    sort_order: int = 0

class SolutionUpdateRequest(BaseModel):
    name: str = ""
    slug: str = ""
    description: str = ""
    product_id: str = ""
    aliases: list = []
    sort_order: int = 0

@app.get("/api/v1/solutions")
async def list_solutions():
    """List all active solutions with parent product info."""
    solutions = [s for s in SOLUTIONS.values()
                 if s.get("tenant_id") == DEFAULT_TENANT_ID and s.get("status") != "deprecated"]
    solutions.sort(key=lambda x: x.get("sort_order", 0))
    for s in solutions:
        parent = PRODUCTS.get(s.get("product_id", ""))
        s["product_name"] = parent["name"] if parent else ""
        s["product_slug"] = parent["slug"] if parent else ""
        # Count docs tagged with this solution slug
        s["related_docs_count"] = sum(
            1 for d in FILE_REGISTRY.values()
            if s["slug"] in (d.get("tags") or [])
        )
    return solutions

@app.get("/api/v1/solutions/{solution_id}")
async def get_solution(solution_id: str):
    """Get solution by ID or slug."""
    solution = SOLUTIONS.get(solution_id)
    if not solution:
        solution = next((s for s in SOLUTIONS.values() if s["slug"] == solution_id), None)
    if not solution:
        raise HTTPException(404, "Solution not found")
    result = {**solution}
    parent = PRODUCTS.get(result.get("product_id", ""))
    result["product_name"] = parent["name"] if parent else ""
    result["product_slug"] = parent["slug"] if parent else ""
    return result

@app.post("/api/v1/solutions")
async def create_solution(req: SolutionCreateRequest):
    """Create a new solution linked to a product."""
    # Validate product exists
    if not PRODUCTS.get(req.product_id):
        raise HTTPException(400, f"Product '{req.product_id}' not found")
    # Auto-generate slug
    slug = req.slug.strip()
    if not slug:
        slug = unicodedata.normalize("NFD", req.name.lower())
        slug = "".join(c for c in slug if not unicodedata.combining(c))
        slug = slug.replace("đ", "d").replace(" ", "_").replace("/", "_")
        slug = slug or f"solution_{uuid.uuid4().hex[:6]}"
    # Check slug uniqueness across products and solutions
    existing_slugs = {p["slug"] for p in PRODUCTS.values()}
    existing_slugs.update(s["slug"] for s in SOLUTIONS.values())
    if slug in existing_slugs:
        raise HTTPException(400, f"Slug '{slug}' already exists")
    now = datetime.utcnow().isoformat()
    sid = str(uuid.uuid4())[:8]
    SOLUTIONS[sid] = {
        "id": sid,
        "tenant_id": DEFAULT_TENANT_ID,
        "name": req.name,
        "slug": slug,
        "description": req.description,
        "product_id": req.product_id,
        "aliases": req.aliases,
        "status": "active",
        "sort_order": req.sort_order,
        "created_at": now,
        "updated_at": now,
    }
    _save_solutions()
    return SOLUTIONS[sid]

@app.put("/api/v1/solutions/{solution_id}")
async def update_solution(solution_id: str, req: SolutionUpdateRequest):
    """Update a solution."""
    solution = SOLUTIONS.get(solution_id)
    if not solution:
        raise HTTPException(404, "Solution not found")
    if req.name:
        solution["name"] = req.name
    if req.slug:
        solution["slug"] = req.slug
    if req.description:
        solution["description"] = req.description
    if req.product_id:
        if not PRODUCTS.get(req.product_id):
            raise HTTPException(400, f"Product '{req.product_id}' not found")
        solution["product_id"] = req.product_id
    if req.aliases is not None:
        solution["aliases"] = req.aliases
    if req.sort_order is not None:
        solution["sort_order"] = req.sort_order
    solution["updated_at"] = datetime.utcnow().isoformat()
    _save_solutions()
    return solution

@app.delete("/api/v1/solutions/{solution_id}")
async def delete_solution(solution_id: str):
    """Soft-delete a solution."""
    solution = SOLUTIONS.get(solution_id)
    if not solution:
        raise HTTPException(404, "Solution not found")
    solution["status"] = "deprecated"
    solution["updated_at"] = datetime.utcnow().isoformat()
    _save_solutions()
    return {"ok": True}


# --- RFI Templates API ---

@app.get("/api/v1/proposals/rfi")
async def list_rfi_templates():
    templates = _load_rfi_templates()
    return {k: {"label": v["label"], "questions_count": len(v.get("questions", []))} for k, v in templates.items()}

@app.get("/api/v1/proposals/rfi/{industry}")
async def get_rfi_template(industry: str):
    templates = _load_rfi_templates()
    if industry not in templates:
        raise HTTPException(404, f"RFI template '{industry}' not found")
    return templates[industry]

@app.post("/api/v1/proposals/rfi")
async def create_rfi_template(body: dict):
    industry = body.get("industry", "")
    label = body.get("label", "")
    questions = body.get("questions", [])
    if not industry or not label:
        raise HTTPException(400, "industry and label are required")
    templates = _load_rfi_templates()
    templates[industry] = {"label": label, "questions": questions}
    _save_rfi_templates(templates)
    return {"status": "ok", "industry": industry}

@app.put("/api/v1/proposals/rfi/{industry}")
async def update_rfi_template(industry: str, body: dict):
    templates = _load_rfi_templates()
    if industry not in templates:
        raise HTTPException(404, f"RFI template '{industry}' not found")
    if "label" in body:
        templates[industry]["label"] = body["label"]
    if "questions" in body:
        templates[industry]["questions"] = body["questions"]
    _save_rfi_templates(templates)
    return {"status": "ok", "industry": industry}

@app.delete("/api/v1/proposals/rfi/{industry}")
async def delete_rfi_template(industry: str):
    templates = _load_rfi_templates()
    if industry not in templates:
        raise HTTPException(404, f"RFI template '{industry}' not found")
    del templates[industry]
    _save_rfi_templates(templates)
    return {"status": "ok"}


# --- Legal Entities ---

@app.get("/api/v1/proposals/entities")
async def get_legal_entities():
    return _get_legal_entities()


# --- Company Lookup ---

class CompanyLookupRequest(BaseModel):
    tax_code: str = ""
    website: str = ""
    company_name: str = ""

@app.post("/api/v1/proposals/lookup-company")
async def lookup_company(req: CompanyLookupRequest):
    """Lookup company info via tax code or website."""
    info = {"company_name": req.company_name, "source": "manual"}

    # Try tax code lookup
    if req.tax_code:
        try:
            async with httpx.AsyncClient(timeout=10.0) as client:
                resp = await client.get(f"https://api.vietqr.io/v2/business/{req.tax_code}")
                if resp.status_code == 200:
                    data = resp.json().get("data", {})
                    if data:
                        info.update({
                            "company_name": data.get("name", req.company_name),
                            "short_name": data.get("shortName", ""),
                            "address": data.get("address", ""),
                            "source": "tax_lookup",
                        })
        except Exception as e:
            print(f"[CompanyLookup] Tax code lookup error: {e}")

    # Try website analysis via Claude
    if req.website and ANTHROPIC_API_KEY:
        try:
            # Fetch website content
            async with httpx.AsyncClient(timeout=15.0, follow_redirects=True) as client:
                web_resp = await client.get(req.website if req.website.startswith("http") else f"https://{req.website}")
                if web_resp.status_code == 200:
                    # Extract text (simple, first 3000 chars)
                    import re as _re
                    text = _re.sub(r'<[^>]+>', ' ', web_resp.text)
                    text = _re.sub(r'\s+', ' ', text)[:3000]

                    # Ask Claude to analyze
                    ai_resp = await client.post(
                        "https://api.anthropic.com/v1/messages",
                        headers={
                            "x-api-key": ANTHROPIC_API_KEY,
                            "anthropic-version": "2023-06-01",
                            "content-type": "application/json",
                        },
                        json={
                            "model": "claude-sonnet-4-20250514",
                            "max_tokens": 500,
                            "system": "Bạn là trợ lý phân tích doanh nghiệp. Trích xuất thông tin từ nội dung website. Trả JSON: {\"industry\": \"ngành nghề\", \"description\": \"mô tả ngắn\", \"products_services\": \"sản phẩm/dịch vụ chính\", \"company_size_estimate\": \"ước tính quy mô\"}. Chỉ trả JSON, không giải thích.",
                            "messages": [{"role": "user", "content": f"Website: {req.website}\n\nNội dung:\n{text}"}],
                        },
                    )
                    if ai_resp.status_code == 200:
                        ai_text = ai_resp.json()["content"][0]["text"].strip()
                        # Try parse JSON
                        try:
                            parsed = json.loads(ai_text)
                            info["website_analysis"] = parsed
                            info["source"] = "website_analysis"
                        except Exception:
                            info["website_analysis"] = {"raw": ai_text}
        except Exception as e:
            print(f"[CompanyLookup] Website analysis error: {e}")

    return info


# --- Parse Brief with AI ---

class ParseBriefRequest(BaseModel):
    brief: str
    industry: str
    products: list = []

@app.post("/api/v1/proposals/parse-brief")
async def parse_brief(req: ParseBriefRequest):
    """AI reads brief text and fills RFI answers."""
    templates = _load_rfi_templates()
    rfi = templates.get(req.industry, templates.get("chung", {}))
    questions = rfi.get("questions", [])

    if not ANTHROPIC_API_KEY:
        return {"answers": {}, "message": "No ANTHROPIC_API_KEY configured"}

    if not req.brief.strip():
        return {"answers": {}, "message": "No brief provided"}

    # Build question list for Claude
    q_list = "\n".join([f"- {q['id']}: {q['label']} (type: {q['type']})" for q in questions])
    products_text = ", ".join(req.products) if req.products else "chưa chọn"

    prompt = f"""Đọc brief sau từ khách hàng và điền vào các trường RFI bên dưới.
Sản phẩm quan tâm: {products_text}

Brief:
\"\"\"
{req.brief}
\"\"\"

Các trường RFI cần điền:
{q_list}

Trả về JSON object mapping question_id -> giá trị.
- Với type "number": trả số
- Với type "text" hoặc "textarea": trả chuỗi
- Với type "select": trả 1 giá trị phù hợp
- Với type "multi_select": trả array các giá trị phù hợp
- Nếu brief không đề cập, trả "" (chuỗi rỗng)
Chỉ trả JSON, không giải thích."""

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            resp = await client.post(
                "https://api.anthropic.com/v1/messages",
                headers={
                    "x-api-key": ANTHROPIC_API_KEY,
                    "anthropic-version": "2023-06-01",
                    "content-type": "application/json",
                },
                json={
                    "model": "claude-sonnet-4-20250514",
                    "max_tokens": 1000,
                    "messages": [{"role": "user", "content": prompt}],
                },
            )
            if resp.status_code == 200:
                text = resp.json()["content"][0]["text"].strip()
                # Extract JSON from response
                if "```" in text:
                    text = text.split("```")[1]
                    if text.startswith("json"):
                        text = text[4:]
                    text = text.strip()
                answers = json.loads(text)
                return {"answers": answers}
            else:
                return {"answers": {}, "message": f"AI error: {resp.status_code}"}
    except Exception as e:
        print(f"[ParseBrief] Error: {e}")
        return {"answers": {}, "message": str(e)}


# --- Proposal Generation (Async Background) ---

class ProposalGenerateRequest(BaseModel):
    customer_name: str
    industry: str = "chung"
    products: list = []
    rfi_answers: dict = {}
    company_info: dict = {}
    legal_entity: str = "omijsc"
    output_format: str = "pptx"  # pptx or docx
    brief_text: str = ""  # Original brief from user


async def _generate_proposal_content(task: dict) -> dict:
    """Call Claude API to generate proposal content with rich product knowledge."""
    # Get rich product info for selected products
    # Resolve solution slugs to parent product IDs
    selected_slugs = set(task.get("products", []))
    resolved_product_ids = set()
    selected_solution_names = []
    for s in SOLUTIONS.values():
        if s["slug"] in selected_slugs and s.get("status") == "active":
            resolved_product_ids.add(s.get("product_id", ""))
            selected_solution_names.append(s["name"])
    # Also match directly by product slug/name
    for p in PRODUCTS.values():
        if p.get("status") == "active" and (p["slug"] in selected_slugs or p["name"] in selected_slugs):
            resolved_product_ids.add(p["id"])
    rich_products = [p for p in PRODUCTS.values()
                     if p["id"] in resolved_product_ids and p.get("status") == "active"]

    # Build detailed product knowledge for Claude
    products_detail = ""
    for p in rich_products:
        products_detail += f"""
### {p['name']} ({p['slug']})
- Mô tả: {p.get('full_description', p.get('short_description', ''))}
- Tính năng: {', '.join(p.get('features', []))}
- Use cases: {', '.join(p.get('use_cases', []))}
- Giá: {p.get('pricing_model', 'Liên hệ')}
- Điểm mạnh: {', '.join(p.get('competitive_advantages', []))}
- Tích hợp: {', '.join(p.get('integration_options', []))}
"""

    # Query knowledge base for relevant documents
    kb_context = ""
    all_solution_slugs = selected_slugs.copy()
    for pid in resolved_product_ids:
        prod = PRODUCTS.get(pid)
        if prod:
            all_solution_slugs.add(prod["slug"])
    # Find docs tagged with any of the selected slugs OR tagged "chung" (global knowledge)
    all_solution_slugs.add("chung")  # Always include global knowledge docs
    relevant_docs = []
    for doc in FILE_REGISTRY.values():
        doc_tags = set(doc.get("tags") or [])
        if doc_tags.intersection(all_solution_slugs) and doc.get("description"):
            relevant_docs.append(doc)
    if relevant_docs:
        kb_context = "\n=== TRI THỨC BỔ SUNG TỪ KHO DỮ LIỆU ===\n"
        for doc in relevant_docs[:10]:  # Limit to 10 docs
            kb_context += f"\n--- {doc.get('title', doc.get('file_name', ''))} ---\n"
            kb_context += f"{doc.get('description', '')}\n"

    # All available products for reference
    all_products = [p for p in PRODUCTS.values() if p.get("status") == "active"]
    all_products_brief = ", ".join([f"{p['name']}" for p in all_products])

    # Brief text from user
    brief_text = task.get("brief_text", "")
    rfi_text = "\n".join([f"- {k}: {v}" for k, v in task.get("rfi_answers", {}).items() if v])
    company_text = json.dumps(task.get("company_info", {}), ensure_ascii=False)

    # Solution names for title
    solution_names_text = ", ".join(selected_solution_names) if selected_solution_names else "giải pháp phù hợp"

    # Get tenant info for branding
    tenant = TENANTS.get(DEFAULT_TENANT_ID, {})
    company_name = tenant.get("name", "ViHAT Group")

    prompt = f"""Bạn là chuyên gia tư vấn giải pháp CPaaS & Contact Center tại {company_name}. Hãy tạo nội dung PROPOSAL CHI TIẾT và CHUYÊN NGHIỆP cho khách hàng.

=== BRIEF TỪ KHÁCH HÀNG / SALES ===
{brief_text if brief_text.strip() else 'Không có brief cụ thể'}

QUAN TRỌNG: Phải ĐỌC KỸ brief và PHÂN TÍCH từng yêu cầu cụ thể của khách hàng. Proposal phải ĐÁP ỨNG TRỰC TIẾP mọi nhu cầu được đề cập trong brief. Không bỏ sót bất kỳ yêu cầu nào.

=== THÔNG TIN KHÁCH HÀNG ===
- Tên công ty: {task['customer_name']}
- Ngành nghề: {task['industry']}
- Thông tin doanh nghiệp: {company_text}

=== RFI (Khảo sát nhu cầu) ===
{rfi_text if rfi_text.strip() else 'Chưa có RFI cụ thể'}

=== THÔNG TIN CHI TIẾT SẢN PHẨM ĐƯỢC CHỌN ===
{products_detail if products_detail.strip() else 'Chưa chọn sản phẩm cụ thể - hãy đề xuất giải pháp phù hợp nhất'}

=== TẤT CẢ SẢN PHẨM CÓ THỂ ĐỀ XUẤT ===
{all_products_brief}
{kb_context}

Trả về JSON với format:
{{
  "cover_title": "Proposal Giải pháp {solution_names_text} cho {task['customer_name']}",
  "cover_subtitle": "{company_name} — Đối tác công nghệ tin cậy",
  "title": "GIẢI PHÁP [tên giải pháp chính]",
  "subtitle": "Dành cho {task['customer_name']}",
  "sections": [
    {{
      "heading": "Tên section",
      "type": "bullets" | "text" | "table" | "two_column",
      "content": [nội dung]
    }}
  ]
}}

CẤU TRÚC BẮT BUỘC — phải có ĐẦY ĐỦ 15 sections sau (đúng thứ tự):

1. "EXECUTIVE SUMMARY" (type: "text")
   — 1 đoạn 5-7 câu tóm gọn: vấn đề → giải pháp → kết quả kỳ vọng. C-level đọc xong hiểu deal.

2. "VỀ {company_name.upper()}" (type: "bullets")
   — 5-7 bullets: năm thành lập, số khách hàng, sản phẩm chính, chứng chỉ/giải thưởng, đối tác chiến lược.

3. "HIỂU KHÁCH HÀNG - {task['customer_name'].upper()}" (type: "bullets")
   — 5-7 bullets phân tích business KH: quy mô, kênh bán, điểm mạnh, thách thức hiện tại. Thể hiện "chúng tôi hiểu bạn".

4. "THÁCH THỨC & PAIN POINTS" (type: "bullets")
   — 5-7 bullets đặt vấn đề CỤ THỂ cho KH, gắn với business impact (mất tiền, mất KH, inefficiency). PHẢI bám sát brief.

5. "GIẢI PHÁP ĐỀ XUẤT - TỔNG QUAN" (type: "text")
   — Đoạn văn 5-8 câu: high-level tổng quan giải pháp, tại sao fit với KH, liên kết trực tiếp với pain points.

6-8. Với MỖI sản phẩm/giải pháp đề xuất, tạo 1 section riêng:
   "CHI TIẾT GIẢI PHÁP: [TÊN SẢN PHẨM]" (type: "bullets")
   — 8-12 tính năng, mỗi cái có mô tả 2-3 câu. PHẢI đề cập tính năng cụ thể liên quan đến brief (VD: nếu brief nói "gọi trên app" → phải có mục về SDK/WebRTC/in-app calling).

9. "KIẾN TRÚC KỸ THUẬT" (type: "bullets")
   — System architecture, integration points, data flow. Mô tả cách các module kết nối.

10. "CASE STUDY" (type: "bullets")
   — 3-4 bullets: KH cùng ngành hoặc quy mô đã triển khai thành công, có số liệu before/after.

11. "ROI & BUSINESS IMPACT" (type: "two_column")
   — Cột trái "Tiết kiệm chi phí": 4-5 items với số liệu %. Cột phải "Nâng cao hiệu quả": 4-5 items.
   Format: {{"left_title": "Tiết kiệm chi phí", "left_items": [...], "right_title": "Nâng cao hiệu quả", "right_items": [...]}}

12. "GÓI GIẢI PHÁP & PRICING" (type: "table")
   — Bảng [Hạng mục, Mô tả, Đơn giá tham khảo, Ghi chú]. Liệt kê: license, setup, training, support.

13. "LỘ TRÌNH TRIỂN KHAI" (type: "table")
   — Bảng [Giai đoạn, Thời gian, Nội dung, Deliverables]. 4-5 phases với timeline cụ thể.

14. "SLA & HỖ TRỢ" (type: "bullets")
   — 5-6 bullets: cam kết uptime, response time, dedicated support model, escalation.

15. "TẠI SAO CHỌN {company_name.upper()}" (type: "bullets")
   — 5-6 lý do differentiators, competitive advantage.

16. "BƯỚC TIẾP THEO & LIÊN HỆ" (type: "bullets")
   — 4-5 action items cụ thể, contact info.

QUAN TRỌNG:
- Viết bằng tiếng Việt, chuyên nghiệp, chi tiết
- BÁM SÁT BRIEF: mọi yêu cầu trong brief PHẢI được đáp ứng trong proposal
- Mỗi bullet point phải có giải thích 2-3 câu, KHÔNG chỉ liệt kê tên
- Chỉ trả về JSON, không thêm text khác"""

    async with httpx.AsyncClient(timeout=180.0) as client:
        resp = await client.post(
            "https://api.anthropic.com/v1/messages",
            headers={
                "x-api-key": ANTHROPIC_API_KEY,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": "claude-sonnet-4-20250514",
                "max_tokens": 12000,
                "messages": [{"role": "user", "content": prompt}],
            },
        )
        if resp.status_code == 200:
            text = resp.json()["content"][0]["text"].strip()
            if "```" in text:
                text = text.split("```")[1]
                if text.startswith("json"):
                    text = text[4:]
                text = text.strip()
            return json.loads(text)
        else:
            raise Exception(f"Claude API error: {resp.status_code}")


def _create_pptx_proposal(content: dict, output_path: Path, legal_entity: str):
    """Create PPTX proposal from AI content, keeping template slides for branding."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    entities = _get_legal_entities()
    entity_info = next((e for e in entities if e["id"] == legal_entity), entities[0] if entities else {"id": "default", "label": "Default", "template": ""})
    template_path = TEMPLATES_DIR / entity_info.get("template", "")

    vihat_blue = RGBColor(0x21, 0x96, 0xF3)
    dark = RGBColor(0x33, 0x33, 0x33)
    white = RGBColor(0xFF, 0xFF, 0xFF)

    if template_path.exists():
        prs = Presentation(str(template_path))
        template_slide_count = len(prs.slides)

        # Customize cover slide (slide 0) with proposal title
        cover_title = content.get("cover_title", f"Proposal cho {content.get('subtitle', '')}")
        cover_subtitle = content.get("cover_subtitle", "")
        if template_slide_count > 0:
            cover = prs.slides[0]
            # Find the largest text shape (likely the title) and second largest (subtitle)
            text_shapes = []
            for shape in cover.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    text_shapes.append(shape)
            # Sort by font size or shape height to find title vs subtitle
            if text_shapes:
                # Replace the first significant text with cover_title
                for shape in text_shapes:
                    tf = shape.text_frame
                    original_text = tf.text.strip().lower()
                    # Detect title-like shapes (usually the main heading)
                    if any(kw in original_text for kw in ["giải pháp", "nâng cao", "cung cấp", "đa kênh", "tổng đài", "omicall", "proposal"]):
                        # Preserve formatting but change text
                        for para in tf.paragraphs:
                            for run in para.runs:
                                run.text = cover_title
                                break
                            break
                        break
                else:
                    # If no match, try the first text shape
                    if text_shapes:
                        tf = text_shapes[0].text_frame
                        for para in tf.paragraphs:
                            for run in para.runs:
                                run.text = cover_title
                                break
                            break

        print(f"[Proposal] Cover: '{cover_title}'. Keeping {template_slide_count} template slides")
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        template_slide_count = 0

    # Determine slide dimensions
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    w_inches = slide_w / 914400  # EMU to inches
    h_inches = slide_h / 914400

    # Helper: Get best layout for content
    def get_layout(layout_name):
        """Find layout by name, fallback to BLANK."""
        for sl in prs.slide_layouts:
            if sl.name.upper() == layout_name.upper():
                return sl
        # Fallback: try TITLE_AND_BODY or BLANK
        for sl in prs.slide_layouts:
            if "BODY" in sl.name.upper() or "TEXT" in sl.name.upper():
                return sl
        return prs.slide_layouts[-1]  # BLANK

    def add_section_header(title_text):
        """Add a section header slide (big title, no body)."""
        layout = get_layout("SECTION_HEADER")
        slide = prs.slides.add_slide(layout)
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                shape.text = title_text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(32)
                        run.font.bold = True
                        run.font.color.rgb = vihat_blue
        return slide

    def add_content_slide(title_text, body_lines, bullet_style=True):
        """Add a slide with title + body text (bullets or paragraph)."""
        layout = get_layout("TITLE_AND_BODY")
        slide = prs.slides.add_slide(layout)

        title_set = False
        body_set = False
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0 and not title_set:
                shape.text = title_text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(24)
                        run.font.bold = True
                        run.font.color.rgb = vihat_blue
                title_set = True
            elif shape.placeholder_format.idx == 1 and not body_set:
                shape.text = ""
                tf = shape.text_frame
                tf.word_wrap = True
                for i, line in enumerate(body_lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = str(line)
                    p.font.size = Pt(13)
                    p.space_after = Pt(4)
                    if bullet_style:
                        p.level = 0
                body_set = True

        # If no placeholders found, use textbox fallback
        if not title_set:
            from pptx.util import Inches as In
            txBox = slide.shapes.add_textbox(In(0.5), In(0.2), In(w_inches - 1), In(0.8))
            p = txBox.text_frame.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = vihat_blue

        if not body_set and body_lines:
            from pptx.util import Inches as In
            txBox = slide.shapes.add_textbox(In(0.5), In(1.2), In(w_inches - 1), In(h_inches - 1.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, line in enumerate(body_lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(line)
                p.font.size = Pt(13)
                p.space_after = Pt(4)

        return slide

    def add_table_slide(title_text, table_data):
        """Add a slide with title + table."""
        from pptx.util import Inches as In

        layout = get_layout("BLANK")
        slide = prs.slides.add_slide(layout)

        # Title
        txBox = slide.shapes.add_textbox(In(0.5), In(0.2), In(w_inches - 1), In(0.8))
        p = txBox.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = vihat_blue

        # Table
        if not table_data:
            return slide
        rows = len(table_data)
        cols = max(len(row) for row in table_data) if table_data else 2
        # Calculate table dimensions
        table_w = w_inches - 1.0
        table_h = min(h_inches - 1.8, rows * 0.5)
        table_shape = slide.shapes.add_table(rows, cols, In(0.5), In(1.2), In(table_w), In(table_h)).table

        for r, row_data in enumerate(table_data):
            for c, cell_val in enumerate(row_data):
                if c < cols:
                    cell = table_shape.cell(r, c)
                    cell.text = str(cell_val)
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(11)
                        if r == 0:
                            para.font.bold = True
                            para.font.color.rgb = white
                    # Header row styling
                    if r == 0:
                        from pptx.oxml.ns import qn
                        tc = cell._tc
                        tcPr = tc.get_or_add_tcPr()
                        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '2196F3'})
                        solidFill.append(srgbClr)
                        tcPr.append(solidFill)

        return slide

    def add_two_column_slide(title_text, col_data):
        """Add a two-column layout slide."""
        from pptx.util import Inches as In

        layout = get_layout("TITLE_AND_TWO_COLUMNS")
        slide = prs.slides.add_slide(layout)

        # Try to use placeholders
        ph_filled = set()
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx == 0:
                shape.text = title_text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(24)
                        run.font.bold = True
                        run.font.color.rgb = vihat_blue
                ph_filled.add(0)
            elif idx == 1:
                # Left column
                shape.text = ""
                tf = shape.text_frame
                lt = col_data.get("left_title", "")
                if lt:
                    p = tf.paragraphs[0]
                    p.text = lt
                    p.font.size = Pt(14)
                    p.font.bold = True
                for item in col_data.get("left_items", []):
                    p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(12)
                    p.space_after = Pt(3)
                ph_filled.add(1)
            elif idx == 2:
                # Right column
                shape.text = ""
                tf = shape.text_frame
                rt = col_data.get("right_title", "")
                if rt:
                    p = tf.paragraphs[0]
                    p.text = rt
                    p.font.size = Pt(14)
                    p.font.bold = True
                for item in col_data.get("right_items", []):
                    p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(12)
                    p.space_after = Pt(3)
                ph_filled.add(2)

        # Fallback: use textboxes if placeholders not available
        if 0 not in ph_filled:
            txBox = slide.shapes.add_textbox(In(0.5), In(0.2), In(w_inches - 1), In(0.8))
            p = txBox.text_frame.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(24)
            p.font.bold = True
        if 1 not in ph_filled:
            col_w = (w_inches - 1.5) / 2
            # Left column
            txBox = slide.shapes.add_textbox(In(0.5), In(1.2), In(col_w), In(h_inches - 2))
            tf = txBox.text_frame
            tf.word_wrap = True
            lt = col_data.get("left_title", "")
            if lt:
                tf.paragraphs[0].text = lt
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.bold = True
            for item in col_data.get("left_items", []):
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)
        if 2 not in ph_filled:
            col_w = (w_inches - 1.5) / 2
            txBox = slide.shapes.add_textbox(In(0.5 + col_w + 0.5), In(1.2), In(col_w), In(h_inches - 2))
            tf = txBox.text_frame
            tf.word_wrap = True
            rt = col_data.get("right_title", "")
            if rt:
                tf.paragraphs[0].text = rt
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.bold = True
            for item in col_data.get("right_items", []):
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)

        return slide

    # ====== ADD A DIVIDER SLIDE before custom content ======
    if template_slide_count > 0:
        divider_title = content.get("title", "ĐỀ XUẤT GIẢI PHÁP").upper()
        add_section_header(divider_title)

    # ====== CONTENT SLIDES FROM AI ======
    for section in content.get("sections", []):
        heading = section.get("heading", "")
        sec_type = section.get("type", "bullets")
        sec_content = section.get("content", [])

        if sec_type == "text" and isinstance(sec_content, str):
            # Long text: split into chunks for readability
            paragraphs = sec_content.split("\n")
            add_content_slide(heading, paragraphs, bullet_style=False)

        elif sec_type == "bullets" and isinstance(sec_content, list):
            items = [str(x) for x in sec_content]
            # Split into pages of 6 items each
            for i in range(0, len(items), 6):
                chunk = items[i:i+6]
                page_heading = heading if i == 0 else f"{heading} (tiếp)"
                add_content_slide(page_heading, [f"• {item}" for item in chunk])

        elif sec_type == "table" and isinstance(sec_content, list) and len(sec_content) > 0:
            # Split large tables across multiple slides
            header = sec_content[0] if sec_content else []
            data_rows = sec_content[1:] if len(sec_content) > 1 else []
            for i in range(0, max(len(data_rows), 1), 8):
                chunk = [header] + data_rows[i:i+8]
                page_heading = heading if i == 0 else f"{heading} (tiếp)"
                add_table_slide(page_heading, chunk)

        elif sec_type == "two_column" and isinstance(sec_content, dict):
            add_two_column_slide(heading, sec_content)

        else:
            # Fallback
            if isinstance(sec_content, list):
                add_content_slide(heading, [str(x) for x in sec_content])
            elif isinstance(sec_content, str):
                add_content_slide(heading, [sec_content], bullet_style=False)
            else:
                add_content_slide(heading, [str(sec_content)])

    # ====== CLOSING SLIDE ======
    layout = get_layout("TITLE")
    slide = prs.slides.add_slide(layout)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = "CẢM ƠN QUÝ KHÁCH"
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(36)
                    run.font.bold = True
                    run.font.color.rgb = vihat_blue
        elif shape.placeholder_format.idx == 1:
            shape.text = f"Liên hệ: {entity_info['label']}\nWebsite: vihat.vn | Hotline: 1900 6181\nEmail: info@vihat.vn"
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(16)

    prs.save(str(output_path))
    print(f"[Proposal] PPTX created: {template_slide_count} template + {len(prs.slides) - template_slide_count} content slides")


def _create_docx_proposal(content: dict, output_path: Path, legal_entity: str):
    """Create DOCX proposal from AI content."""
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    entities = _get_legal_entities()
    entity_info = next((e for e in entities if e["id"] == legal_entity), entities[0] if entities else {"id": "default", "label": "Default", "template": ""})
    doc = Document()

    # Styles
    style = doc.styles['Heading 1']
    style.font.color.rgb = RGBColor(0x21, 0x96, 0xF3)
    style.font.size = Pt(20)

    style2 = doc.styles['Heading 2']
    style2.font.color.rgb = RGBColor(0x21, 0x96, 0xF3)
    style2.font.size = Pt(16)

    # Cover page
    title = content.get("title", "Proposal")
    subtitle = content.get("subtitle", "")

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"\n\n\n\n{entity_info['label']}\n\n")
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(title)
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x21, 0x96, 0xF3)

    if subtitle:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(subtitle)
        run.font.size = Pt(14)

    doc.add_page_break()

    # Content sections
    for section in content.get("sections", []):
        heading = section.get("heading", "")
        sec_type = section.get("type", "bullets")
        sec_content = section.get("content", [])

        doc.add_heading(heading, level=1)

        if sec_type == "text" and isinstance(sec_content, str):
            doc.add_paragraph(sec_content)
        elif sec_type == "bullets" and isinstance(sec_content, list):
            for item in sec_content:
                doc.add_paragraph(str(item), style='List Bullet')
        elif sec_type == "table" and isinstance(sec_content, list) and len(sec_content) > 0:
            cols = max(len(row) for row in sec_content)
            table = doc.add_table(rows=len(sec_content), cols=cols)
            table.style = 'Light Grid Accent 1'
            for r, row_data in enumerate(sec_content):
                for c, cell_val in enumerate(row_data):
                    if c < cols:
                        table.cell(r, c).text = str(cell_val)
        else:
            lines = sec_content if isinstance(sec_content, list) else [str(sec_content)]
            for line in lines:
                doc.add_paragraph(str(line))

        doc.add_paragraph()  # spacing

    # Footer
    doc.add_page_break()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"\n\n{entity_info['label']}\nWebsite: vihat.vn | Hotline: 1900 6181")
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    doc.save(str(output_path))


async def _run_proposal_generation(task_id: str):
    """Background worker for proposal generation."""
    task = PROPOSAL_TASKS[task_id]
    try:
        # Phase 1: Generate content with AI
        task["status"] = "generating_content"
        content = await _generate_proposal_content(task)

        # Phase 2: Create document
        task["status"] = "creating_document"
        safe_name = re.sub(r'[^\w\s-]', '', task["customer_name"]).strip().replace(' ', '_')
        filename = f"proposal_{safe_name}_{task_id}.{task['output_format']}"
        output_path = PROPOSALS_DIR / filename

        if task["output_format"] == "pptx":
            await asyncio.to_thread(_create_pptx_proposal, content, output_path, task.get("legal_entity", "omijsc"))
        else:
            await asyncio.to_thread(_create_docx_proposal, content, output_path, task.get("legal_entity", "omijsc"))

        task["status"] = "completed"
        task["file_path"] = str(output_path)
        task["file_name"] = filename
        task["completed_at"] = datetime.now().isoformat()
        _save_proposal_tasks()
        print(f"[Proposal] Task {task_id} completed: {filename}")

    except Exception as e:
        import traceback
        error_detail = f"{type(e).__name__}: {e}"
        traceback_str = traceback.format_exc()
        task["status"] = "error"
        task["error"] = error_detail
        task["completed_at"] = datetime.now().isoformat()
        _save_proposal_tasks()
        print(f"[Proposal] Task {task_id} error: {error_detail}\n{traceback_str}")


@app.post("/api/v1/proposals/generate")
async def generate_proposal(req: ProposalGenerateRequest):
    """Start async proposal generation."""
    if not req.customer_name.strip():
        raise HTTPException(400, "Tên khách hàng là bắt buộc")
    if not ANTHROPIC_API_KEY:
        raise HTTPException(500, "ANTHROPIC_API_KEY not configured")

    task_id = str(uuid.uuid4())[:8]
    PROPOSAL_TASKS[task_id] = {
        "task_id": task_id,
        "status": "generating_content",
        "customer_name": req.customer_name,
        "industry": req.industry,
        "products": req.products,
        "rfi_answers": req.rfi_answers,
        "company_info": req.company_info,
        "legal_entity": req.legal_entity,
        "output_format": req.output_format,
        "brief_text": req.brief_text,
        "started_at": datetime.now().isoformat(),
        "completed_at": None,
        "file_name": None,
        "error": None,
    }

    _save_proposal_tasks()
    asyncio.create_task(_run_proposal_generation(task_id))

    return {"task_id": task_id, "status": "generating_content", "message": "Đang tạo proposal..."}


@app.get("/api/v1/proposals/tasks")
async def list_proposal_tasks():
    """List all proposal tasks."""
    tasks = sorted(PROPOSAL_TASKS.values(), key=lambda t: t.get("started_at", ""), reverse=True)
    # Return safe subset (no file_path)
    return [
        {k: v for k, v in t.items() if k != "file_path"}
        for t in tasks
    ]


@app.get("/api/v1/proposals/{task_id}/download")
async def download_proposal(task_id: str):
    """Download generated proposal file."""
    from fastapi.responses import FileResponse

    task = PROPOSAL_TASKS.get(task_id)
    if not task:
        raise HTTPException(404, "Task not found")
    if task["status"] != "completed":
        raise HTTPException(400, f"Task is {task['status']}, not ready for download")

    file_path = task.get("file_path")
    if not file_path or not Path(file_path).exists():
        raise HTTPException(404, "File not found")

    media_types = {
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    }

    return FileResponse(
        file_path,
        filename=task.get("file_name", f"proposal.{task['output_format']}"),
        media_type=media_types.get(task["output_format"], "application/octet-stream"),
    )


# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/health")
async def health():
    dify_ok = bool(DIFY_API_KEY)
    return {"status": "ok", "service": "vihat-knowledge-api", "mode": "dev", "dify_connected": dify_ok}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("dev_mock_server:app", host="0.0.0.0", port=8000, reload=True)
